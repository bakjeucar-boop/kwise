"""PPT 보고서 (36세션).

Word 다섯 장을 **옮기는 것이 아니라 구조를 다시 짠다.** 문서의 절과 슬라이드
한 장은 담는 양이 다르다 — 글은 줄이고 그림을 키운다. **한 장에 한 메시지다.**

    표지 · 목차
    건물현황 및 계약정보 · 전력사용현황 · 부하패턴 및 피크특성 · 피크특성
    현재 요금 구조 · 개선안별 요약
    검토한 수단별 1장씩 (**켠 수단만.** 0개면 그 장이 통째로 빠진다)
    조합구성 및 합산효과 · 부록 산출근거 상세

**색·크기는 여기서 정하지 않는다.** 전부 ``data\\ppt_design.json`` 에서 오고
:mod:`kwise.report.design` 이 읽는다 (36세션 3절). 그림 png 도 같은 팔레트를
쓴다 (:func:`kwise.report.figures.chart_palette`).

**재료는 Word 와 같은 :class:`~kwise.report.document.DocumentSections` 다.** 두
산출물이 같은 값을 봐야 나란히 놓았을 때 어긋나지 않는다. 다른 것은 무엇을
싣느냐뿐이다 — PPT 부록은 **산출 근거만** 담는다 (36세션 6절).

**레이아웃을 섞는다** (36세션 3-4). 카드 기반으로만 만들지 않는다 — 표·리스트·
여백만인 슬라이드를 함께 둔다. :data:`SlideSpec.layout` 이 그 형태이고,
시험이 형태별 개수를 센다.

파일명에 날짜·시각 접미사를 붙인다 — PowerPoint 가 파일을 열고 있으면
덮어쓰기가 실패한다 (Excel·Word 와 같은 이유).
"""

from __future__ import annotations

import datetime as dt
import io
import math
from collections.abc import Callable, Sequence
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.parts.image import Image
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Emu, Inches, Pt

from kwise import money
from kwise.diagnose import ChargeStructure
from kwise.measures import BASE_FEE_UNCHANGED
from kwise.report import figures, narrative
from kwise.report.design import DesignGuide, load_design_guide
from kwise.report.document import DocumentSections, MeasureEntry, MeasureFigure
from kwise.report.narrative import GLOSSARY_KEYS
from kwise.report.notices import (
    NOT_INCLUDED_NOTICE,
    TRUNCATION_FOOTNOTE,
    plain_text,
    rules_basis_line,
)
from kwise.report.worksheet import COLUMNS
from kwise.tariff.labels import SEASON_LABELS

__all__ = [
    "ANNUAL_BASIS_NOTE",
    "APPENDIX_SLIDE_TITLE",
    "CLOSING_SLIDE_TITLE",
    "DECK_TITLE",
    "DEFAULT_OUTPUT_DIR",
    "FULL_FIGURE",
    "FULL_FIGURE_WITH_LEGEND",
    "HALF_FIGURE",
    "HALF_FIGURE_WITH_LEGEND",
    "IMMEDIATE",
    "LAYOUTS",
    "MEASURE_AGENDA_ITEM",
    "NEXT_STEPS",
    "NEXT_STEPS_HEADLINE",
    "NOTE_MARK",
    "NO_INVESTMENT",
    "SLIDE_TITLES",
    "SPEC_TABLE_WIDTHS",
    "AppendixPage",
    "SlideSpec",
    "agenda_items",
    "appendix_chunks",
    "appendix_pages",
    "build_slides",
    "export_slides",
    "mark_note",
    "measure_slide_title",
    "plain_text",
    "season_pairs",
    "slide_investment",
    "slide_payback",
    "slide_specs",
    "slides_bytes",
    "slides_path",
    "split_reason",
]

DECK_TITLE = "전력 비용 진단 보고서"
DEFAULT_OUTPUT_DIR = Path("output")
APPENDIX_SLIDE_TITLE = "부록 산출근거 상세"
CLOSING_SLIDE_TITLE = "다음 단계"

#: 슬라이드 제목 — **지시서의 괄호 하나가 한 장이다** (36세션 2절).
SLIDE_TITLES: dict[str, str] = {
    "cover": DECK_TITLE,
    "agenda": "목차",
    "building": "건물현황 및 계약정보",
    "usage_pattern": "전력사용현황 및 부하패턴",
    "peak_summary": "피크특성 (1/2)",
    "peak_detail": "피크특성 (2/2)",
    "structure": "현재 요금 구조",
    "measure_summary": "개선안별 요약",
    "surplus": "잉여 활용",
    "combination": "조합구성 및 합산효과",
    "appendix": APPENDIX_SLIDE_TITLE,
    "closing": CLOSING_SLIDE_TITLE,
}

#: 레이아웃 형태 (36세션 3-4). **같은 형태를 반복하지 않는다.**
#:
#:     cover        전체 다크 배경 밴드 — 샌드위치의 바깥
#:     agenda       리스트형. 구분선으로 나눈다
#:     table        표형. **구분선 기반이고 카드가 아니다**
#:     chart        그림 한 장을 크게. 글은 캡션 한 줄
#:     stat_chart   통계 강조형 — 지표 + 차트
#:     chart_pair   같은 크기의 그림 둘을 좌우로
#:     split        좌우가 서로 다른 것을 담는다 (지표+차트 | 격자)
#:     stat_table_chart  지표 + 폭 전체 표 + 그 아래 그림 (46세션. ESS)
#:     fact_stat_chart   **근거 지표**가 먼저, 결과 지표가 뒤 (53세션. 계약전력)
#:     compare      비교형 — 좌우 2단에 중앙 얇은 구분선
#:     closing      마무리 — 표지와 짝이 되는 다크 배경 밴드 (37세션)
LAYOUTS: tuple[str, ...] = (
    "cover",
    "agenda",
    "table",
    "chart",
    "stat_chart",
    "chart_pair",
    "split",
    "compare",
    "closing",
    "stat_table_chart",
    "fact_stat_chart",
)

_UNPRICED = "미산출"

#: **참고용 작은 글씨 앞에 붙이는 표식** (53세션 1-1).
#:
#: 툴팁에서 옮긴 용어 풀이, 전제·한계 각주, 표 아래 참고 한 줄이 전부 이것을
#: 단다. **그림 캡션은 제외한다** — 캡션은 그림이 무엇인지 말하는 이름이지
#: 참고가 아니다.
NOTE_MARK = "※ "


def mark_note(line: str) -> str:
    """참고 한 줄에 :data:`NOTE_MARK` 를 붙인다. **이미 붙었으면 그대로 둔다.**"""
    text = line.strip()
    if not text or text.startswith(NOTE_MARK.strip()):
        return text
    return f"{NOTE_MARK}{text}"


#: 「미산출 — 사유」 를 가르는 표식.
_REASON_MARK = " — "


#: 「값 — 사유」 꼴로 나오는 머리말들 (39세션 2-1 · 53세션 6절).
#:
#: 「기본요금 변화없음」 이 53세션에 붙었다 — 48세션이 「0원」 대신 이 말을
#: 세우면서 사유가 통째로 큰 글씨 자리에 들어와, 계약전력 장의 지표 칸 하나가
#: **두 줄로 흘러** 옆 셋과 위계가 어긋났다.
_REASON_HEADS: tuple[str, ...] = (_UNPRICED, BASE_FEE_UNCHANGED)


def split_reason(value: str) -> tuple[str, str]:
    """「미산출 — 투자비 미입력」 을 (「미산출」, 「투자비 미입력」) 로 가른다 (39세션 2-1).

    **지표 칸에는 값만 둔다.** 큰 글씨 자리에 사유가 들어가면 두 줄로 흐르고,
    넷을 나란히 놓았을 때 하나만 길어 위계가 무너진다. 사유는 각주가 받는다 —
    무엇을 넣으면 값이 나오는지가 고객에게 필요한 정보이므로 지우지는 않는다.

    :data:`_REASON_HEADS` 로 시작하지 않는 값은 그대로 돌려준다.
    """
    if _REASON_MARK not in value or not value.startswith(_REASON_HEADS):
        return value, ""
    head, _, reason = value.partition(_REASON_MARK)
    return head.strip(), reason.strip()


#: 글자 하나가 차지하는 폭 — 글꼴 크기에 대한 비율.
#:
#: **한글은 전각, 숫자·영문은 반각에 가깝다.** 정확한 값은 글꼴이 쥐고 있어
#: 여기서 알 수 없으므로 **넉넉히 잡는다** — 잘못 잡으면 큰 글씨가 두 줄로
#: 넘쳐 아래 그림을 덮는다. 좁게 잡아 빈자리가 남는 편이 낫다.
_WIDE_GLYPH = 1.0
_NARROW_GLYPH = 0.6


def _text_width_in(text: str, size_pt: float) -> float:
    """글자열이 차지할 대략의 폭 (in)."""
    units = sum(_WIDE_GLYPH if ord(char) > 0x2000 else _NARROW_GLYPH for char in text)
    return units * size_pt / 72.0


def _fitting_size(values: Sequence[str], *, span: float, ladder: Sequence[float]) -> float:
    """칸 폭에 **한 줄로** 들어가는 가장 큰 글자 크기. 없으면 사다리의 끝이다."""
    for size in ladder:
        if all(_text_width_in(value, size) <= span for value in values):
            return size
    return ladder[-1]


def _fitting_lines(values: Sequence[str], *, span: float, size: float) -> int:
    """그 크기로 적었을 때 **가장 긴 값이 차지할 줄 수** (38세션 4절).

    사다리 끝에서도 한 줄에 못 들어가는 값이 있다 — 「9,050,000원 (12개월 환산
    9,050,000원)」 처럼 한 칸에 두 값을 담은 문구다. 그런 값은 두 줄로 흐르는데,
    **블록 높이를 한 줄로 잡아 두면 아래 그림 자리를 먹는다.** 자르지 않고
    자리를 내주되, 몇 줄인지를 세어 그만큼 높이를 잡는다.

    셋 이상은 세지 않는다 — 그만한 값이 지표 칸에 들어갈 자리가 아니다.
    """
    if span <= 0:
        return 1
    longest = max((_text_width_in(value, size) for value in values), default=0.0)
    return min(2, max(1, int(-(-longest // span))))


@dataclass(frozen=True)
class SlideSpec:
    """슬라이드 한 장의 자리표.

    **덱은 이 목록대로 만들어진다.** 구성과 실물이 갈라지지 않게 하려면 차례를
    쥔 자리가 하나여야 한다 — 시험도 이것을 본다.
    """

    key: str
    title: str
    layout: str
    measure: int | None = None
    """수단별 장이면 :attr:`DocumentSections.measures` 의 자리."""
    page: int | None = None
    """부록처럼 여러 장으로 나뉘는 자리의 쪽 번호 (0부터)."""


def slide_specs(sections: DocumentSections) -> tuple[SlideSpec, ...]:
    """덱의 차례 (36세션 2절).

    **수단을 하나도 켜지 않으면 수단별 장만 빠진다.** 나머지는 그대로다 —
    진단만 보고 받아 가는 것이 정상 경로다 (Word 와 같은 규약).

    표지·목차·**마무리**는 자료가 무엇이든 늘 나온다. 셋 다 내용이 아니라
    덱의 뼈대이기 때문이다 (37세션).
    """
    specs = [
        SlideSpec("cover", SLIDE_TITLES["cover"], "cover"),
        SlideSpec("agenda", SLIDE_TITLES["agenda"], "agenda"),
        SlideSpec("building", SLIDE_TITLES["building"], "table"),
        SlideSpec("usage_pattern", SLIDE_TITLES["usage_pattern"], "stat_chart"),
        SlideSpec("peak_summary", SLIDE_TITLES["peak_summary"], "stat_chart"),
        SlideSpec("peak_detail", SLIDE_TITLES["peak_detail"], "chart_pair"),
        SlideSpec("structure", SLIDE_TITLES["structure"], "split"),
        SlideSpec("measure_summary", SLIDE_TITLES["measure_summary"], "table"),
    ]
    # **그림이 둘이면 좌우로 나눈다** (38세션 3절). 형태를 자리표가 쥐고 있어야
    # 시험이 「어느 장이 무엇으로 서는가」 를 셀 수 있다.
    for index, entry in enumerate(sections.measures):
        specs.append(
            SlideSpec(
                f"measure_{entry.kind.key}",
                measure_slide_title(entry),
                _measure_layout(entry),
                measure=index,
            )
        )
        # **잉여 활용은 태양광 다음 한 장이다** (53세션 3-1). 개선안이 아니라
        # 태양광의 결과이므로 그 장 바로 뒤에 붙는다. **잉여가 0 이면 없다.**
        if entry.kind.key == "solar" and sections.surplus is not None:
            specs.append(SlideSpec("surplus", SLIDE_TITLES["surplus"], "split"))
    specs.append(SlideSpec("combination", SLIDE_TITLES["combination"], "compare"))
    # **부록은 수단마다 한 장 이상이다** (39세션 5절). 한 장에 눌러 담고 「자리가
    # 모자라 뺐다」 고 적으면 여섯 수단이 통째로 사라진다.
    specs.extend(
        SlideSpec("appendix", page.title, "table", page=index)
        for index, page in enumerate(appendix_pages(sections))
    )
    specs.append(SlideSpec("closing", SLIDE_TITLES["closing"], "closing"))
    return tuple(specs)


#: 수단별 장을 가리키는 목차 한 줄 (38세션 1-1).
#:
#: **수단을 나열하지 않는다.** 일곱을 이어 적으면 줄이 두 줄로 넘쳐 다음 항목
#: 자리를 덮는다 — 07번이 08번을 가리고 있었다. 어느 수단을 보았는지는 바로
#: 다음 장인 「개선안별 요약」 표가 낸다.
MEASURE_AGENDA_ITEM = "검토한 수단별 상세"


def measure_slide_title(entry: MeasureEntry) -> str:
    """수단 장 제목 — **절 번호를 뗀다** (38세션 1-3).

        7.5 태양광  →  태양광

    「7.」 이 무엇인지 덱 어디에도 적혀 있지 않다. 화면은 27세션에 순번(1~7)으로
    바꿨는데 PPT 만 요구사항서 번호를 그대로 내고 있었다 — 처음 받아 보는
    사람에게는 없는 7장을 찾게 만드는 표시다.

    **순번도 붙이지 않는다.** 목차가 수단을 나열하지 않게 되었고(1-1) 장은 차례로
    넘어가므로, 번호가 지는 몫이 없다. 켠 수단만 실리는 덱이라 「5.」 로 시작하는
    일도 생긴다.

    **Word·Excel 의 7.x 는 그대로다** — 거기서는 요구사항서와 맞물린 번호다.
    정본(:attr:`kwise.measures.MeasureKind.title`)을 고치지 않고 낼 때만 바꾼다.
    """
    return entry.kind.label


def agenda_items(sections: DocumentSections) -> tuple[str, ...]:
    """목차에 적을 줄. **수단별 장은 한 줄로 묶는다.**

    슬라이드마다 한 줄을 적으면 수단이 일곱일 때 목차가 열두 줄이 된다 — 목차는
    어디를 보는지 알려 주는 자리이지 슬라이드 색인이 아니다.
    """
    lines = [
        SLIDE_TITLES["building"],
        SLIDE_TITLES["usage_pattern"],
        SLIDE_TITLES["peak_summary"],
        SLIDE_TITLES["peak_detail"],
        SLIDE_TITLES["structure"],
        SLIDE_TITLES["measure_summary"],
    ]
    if sections.measures:
        lines.append(MEASURE_AGENDA_ITEM)
    # **잉여 활용은 수단이 아니라 결과다** (53세션 3-3). 「검토한 수단별 상세」
    # 안에 넣으면 개선안이 일곱으로 읽힌다 — 41세션이 여섯으로 줄인 그 수다.
    if sections.surplus is not None:
        lines.append(SLIDE_TITLES["surplus"])
    lines.append(SLIDE_TITLES["combination"])
    lines.append(SLIDE_TITLES["appendix"])
    lines.append(SLIDE_TITLES["closing"])
    return tuple(lines)


def season_pairs(structure: ChargeStructure) -> tuple[tuple[str, str | None], ...]:
    """도넛 넷의 갈래 (36세션 4절). **화면과 같은 차례다.**

    자료에 없는 계절은 두지 않는다 — 빈 원은 「그 계절에 안 썼다」 로 읽힌다
    (화면의 ``season_choices`` 와 같은 규칙이고, 이유도 같다).
    """
    available = {str(key) for key in structure.band_season_kwh.index}
    pairs: list[tuple[str, str | None]] = [("전체", None)]
    pairs.extend((label, key) for key, label in SEASON_LABELS.items() if key in available)
    return tuple(pairs)


# ===================================================================== 그리기 도구
#
# **색·크기를 여기서 정하지 않는다.** 전부 :class:`DesignGuide` 에서 온다.


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#").upper())


def _apply_font(run: object, guide: DesignGuide, size: float, color: str, *, bold: bool) -> None:
    """글꼴·크기·색. **동아시아 글꼴을 따로 지정해야** PowerPoint 가 한글에 쓴다."""
    font = run.font  # type: ignore[attr-defined]
    name = guide.typography.primary
    font.name = name
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = _rgb(color)
    rpr = run._r.get_or_add_rPr()  # type: ignore[attr-defined]
    for tag in ("a:ea", "a:cs"):
        element = rpr.find(qn(tag))
        if element is None:
            element = rpr.makeelement(qn(tag), {})
            rpr.append(element)
        element.set("typeface", name)


def _text(
    slide: Slide,
    guide: DesignGuide,
    lines: Sequence[str],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    size: float,
    color: str,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    spacing: float = 1.25,
) -> None:
    """텍스트 상자 하나. **본문은 항상 좌측 정렬이다** (36세션 3-3).

    높이에 :meth:`~kwise.report.design.SlideGeometry.slack` 만큼 여유를 둔다 —
    렌더러마다 글자 폭이 달라 딱 맞춘 상자는 다른 PowerPoint 에서 넘친다.
    """
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(guide.slide.slack(height))
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.line_spacing = spacing
        run = paragraph.add_run()
        run.text = plain_text(line)
        _apply_font(run, guide, size, color, bold=bold)


def _line(
    slide: Slide,
    guide: DesignGuide,
    *,
    left: float,
    top: float,
    right: float,
    bottom: float,
    color: str,
    width_pt: float | None = None,
) -> None:
    """얇은 선 하나. **테마 효과를 뗀다.**

    ``add_connector`` 는 테마의 ``effectRef`` 를 달고 나온다 — 옅은 그림자다.
    가이드가 무거운 그림자를 금지하므로(3-5) 애초에 효과를 걸지 않는다.
    """
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(left), Inches(top), Inches(right), Inches(bottom)
    )
    style = line._element.find(qn("p:style"))
    if style is not None:
        line._element.remove(style)
    line.line.color.rgb = _rgb(color)
    line.line.width = Pt(width_pt if width_pt is not None else guide.slide.rule_pt)


def _rule(
    slide: Slide,
    guide: DesignGuide,
    *,
    left: float,
    top: float,
    length: float,
    color: str,
    width_pt: float | None = None,
) -> None:
    """가로 구분선. **슬라이드 폭을 가로지르는 색 바가 아니다** (36세션 3-5).

    표와 리스트의 줄을 가르는 얇은 선이며, 제목 아래에는 긋지 않는다.
    """
    _line(
        slide,
        guide,
        left=left,
        top=top,
        right=left + length,
        bottom=top,
        color=color,
        width_pt=width_pt,
    )


def _vrule(
    slide: Slide, guide: DesignGuide, *, left: float, top: float, length: float, color: str
) -> None:
    """세로 구분선. 비교형의 **중앙 얇은 선**과 지표 사이를 가른다 (36세션 3-4)."""
    _line(slide, guide, left=left, top=top, right=left, bottom=top + length, color=color)


def _title(slide: Slide, guide: DesignGuide, text: str) -> float:
    """슬라이드 제목 하나. **큰 타이틀은 한 장에 하나다** (36세션 3-3).

    **아래에 강조선을 긋지 않는다** (3-5). 돌려주는 값은 본문이 시작할 y 다.
    """
    geometry = guide.slide
    height = 0.62
    _text(
        slide,
        guide,
        [text],
        left=geometry.margin_in,
        top=geometry.margin_in,
        width=geometry.content_width_in,
        height=height,
        size=guide.type_scale.slide_title,
        color=guide.colors.ink,
        bold=True,
    )
    return geometry.margin_in + height + geometry.title_gap_in


#: 본문 한 줄이 차지하는 높이 (in). 글자 크기가 아니라 **줄 간격**이다.
_LINE_HEIGHT = 0.28

#: 글 덩어리 아래에 두는 숨 (in).
_TEXT_GAP = 0.16


def _wrapped_height(text: str, guide: DesignGuide, *, width: float, size: float) -> float:
    """글이 차지할 높이 (in) — **줄 수를 세어 잡는다** (39세션 7절).

    해석 한 줄과 결론 한 줄이 붙으면서 두 줄로 흐르는 자리가 생겼다. 높이를
    한 줄로 고정해 두면 넘친 줄이 아래 지표 라벨과 겹친다 — 경제성DR 장이
    그랬다. 폭으로 재어 그만큼 내려 준다.

    **넉넉히 잡는다** (:data:`_WIDE_GLYPH` 와 같은 이유). 좁게 잡아 빈자리가
    남는 편이 겹치는 것보다 낫다.
    """
    if not text or width <= 0:
        return 0.0
    lines = max(1, math.ceil(_text_width_in(text, size) / width))
    return _LINE_HEIGHT * min(lines, 3) + _TEXT_GAP


def _lead(slide: Slide, guide: DesignGuide, text: str, *, top: float) -> float:
    """제목 아래 **해석 한 줄** (39세션 1-1).

    그림만 나열하면 무엇을 보아야 하는지 알 수 없다. 문장은
    :mod:`kwise.report.narrative` 가 짓는다 — 진단이 이미 내린 판정을 옮기거나,
    값의 크기로 고르거나, 고정 문장이다.

    **돌려주는 값은 본문이 시작할 y 다.** 빈 글이면 자리를 먹지 않는다.
    """
    if not text:
        return top
    geometry = guide.slide
    size = guide.type_scale.body
    height = _wrapped_height(text, guide, width=geometry.content_width_in, size=size)
    _text(
        slide,
        guide,
        [text],
        left=geometry.margin_in,
        top=top,
        width=geometry.content_width_in,
        height=height,
        size=size,
        color=guide.colors.ink,
    )
    return top + height


def _note(slide: Slide, guide: DesignGuide, *lines: str) -> None:
    """슬라이드 **맨 아래 작은 글씨** — 용어 풀이와 각주 (39세션 1-2).

    화면 툴팁의 산식 줄을 그대로 깐다. 그 장에 나오는 용어만 고르므로 대개 한
    줄에 앉는다 — 전부 깔면 각주가 본문이 된다.

    줄을 여럿 주면 위에서부터 쌓는다. **용어 풀이가 먼저다** — 표를 읽는 데
    바로 쓰이고, 제도 각주는 그 아래에서 받는다.

    **앞에 :data:`NOTE_MARK` 를 붙인다** (53세션 1-1). 작은 회색 글씨라는 것만으로는
    본문의 끝인지 참고인지 갈리지 않는다 — 표식 하나가 그 자리를 정한다.
    """
    kept, height = _note_block(guide, lines)
    if not kept:
        return
    geometry = guide.slide
    size = guide.type_scale.caption
    _text(
        slide,
        guide,
        kept,
        left=geometry.margin_in,
        top=_note_top(guide, *lines),
        width=geometry.content_width_in,
        height=height,
        size=size,
        color=guide.colors.muted,
    )


def _note_block(guide: DesignGuide, lines: Sequence[str]) -> tuple[list[str], float]:
    """각주로 실을 줄과 그 높이 (in) (53세션 9절).

    **줄 수를 세어 잡는다.** 각주는 슬라이드 아래에 붙어 있으므로 넘친 줄이
    밖으로 나가고, 위쪽 표는 그 사실을 모른 채 자리를 다 쓴다 — 개선안별 요약이
    그랬다: 각주가 세 줄인데 예약은 두 줄어치라 **표 마지막 줄과 겹쳤다.**
    """
    kept = [mark_note(line) for line in lines if line]
    if not kept:
        return [], 0.0
    size = guide.type_scale.caption
    width = guide.slide.content_width_in
    wrapped = sum(
        max(1, math.ceil(_text_width_in(line, size) / width)) for line in kept
    )
    return kept, 0.2 * wrapped


def _note_top(guide: DesignGuide, *lines: str) -> float:
    """각주 띠가 시작하는 y. **각주가 없으면 본문 바닥이다.**

    **자리를 못박지 않고 줄 수를 센다** (53세션 9절 · 59세션 1절). 53세션까지는
    ``_body_bottom`` 이 각주 높이를 0.58in 으로 고정해 두었는데, 그것은 두 줄어치다
    — 세 줄로 흐르면 위쪽 덩어리가 0.24in 을 침범했다. 표는 53세션에 이 함수로
    옮겼고, **그림 덩어리 넷(4·5·6·7장)이 남아 있어 59세션에 함께 옮겼다.**
    """
    geometry = guide.slide
    _kept, height = _note_block(guide, lines)
    if not height:
        return geometry.height_in - geometry.margin_in
    return geometry.height_in - geometry.margin_in - 0.42 - height + 0.2


def _caption(
    slide: Slide,
    guide: DesignGuide,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """작은 회색 한 줄. **그림 캡션은 가운데다** (53세션 1-2).

    왼쪽에 붙여 두면 캡션이 그림 아래가 아니라 칸의 왼쪽 귀퉁이에 매달린 것처럼
    보인다 — 그림이 칸 가운데 앉기 때문이다(:func:`_picture_block`). 부르는 쪽이
    정렬을 고르되 **기본은 왼쪽**이다: 부록 각주처럼 그림과 무관한 줄도 이 함수를
    쓴다.
    """
    _text(
        slide,
        guide,
        [text],
        left=left,
        top=top,
        width=width,
        height=0.22,
        size=guide.type_scale.caption,
        color=guide.colors.muted,
        align=align,
    )


#: 캡션 한 줄이 차지하는 높이 (in). 그림 덩어리의 높이 계산에 함께 든다.
_CAPTION_HEIGHT = 0.32


def _table_room_above(guide: DesignGuide, *, top: float, note_top: float) -> float:
    """각주 띠 **바로 위까지** 표에 내줄 수 있는 높이 (53세션 9절).

    :func:`_table_room` 은 각주 자리를 0.65in 로 못박는데, 각주가 세 줄이면
    0.82in 을 먹는다 — 그 차이가 그대로 겹침이 됐다. 실측한 자리를 받는다.
    """
    geometry = guide.slide
    room = note_top - geometry.block_gap_in - top
    return max(0.0, room) / (1.0 + geometry.text_slack)


def _table_room(guide: DesignGuide, *, top: float, footnote: bool) -> float:
    """표에 내줄 수 있는 높이 (38세션 4절).

    **행 높이의 여유(slack)를 미리 빼 둔다.** :func:`_table` 은 줄마다 10% 여유를
    두므로(36세션 3-1) 넘긴 높이보다 실제 표가 그만큼 크다 — 빼지 않고 칸을 꽉
    채워 넘기면 마지막 줄이 아래 각주에 닿는다. 부록 장이 그러고 있었다.
    """
    geometry = guide.slide
    room = geometry.height_in - geometry.margin_in - top
    if footnote:
        room -= 0.3 + geometry.block_gap_in
    return max(0.0, room) / (1.0 + geometry.text_slack)


def _aspect(png: bytes) -> float:
    """그림의 세로÷가로. **꾸러미에 넣지 않고** 바이트에서 읽는다."""
    width, height = Image.from_blob(png).size
    return height / width if width else 1.0


# ===================================================================== 그림 비율
#
# **Word 의 그림 비율을 그대로 쓰지 않는다.** 9:3.6 은 문서 한 단(6.3in)에 맞춘
# 것이라 16:9 슬라이드의 반 칸에 넣으면 위아래가 통째로 남는다. 슬라이드가
# 자기 칸에 맞는 크기를 **부를 때 정해서** 넘긴다 (:mod:`kwise.report.figures`
# 의 ``size``). 값은 칸의 가로세로에서 나온 것이라 여기 둔다.

#: 반 폭 칸(좌우 2단)에 들어갈 그림.
HALF_FIGURE = (6.0, 4.4)

#: 반 폭이되 **범례가 바깥에 붙는** 그림. 범례가 폭을 늘려 비율이 납작해지므로
#: 그만큼 미리 좁혀 잡는다 — 그러지 않으면 옆 그림보다 작게 앉는다.
HALF_FIGURE_WITH_LEGEND = (5.4, 4.6)

#: 좌우 2단의 넓은 쪽.
WIDE_FIGURE = (7.2, 4.4)

#: **폭을 다 쓰는 칸**의 그림 (38세션 4절). 5장이 이 크기로 굽는다.
#:
#: 기본 비율(9:3.6)로 구우면 슬라이드에서 높이가 먼저 차서 폭의 8할만 쓰고
#: 좌우가 통째로 빈다 — 칸의 가로세로에서 나온 값이라 여기 둔다.
FULL_FIGURE = (12.0, 3.4)

#: **범례가 아래에 붙는** 폭 전체 그림 (53세션 7-3). 4장의 기온 그래프다.
#:
#: 범례 줄이 그림 높이에 얹혀 비율이 0.29 → 0.36 으로 두꺼워졌고, 그만큼
#: 슬라이드에서 폭을 못 썼다 (9.4 → 7.8in). **범례가 먹는 만큼 미리 낮게 굽는다**
#: — :data:`HALF_FIGURE_WITH_LEGEND` 와 같은 이치다.
FULL_FIGURE_WITH_LEGEND = (12.0, 2.8)


def _picture(
    slide: Slide,
    png: bytes,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
) -> float:
    """그림을 상자 안에 **비율을 지켜** 앉힌다. 돌려주는 값은 그림의 아랫변이다.

    **모서리를 둥글리지 않는다** (36세션 3-5). 라운드를 준다면 0.08in 이상이라야
    하는데, 그만한 라운드는 이 크기의 차트에서 축 눈금을 갉아먹는다.
    """
    stream = io.BytesIO(png)
    picture = slide.shapes.add_picture(stream, Inches(left), Inches(top))
    ratio = picture.height / picture.width
    if width * ratio <= height:
        picture.width = Emu(int(Inches(width)))
        picture.height = Emu(int(Inches(width) * ratio))
    else:
        picture.height = Emu(int(Inches(height)))
        picture.width = Emu(int(Inches(height) / ratio))
    picture.left = Emu(int(Inches(left) + (Inches(width) - picture.width) / 2))
    picture.top = Emu(int(top * 914_400))
    return top + picture.height / 914_400


def _picture_block(
    slide: Slide,
    guide: DesignGuide,
    png: bytes,
    caption: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    """그림과 캡션을 **한 덩어리로 남는 높이 가운데** 앉힌다.

    차트는 가로로 길어서 폭이 먼저 차고, 그러면 아래에 빈 띠가 남는다. 위로
    붙여 두면 슬라이드가 덜 만들어진 것처럼 보인다 — 남는 여백을 위아래로
    나누면 **여백이 자리가 된다.**
    """
    # **넣어 보고 지우지 않는다.** ``add_picture`` 는 그림 파트를 꾸러미에 넣고
    # 관계를 맺는다 — 도형만 지우면 그 파트가 남아 파일이 두 배로 커진다.
    # 비율은 바이트에서 바로 읽는다.
    ratio = _aspect(png)
    room = height - _CAPTION_HEIGHT
    drawn_width = min(width, room / ratio) if ratio else width
    drawn_height = drawn_width * ratio
    offset = max(0.0, (room - drawn_height) / 2)
    _picture(
        slide,
        png,
        left=left + (width - drawn_width) / 2,
        top=top + offset,
        width=drawn_width,
        height=drawn_height,
    )
    # **캡션은 칸의 밑변에 붙인다.** 그림 바로 밑에 두면 좌우로 나란한 두 그림의
    # 캡션이 서로 다른 높이에 앉는다 — 그림 높이가 범례 유무로 갈리기 때문이다.
    _caption(
        slide,
        guide,
        caption,
        left=left,
        top=top + height - _CAPTION_HEIGHT,
        width=width,
        align=PP_ALIGN.CENTER,
    )


#: 표에 쓸 스타일 — **테두리도 띠도 없는 것** (36세션 3-4). 줄은 우리가 긋는다.
_PLAIN_TABLE_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"


def _table(
    slide: Slide,
    guide: DesignGuide,
    rows: Sequence[Sequence[str]],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    widths: Sequence[float] | None = None,
) -> None:
    """표형 슬라이드의 표 — **구분선 기반이고 카드가 아니다** (36세션 3-4).

    칸을 칠하지 않고 줄 밑에만 얇은 선을 둔다. 기본 표 스타일은 띠를 칠하므로
    **꺼서 넣는다** — 칠한 띠는 카드와 같은 무게를 갖고, 가이드가 금지한
    「같은 형태의 반복」이 표에서 되살아난다.

    표 자체가 이 슬라이드의 **시각 요소**다 (36세션 5절).
    """
    frame = slide.shapes.add_table(
        len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height)
    )
    table = frame.table
    properties = table._tbl.find(qn("a:tblPr"))
    properties.set("firstRow", "0")
    properties.set("bandRow", "0")
    style = properties.find(qn("a:tableStyleId"))
    if style is None:
        style = properties.makeelement(qn("a:tableStyleId"), {})
        properties.append(style)
    style.text = _PLAIN_TABLE_STYLE

    if widths is not None:
        # ``_ColumnCollection`` 은 Iterable 로 선언되어 있지 않아 zip 이 물린다.
        for index, share in enumerate(widths):
            table.columns[index].width = Emu(int(Inches(width) * share))

    scale = guide.type_scale
    colors = guide.colors
    for row_index, row in enumerate(rows):
        table.rows[row_index].height = Inches(guide.slide.slack(height / len(rows)))
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.fill.background()
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            run = paragraph.add_run()
            run.text = plain_text(str(value))
            head = row_index == 0
            _apply_font(
                run,
                guide,
                scale.caption if not head else scale.caption + 0.5,
                colors.muted if head else colors.ink,
                bold=head,
            )
            _cell_rule(cell, colors.rule if not head else colors.ink, guide.slide.rule_pt)


def _cell_rule(cell: object, color: str, width_pt: float) -> None:
    """칸 **아래 줄만** 긋는다. 격자를 두르면 표가 상자가 된다.

    **``a:lnB`` 는 ``a:tcPr`` 의 맨 앞에 와야 한다.** DrawingML 이 자식 차례를
    강제해서(테두리 → 채우기 순), 채우기 뒤에 붙이면 PowerPoint 가 조용히
    무시한다 — 줄이 안 그려지는데 파일은 멀쩡히 열린다.
    """
    properties = cell._tc.get_or_add_tcPr()  # type: ignore[attr-defined]
    existing = properties.find(qn("a:lnB"))
    if existing is not None:
        properties.remove(existing)
    element = properties.makeelement(
        qn("a:lnB"), {"w": str(int(width_pt * 12_700)), "cap": "flat", "cmpd": "sng"}
    )
    fill = element.makeelement(qn("a:solidFill"), {})
    srgb = fill.makeelement(qn("a:srgbClr"), {"val": color.lstrip("#").upper()})
    fill.append(srgb)
    element.append(fill)
    properties.insert(0, element)


def _stats(
    slide: Slide,
    guide: DesignGuide,
    items: Sequence[tuple[str, str]],
    *,
    left: float,
    top: float,
    width: float,
) -> float:
    """통계 강조 — **큰 수와 이름, 그 사이를 세로선이 가른다** (36세션 3-4).

    **카드로 만들지 않는다.** 칠한 상자를 늘어놓으면 슬라이드마다 같은 형태가
    되고, 모서리 색 띠는 가이드가 금지한다 (3-5). 돌려주는 값은 아랫변이다.
    """
    if not items:
        return top
    scale = guide.type_scale
    colors = guide.colors
    span = width / len(items)
    # **칸에 안 들어가면 글자를 줄인다.** 「28억 9,828만원」·「1,234,000원 (12개월
    # 환산 …)」 처럼 늘어지는 값이 있는데, 큰 글씨 그대로 두면 두 줄로 넘쳐 아래
    # 그림과 겹친다. 길이만 보고 정하면 칸 수에 따라 또 넘친다 — **폭으로 잰다.**
    values = [value for _label, value in items]
    inner = span - 0.24
    size = _fitting_size(values, span=inner, ladder=(scale.card_title, scale.body + 2, scale.body))
    # **두 줄이 될 값이 있으면 그만큼 높이를 잡는다** (38세션 4절). 예전에는 한
    # 줄 높이로 고정해 두어, 넘친 줄이 아래 그림의 윗여백을 먹고 지표 넷 가운데
    # 하나만 아래로 처져 보였다 — 역률·태양광·ESS 세 장에서 같은 모양이었다.
    lines = _fitting_lines(values, span=inner, size=size)
    block = (0.86 if size >= scale.card_title else 0.94) + (size / 72.0) * 1.5 * (lines - 1)
    for index, (label, value) in enumerate(items):
        x = left + span * index
        _text(
            slide,
            guide,
            [label],
            left=x + (0.16 if index else 0.0),
            top=top,
            width=span - 0.24,
            height=0.22,
            size=scale.caption,
            color=colors.muted,
        )
        _text(
            slide,
            guide,
            [value],
            left=x + (0.16 if index else 0.0),
            top=top + 0.26,
            width=inner,
            height=block - 0.32,
            size=size,
            color=colors.ink,
            bold=True,
        )
        if index:
            _vrule(slide, guide, left=x, top=top, length=block - 0.08, color=colors.rule)
    return top + block


# ===================================================================== 값 다듬기


def _won(value: float | None, *, reason: str | None = None) -> str:
    if value is None:
        return reason if reason is not None else _UNPRICED
    return money.won_short(value, reason=_UNPRICED)


def _payback(years: float | None, investment_won: float | None) -> str:
    if investment_won is None:
        return _UNPRICED
    if not investment_won:
        return IMMEDIATE
    return f"{years:,.1f}년" if years is not None else _UNPRICED


def _pct(value: float | None) -> str:
    return f"{value * 100:,.1f}%" if value is not None else "—"


#: 투자가 없는 수단의 투자비 칸 (53세션 1-5).
#:
#: **「0 원」 은 값이 아니라 없음이다.** 큰 글씨 자리에 0 이 서면 「계산해 보니
#: 0 원」 으로 읽히는데, 요금제 전환·계약전력 조정은 애초에 살 물건이 없다.
NO_INVESTMENT = "—"

#: 투자가 없는 수단의 회수기간 칸. **괄호를 뗀다** — 투자비 칸이 이미 없다고
#: 말하므로 「(투자 없음)」 이 같은 말을 한 번 더 한다.
IMMEDIATE = "즉시"


def _is_zero_won(text: str) -> bool:
    """「0원」 인가. **숫자를 되파싱하지 않는다** — 0 하나만 본다."""
    return text.replace(",", "").rstrip("원").strip() == "0"


def slide_investment(value: str) -> str:
    """슬라이드 투자비 칸 (53세션 1-5). 0 이면 :data:`NO_INVESTMENT` 다."""
    head = split_reason(value)[0]
    return NO_INVESTMENT if _is_zero_won(head) else head


def slide_payback(value: str) -> str:
    """슬라이드 회수기간 칸 (53세션 1-5). 투자가 없으면 :data:`IMMEDIATE` 다."""
    head = split_reason(value)[0]
    return IMMEDIATE if head.startswith(IMMEDIATE) else head


# ===================================================================== 장별


def _build_cover(
    slide: Slide, guide: DesignGuide, sections: DocumentSections, _spec: SlideSpec
) -> None:
    """표지 — **전체 배경 밴드가 다크다** (36세션 3-2).

    딥그린은 여기와 같은 **전체 배경**으로만 쓴다. 슬라이드 폭을 가로지르는
    색 바로 쓰면 3-5 가 금지한 그것이 된다.
    """
    geometry = guide.slide
    colors = guide.colors
    scale = guide.type_scale
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(colors.cover)

    # 표지의 시각 요소 — **작은 포인트 하나.** 코랄은 이런 자리에만 쓴다 (3-2).
    # 슬라이드를 가로지르지 않는 짧은 표식이라 3-5 의 「색 바」 와 다르다.
    left = geometry.margin_in + 0.5
    _rule(slide, guide, left=left, top=2.86, length=0.9, color=colors.coral, width_pt=3.0)
    bill = sections.bill
    _text(
        slide,
        guide,
        [DECK_TITLE],
        left=left,
        top=3.14,
        width=geometry.content_width_in - 1.0,
        height=1.0,
        size=scale.cover,
        color=colors.on_dark,
        bold=True,
    )
    _text(
        slide,
        guide,
        [
            sections.building,
            f"{bill.contract_label} {bill.voltage_label} 선택{bill.selection.option}",
            f"분석 기간 {bill.period_start:%Y-%m-%d} ~ {bill.period_end:%Y-%m-%d}"
            f" · 작성일 {sections.prepared:%Y-%m-%d}",
        ],
        left=left,
        top=4.52,
        width=geometry.content_width_in - 1.0,
        height=1.2,
        size=scale.body,
        color=colors.on_dark_muted,
    )


def _build_agenda(
    slide: Slide, guide: DesignGuide, sections: DocumentSections, _spec: SlideSpec
) -> None:
    """목차 — 리스트형. **번호와 구분선이 시각 요소다.**"""
    geometry = guide.slide
    colors = guide.colors
    scale = guide.type_scale
    top = _title(slide, guide, SLIDE_TITLES["agenda"])
    items = agenda_items(sections)
    step = min(0.62, (geometry.height_in - geometry.margin_in - top) / max(len(items), 1))
    for index, line in enumerate(items):
        y = top + step * index
        _text(
            slide,
            guide,
            [f"{index + 1:02d}"],
            left=geometry.margin_in,
            top=y + 0.04,
            width=0.7,
            height=0.32,
            size=scale.body,
            color=colors.coral,
            bold=True,
        )
        _text(
            slide,
            guide,
            [line],
            left=geometry.margin_in + 0.8,
            top=y,
            width=geometry.content_width_in - 0.8,
            height=0.36,
            size=scale.card_title,
            color=colors.ink,
        )
        _rule(
            slide,
            guide,
            left=geometry.margin_in,
            top=y + step - 0.08,
            length=geometry.content_width_in,
            color=colors.rule,
        )


def _build_building(
    slide: Slide, guide: DesignGuide, sections: DocumentSections, spec: SlideSpec
) -> None:
    """건물현황 및 계약정보 — **표형** (36세션 3-4). 카드가 아니다.

    **결측·정전이 결과에 어떻게 걸리는지 적는다** (39세션 6-2). 「972구간(2.8%)」
    만으로는 그 값을 어디까지 믿을지 알 수 없다 — 결측이 몰린 달과 정전 지속
    시간이 화면에는 있는데 슬라이드에는 건수뿐이었다.
    """
    geometry = guide.slide
    top = _title(slide, guide, spec.title)
    bill = sections.bill
    meta = sections.usage.meta
    quality = sections.diagnosis.quality if sections.diagnosis is not None else None
    # **해석 한 줄을 각주로 내렸다** (53세션 7-1). 이 장의 본체는 표다 —
    # 자료가 얼마나 성한지는 그 표를 **읽는 데 붙는 단서**이지 제목 다음에 와야
    # 할 결론이 아니다. 표가 그만큼 넓게 앉는다.
    note = narrative.building_lead(quality)
    rows = [
        ["항목", "내용"],
        ["건물명", sections.building],
        ["계약종별", f"{bill.contract_label} {bill.voltage_label} 선택{bill.selection.option}"],
        [
            "분석 기간",
            f"{bill.period_start:%Y-%m-%d} ~ {bill.period_end:%Y-%m-%d}"
            f" ({bill.period_days:.0f}일 · 기본요금 {bill.base_fee_months:.2f}개월분)",
        ],
        ["검침 간격", f"{meta.interval_minutes}분 · {meta.expected_rows:,}구간"],
        ["총 사용량", f"{meta.total_kwh / 1000:,.1f} MWh"],
        [
            "결측",
            f"{meta.missing_rows:,}구간 ({meta.missing_ratio:.1%}) — 보간하지 않고 뺐습니다",
        ],
        ["적용 요금표 시행일", f"{bill.effective_date}"],
        # **어느 기준 데이터로 계산했는가** (56세션 3절). 실물과 재현이 갈릴 때
        # 이 줄이 조건을 말한다 — 전문은 Excel 부록 「기준 데이터」 에 있다.
        ["기준 데이터", rules_basis_line()],
        ["작성일", f"{sections.prepared:%Y-%m-%d}"],
    ]
    if quality is not None:
        # **건수만으로는 규모를 알 수 없다** (30세션 1-3 · 39세션 6-2). 15분 하나가
        # 1건이고 며칠이 통째로 빈 것도 1건이라, 화면과 같이 지속 시간을 함께 적는다.
        hours = sum(event.duration_hours for event in quality.outages)
        outage = f"{len(quality.outages)}건"
        if quality.outages:
            outage += f" · 합 {hours:,.1f}시간 — 그 동안은 피크가 날 수 없어 편중 판정에서 뺐습니다"
        rows.insert(-2, ["정전 추정", outage])
    _table(
        slide,
        guide,
        rows,
        left=geometry.margin_in,
        top=top,
        width=geometry.content_width_in,
        height=min(
            _table_room_above(guide, top=top, note_top=_note_top(guide, note)),
            0.44 * len(rows),
        ),
        widths=(0.26, 0.74),
    )
    _note(slide, guide, note)


def _build_usage_pattern(
    slide: Slide, guide: DesignGuide, sections: DocumentSections, spec: SlideSpec
) -> None:
    """전력사용현황 및 부하패턴 — **두 장을 하나로 합쳤다** (38세션 2-1).

    **화면 1단계의 차례를 그대로 따른다.** 화면은 머릿수 지표 자리에 그림이 없고
    첫 그림이 「부하 패턴」 절에 나온다. 36세션의 덱은 그 구조와 어긋나 4장에
    사용량 그래프를 억지로 넣고, 5장이 같은 이야기를 지표만 바꿔 되풀이했다.

    **지표 넷을 고른 자리다.** 화면이 두 절에 걸쳐 낸 일곱 가운데 넷만 선다.

        분석 기간      3장 「건물현황 및 계약정보」 표에 이미 있다 — 뺀다
        최대수요       다음 장 「피크특성 (1/2)」 의 첫 지표다 — 그쪽으로 넘긴다
        요금적용전력   같은 이유로 넘긴다
        주말 부하 비율 「운영시간 외 부하 비중」 이 주말을 품고(주말은 전부 밖이다)
                       주말 이야기는 다음 장의 「상위 구간 주말 비중」 이 잇는다

    남는 넷이 **규모 하나(연간 사용량)와 모양 셋**이다.

    그림은 **일별 사용량에 일평균 기온을 겹친 한 장**이다 — 냉난방이 부하의
    얼마를 차지하는지가 태양광·ESS 판단을 가른다 (30세션 4절). 기온이 없으면
    사용량만 그린다 (화면과 같은 규칙).
    """
    geometry = guide.slide
    top = _title(slide, guide, spec.title)
    meta = sections.usage.meta
    diagnosis = sections.diagnosis
    if diagnosis is None:  # pragma: no cover - 진단 없이 부르지 않는다
        return
    pattern = diagnosis.pattern
    top = _lead(slide, guide, narrative.pattern_lead(pattern), top=top)
    # 1년치가 아닌 자료를 「연간」 이라 적으면 그 자체가 오독이다 (화면과 같다).
    span_label = "연간 사용량" if (meta.period_days or 0) >= 350 else "기간 사용량"
    bottom = _stats(
        slide,
        guide,
        [
            (span_label, f"{meta.total_kwh / 1000:,.0f} MWh"),
            ("부하율", _pct(pattern.load_factor)),
            ("기저부하 비율", _pct(pattern.base_load_ratio)),
            ("운영시간 외 부하 비중", _pct(pattern.off_hours_energy_share)),
        ],
        left=geometry.margin_in,
        top=top,
        width=geometry.content_width_in,
    )
    body = bottom + geometry.block_gap_in
    png, caption = _usage_figure(sections)
    note = narrative.glossary_note(GLOSSARY_KEYS["usage_pattern"], pattern)
    _picture_block(
        slide,
        guide,
        png,
        caption,
        left=geometry.margin_in,
        top=body,
        width=geometry.content_width_in,
        height=_note_top(guide, note) - body,
    )
    _note(slide, guide, note)


#: 기온을 곁들인 그림의 캡션과, 곁들이지 못했을 때의 캡션.
#:
#: **「결측일은 그리지 않았습니다」 를 뺐다** (53세션 7-2). 그 사실은 3장이
#: 표 한 줄과 각주로 이미 두 번 말한다 — 캡션은 **무엇을 그렸나**를 적는
#: 자리이지 전제를 되풀이하는 자리가 아니다.
_TEMPERATURE_CAPTION = "일별 사용량과 일평균 기온"
_USAGE_ONLY_CAPTION = "일별 사용량"


def _usage_figure(sections: DocumentSections) -> tuple[bytes, str]:
    """4장의 그림 하나. **기온이 없으면 사유를 캡션에 적고 사용량만 그린다.**

    빈 축을 남기지 않는 것이 화면의 규칙이고(30세션 4절), 슬라이드는 지역 입력
    없이도 한 장이 채워져야 한다.
    """
    temperature = sections.temperature
    if temperature is not None and len(temperature):
        try:
            png = figures.daily_temperature_png(
                sections.usage, temperature, size=FULL_FIGURE_WITH_LEGEND
            )
            return png, _TEMPERATURE_CAPTION
        except Exception:  # pragma: no cover - 그림 하나 때문에 덱을 잃지 않는다
            pass
    return figures.daily_usage_png(sections.usage, size=FULL_FIGURE), _USAGE_ONLY_CAPTION


def _peak_stats(sections: DocumentSections) -> list[tuple[str, str]]:
    """피크 지표 — **화면 「피크 특성」 절과 같은 갈림이다** (38세션 2-2).

    관측 최대와 요금적용 대상 최대가 같은 값이면 한 칸으로 접는다 — 두 칸을
    나란히 두면 둘이 다른 개념인 줄 알고 차이를 찾게 된다 (13세션). 야간
    피크형에서만 갈린다.

    **화면은 칸이 셋이라 갈릴 때 정오 비중을 밀어냈다.** 슬라이드는 넷이므로
    밀어낼 것이 없다 — 정오 비중이 태양광 판정의 근거라 어느 갈래에서도 남는다.
    """
    diagnosis = sections.diagnosis
    assert diagnosis is not None
    peak = diagnosis.peak
    split = peak.billing_demand_kw < peak.peak_kw * 0.99
    items: list[tuple[str, str]] = []
    if split:
        items.append(("관측 최대수요", f"{peak.peak_kw:,.0f} kW"))
        items.append(("요금적용전력", f"{peak.billing_demand_kw:,.0f} kW"))
    else:
        items.append(("최대수요 = 요금적용전력", f"{peak.peak_kw:,.0f} kW"))
    items.append(("상위 구간 정오 비중", _pct(diagnosis.summary.pv_midday_share)))
    items.append(
        (
            "상위 구간 주말 비중",
            _pct(peak.weekend_slots / peak.top_n if peak.top_n else None),
        )
    )
    return items


def _build_peak_summary(
    slide: Slide, guide: DesignGuide, sections: DocumentSections, spec: SlideSpec
) -> None:
    """피크특성 (1/2) — **지표 셋과 월별 최대수요** (38세션 2-2).

    **해석 한 줄이 그림과 같은 것을 말한다** (53세션 4-3). 39세션은 여기에 정오
    비중 판정을 적었는데 그림은 월별 최대수요였다 — 문장과 그림이 다른 것을
    말하면 읽는 사람이 둘을 잇지 못한다. 판정은 그 근거 그림이 있는 6장으로
    옮겼다.
    """
    geometry = guide.slide
    top = _title(slide, guide, spec.title)
    diagnosis = sections.diagnosis
    if diagnosis is None:  # pragma: no cover
        return
    quality = diagnosis.quality
    top = _lead(
        slide,
        guide,
        narrative.peak_month_lead(diagnosis, quality, sections.tariff_table),
        top=top,
    )
    bottom = _stats(
        slide,
        guide,
        _peak_stats(sections),
        left=geometry.margin_in,
        top=top,
        width=geometry.content_width_in,
    )
    body = bottom + geometry.block_gap_in
    note = narrative.glossary_note(GLOSSARY_KEYS["peak_summary"], diagnosis.pattern)
    _picture_block(
        slide,
        guide,
        figures.monthly_peak_png(diagnosis.peak, size=FULL_FIGURE),
        "월별 최대수요 — 붉은 점선이 기본요금을 매기는 요금적용전력입니다.",
        left=geometry.margin_in,
        top=body,
        width=geometry.content_width_in,
        height=_note_top(guide, note) - body,
    )
    _note(slide, guide, note)


def _build_peak_detail(
    slide: Slide, guide: DesignGuide, sections: DocumentSections, spec: SlideSpec
) -> None:
    """피크특성 (2/2) — **그림 둘을 좌우로.** 시간대별 평균 부하와 상위 구간 시각.

    **그림 크기를 칸에 맞춰 부른다** (:data:`HALF_FIGURE`). Word 의 가로 긴
    비율을 그대로 반 칸에 넣으면 눈금이 겹치고 위아래가 남는다.
    """
    geometry = guide.slide
    top = _title(slide, guide, spec.title)
    diagnosis = sections.diagnosis
    if diagnosis is None:  # pragma: no cover
        return
    top = _lead(slide, guide, narrative.peak_detail_lead(diagnosis), top=top)
    peak = diagnosis.peak
    gap = geometry.block_gap_in
    half = (geometry.content_width_in - gap) / 2
    height = _note_top(guide) - top
    # **문장이 말하는 그림이 왼쪽이다** (53세션 4-5). 해석 한 줄이 「상위 ○○구간」
    # 을 말하는데 그 그림이 오른쪽에 있어, 읽는 눈이 문장에서 오른쪽으로 건너뛴
    # 뒤 다시 왼쪽으로 돌아와야 했다.
    for index, (png, caption) in enumerate(
        (
            (
                figures.top_hour_png(peak, size=HALF_FIGURE_WITH_LEGEND),
                f"최대수요 상위 {peak.top_n}구간이 발생한 시각",
            ),
            (
                figures.hourly_profile_png(peak, size=HALF_FIGURE),
                "하루 24시간의 평균 부하 모양",
            ),
        )
    ):
        _picture_block(
            slide,
            guide,
            png,
            caption,
            left=geometry.margin_in + (half + gap) * index,
            top=top,
            width=half,
            height=height,
        )


def _build_structure(
    slide: Slide, guide: DesignGuide, sections: DocumentSections, spec: SlideSpec
) -> None:
    """현재 요금 구조 — 월별 요금 막대와 **계시별 도넛 2×2** (36세션 2·4절)."""
    geometry = guide.slide
    top = _title(slide, guide, spec.title)
    diagnosis = sections.diagnosis
    structure = diagnosis.structure if diagnosis is not None else None
    if structure is None:  # pragma: no cover - 계약 정보가 없으면 부르지 않는다
        return
    top = _lead(slide, guide, narrative.structure_lead(structure), top=top)
    pattern = diagnosis.pattern if diagnosis is not None else None
    note = narrative.glossary_note(GLOSSARY_KEYS["structure"], pattern)
    gap = geometry.block_gap_in
    left_width = geometry.content_width_in * 0.56
    right_width = geometry.content_width_in - left_width - gap
    right_left = geometry.margin_in + left_width + gap
    bottom_y = _note_top(guide, note)
    height = bottom_y - top

    base_won = structure.base_won + structure.bill.total_power_factor_won
    chart_top = (
        _stats(
            slide,
            guide,
            [
                ("기본요금", _won(base_won)),
                ("전력량요금", _won(structure.energy_won)),
                (
                    "기본요금 비중",
                    _pct(base_won / structure.total_won if structure.total_won else None),
                ),
            ],
            left=geometry.margin_in,
            top=top,
            width=left_width,
        )
        + gap
    )
    _picture_block(
        slide,
        guide,
        figures.monthly_charge_png(structure, size=WIDE_FIGURE),
        "월별 요금 구성 — 기본요금은 직전 12개월 최대수요로 정해져 매달 같습니다.",
        left=geometry.margin_in,
        top=chart_top,
        width=left_width,
        height=bottom_y - chart_top,
    )
    _picture_block(
        slide,
        guide,
        figures.band_donut_grid_png(structure, season_pairs(structure)),
        "계절별 계시 시간대 사용량 구성",
        left=right_left,
        top=top,
        width=right_width,
        height=height,
    )
    _note(slide, guide, note)


def _build_measure_summary(
    slide: Slide, guide: DesignGuide, sections: DocumentSections, spec: SlideSpec
) -> None:
    """개선안별 요약 — **표형.** 켠 수단을 한 자리에 모은다.

    **투자 없이 가능한 절감액을 여기서 낸다** (39세션 6-1). 화면 1단계 맨 위의
    숫자인데 슬라이드에 없었다 — 고객이 가장 먼저 보고 싶은 값이고, 개선 수단이
    처음 나오는 이 장이 그 자리다.

    **금액은 12개월 환산 한 값이다** (2-2). 같은 값을 괄호로 한 번 더 적지 않고,
    환산 기준이라는 사실은 아래 각주가 한 번 말한다.
    """
    geometry = guide.slide
    top = _title(slide, guide, spec.title)
    diagnosis = sections.diagnosis
    if diagnosis is not None:
        saving = _won(diagnosis.summary.no_investment_saving_won)
        top = _lead(slide, guide, narrative.measure_summary_lead(diagnosis, saving), top=top)
        top = (
            _stats(
                slide,
                guide,
                [
                    ("투자 없이 가능한 절감액", saving),
                    ("검토한 수단", f"{len(sections.measures)}개"),
                ],
                left=geometry.margin_in,
                top=top,
                width=geometry.content_width_in / 2,
            )
            + geometry.block_gap_in
        )
    # **확실성 열을 뺐다** (53세션 1-4). 무엇에 대한 등급인지 이름에 없어
    # 화면에서 28세션에 걷어냈는데 산출물에만 남아 있었다. 계산은 그대로다 —
    # :class:`~kwise.measures.Certainty` 도 :attr:`MeasureEntry.certainty` 도 산다.
    rows = [["개선 수단", "절감액", "투자비", "회수기간"]]
    # **절 번호를 뗀다** (38세션 1-3). 장 제목과 같은 이름이어야 표에서 고른
    # 줄을 뒤에서 찾을 수 있다.
    rows.extend(
        [
            measure_slide_title(entry),
            split_reason(entry.slide_saving)[0],
            slide_investment(entry.investment),
            slide_payback(entry.payback),
        ]
        for entry in sections.measures
    )
    if len(rows) == 1:
        rows.append(["검토한 수단이 없습니다", "—", "—", "—"])
    # **각주가 실제로 먹는 자리를 재서 그 위까지만 쓴다** (53세션 9절).
    # 아래 각주는 세 줄로 흐르는데 예약은 두 줄어치였다 — 표 마지막 줄이
    # 「※ 회수기간 = …」 위에 겹쳐 있었다.
    notes = (
        narrative.glossary_note(GLOSSARY_KEYS["measure_summary"]),
        f"{ANNUAL_BASIS_NOTE} {NOT_INCLUDED_NOTICE} {TRUNCATION_FOOTNOTE}",
    )
    _table(
        slide,
        guide,
        rows,
        left=geometry.margin_in,
        top=top,
        width=geometry.content_width_in,
        height=min(
            _table_room_above(guide, top=top, note_top=_note_top(guide, *notes)),
            0.5 * len(rows),
        ),
        widths=(0.26, 0.34, 0.24, 0.16),
    )
    _note(slide, guide, *notes)


#: 금액 기준을 **한 번만** 적는다 (39세션 2-2). 값마다 「(12개월 환산 ○○원)」 을
#: 붙이면 같은 값이 두 번 나오고 지표 칸이 두 줄로 흐른다.
ANNUAL_BASIS_NOTE = "금액은 12개월 환산 기준입니다."

#: 수단 한 장에 실을 주의사항 개수. **슬라이드는 읽는 자리가 아니라 보는 자리다** —
#: 전문은 Word 3장에 그대로 있다.
CAUTION_LIMIT = 3


def _cautions(entry: MeasureEntry) -> tuple[str, ...]:
    """주의사항을 **겹치지 않게** 추리고 개수를 자른다.

    같은 문장이 두 번 실리는 경우가 있다 — 카드가 낸 주의와 안내가 낸 근거가
    같은 말일 때다. 문서에서는 눈에 덜 띄지만 슬라이드에서는 한 화면에 나란히
    놓여 바로 보인다.

    **결론과 무관하면 싣지 않는다** (39세션 4-2). 하향 여지가 없는 계약전력에
    하향 시 주의사항 셋을, 잉여가 0인 수단에 판매 자격요건을 다는 것은 하지도
    않을 일을 조심하라는 말이라 결론보다 길어진다. 그 자리는
    :attr:`MeasureEntry.facts` 가 받는다 — **왜 여지가 없는지 보이는 숫자**다.
    """
    if not entry.actionable:
        return ()
    seen: dict[str, None] = {}
    for line in entry.cautions:
        seen.setdefault(line.strip(), None)
    return tuple(seen)[:CAUTION_LIMIT]


def _trim_repeat(reason: str, head: str) -> str:
    """사유 꼬리에 머리말이 되풀이되면 뗀다 (59세션 8절).

    「하한 규정 미확인 — 금액 미산출」 처럼 사유 자체가 「…미산출」 로 끝나는
    자리가 있다. 앞에 「절감액 미산출」 을 붙이면 한 줄에 같은 말이 두 번 선다.
    **떼는 것은 표시뿐이고 계산이 낸 사유 문자열은 그대로 둔다.**
    """
    tail = reason.rsplit(_REASON_MARK, 1)[-1].strip()
    if _REASON_MARK in reason and tail.endswith(head):
        return reason.rsplit(_REASON_MARK, 1)[0].strip()
    return reason


def _measure_note(entry: MeasureEntry) -> str:
    """수단 장 맨 아래 작은 글씨 (39세션 2-1).

    **지표 칸에서 뗀 미산출 사유가 여기로 온다.** 값 자리에는 「미산출」 만 두되,
    무엇을 넣으면 값이 나오는지는 고객에게 필요한 정보라 지우지 않는다.
    """
    parts: list[str] = []
    for label, value in (
        ("절감액", entry.slide_saving),
        ("투자비", entry.investment),
        ("회수기간", entry.payback),
    ):
        head, reason = split_reason(value)
        if reason:
            # **머리말을 그대로 옮긴다** (53세션 9절). 「미산출」 로 못박아 두어
            # 「기본요금 변화없음」 인 값이 각주에서 「절감액 **미산출**」 이 됐다.
            #
            # **「미산출」 이 두 번 나오지 않게 한다.** 태양광 투자비 사유가
            # 「미산출 — 태양광 설치 단가 미입력」 꼴이라, 앞에 라벨을 붙이면
            # 「투자비 미산출 — 미산출 — 태양광 …」 이 됐다.
            #
            # **꼬리에도 붙는다** (59세션 8절). 계약전력의 사유는 「하한 규정
            # 미확인 — 금액 미산출」 이라 앞이 아니라 **뒤**에 같은 말이 있었다
            # — 「절감액 미산출 — 하한 규정 미확인 — 금액 미산출」.
            tail = _trim_repeat(split_reason(reason)[1] or reason, head)
            parts.append(f"{label} {head} — {tail}")
    return " · ".join(parts)


def _measure_layout(entry: MeasureEntry) -> str:
    """수단 장의 형태. **자리표가 쥐고 있어야 시험이 셀 수 있다** (38세션 3절).

    46세션에 ``stat_table_chart`` 가 붙었다 — ESS 가 회수기간 곡선 대신 목표별
    사양 표를 싣는다. 열이 아홉이라 반 칸에 못 넣으므로 **표가 폭 전체로 위,
    그림이 아래**다.
    """
    if entry.spec_table:
        return "stat_table_chart"
    if entry.facts_first and entry.facts:
        return "fact_stat_chart"
    return "chart_pair" if len(entry.slide_figures) > 1 else "stat_chart"


def _measure_pictures(
    slide: Slide,
    guide: DesignGuide,
    drawings: Sequence[MeasureFigure],
    *,
    top: float,
    height: float,
) -> None:
    """수단 장의 그림들 — **둘이면 좌우로, 하나면 폭 전체로** (38세션 3절).

    화면이 그림 둘로 답하는 물음을 슬라이드가 하나로 줄이면, 남은 하나가 무엇의
    근거인지 알 수 없어진다 — ESS 의 「목표 5,170 kW 는 어디서 나왔나」 가 그
    자리였다. 셋 이상은 두지 않는다: 반 칸 아래로 내려가면 축 눈금이 뭉개진다.
    """
    geometry = guide.slide
    gap = geometry.block_gap_in
    count = min(len(drawings), 2)
    width = (geometry.content_width_in - gap * (count - 1)) / count
    for index, item in enumerate(drawings[:count]):
        _picture_block(
            slide,
            guide,
            item.png,
            item.caption,
            left=geometry.margin_in + (width + gap) * index,
            top=top,
            width=width,
            height=height,
        )


#: 사양 표 한 줄이 실제로 차지하는 높이 (in) — **머리글도 한 줄이다** (53세션 2절).
#:
#: 46세션은 「줄당 0.4in」 을 상한으로만 쓰고 실제 배분은 ``height * 0.55`` 로
#: 잡았다. 50세션에 격자가 붙어 표가 **다섯 줄에서 여섯 줄**이 되자 그 비율이
#: 그대로 표를 키워, 남은 자리에 그림이 **0.85 × 0.41in** 로 우겨 넣어졌다 —
#: 손톱만 한 차트다. **비율이 아니라 줄 수로 잡는다.**
_SPEC_ROW_HEIGHT = 0.30

#: 그림을 뺐을 때 사양 표가 늘어날 수 있는 **줄당 상한** (in) (53세션 9절).
#: 남는 자리를 표가 다 쓰되, 여섯 줄짜리 표가 슬라이드를 세로로 가르지 않게 한다.
_SPEC_ROW_MAX = 0.55

#: 사양 표와 그 아래 참고 줄 사이의 숨 (in).
#:
#: **0 이면 캡션이 마지막 줄에 얹힌다** (53세션 9절). 표 높이는 PowerPoint 에서
#: 최소값일 뿐이라 실제로는 조금 더 벌어지는데, 캡션을 표 밑변에 딱 붙여 두어
#: 그 차이만큼 겹쳤다 — 실물을 렌더해서야 보였다.
_SPEC_CAPTION_GAP = 0.10

#: 표 아래 그림에 남겨야 하는 최소 높이 (in). 캡션 0.32 + 그림 0.7 남짓이다.
#:
#: **이보다 좁으면 그림을 싣지 않는다.** 뭉갠 차트는 없는 것만 못하고, 자리를
#: 비우면 표가 그 자리를 받아 마지막 줄이 캡션에 닿지 않는다.
#:
#: **53세션 9절에 1.02 → 1.35 로 올렸다.** 실물을 렌더해 보니 1.0in 짜리 칸에
#: 앉은 ESS 대표일 곡선은 축 눈금이 읽히지 않는 손톱만 한 그림이었다 — 그
#: 자리를 표에 주면 아홉 열이 넉넉히 앉는다. 같은 그림은 화면과 Word 에 있다.
_MIN_FIGURE_BLOCK = 1.35

#: 사양 표 아래 참고 **한 줄**이 먹는 높이 (in). 줄 수만큼 곱한다.
_SPEC_CAPTION_HEIGHT = 0.24

#: 수단 장 본문 아래에 남기는 숨 (in) (53세션 2절).
#:
#: 46세션까지 0.32 였는데 :func:`_picture_block` 이 **이미 캡션 높이를 제 안에서
#: 빼고 있어** 같은 몫을 두 번 뺐다 — 모든 수단 장이 아래를 0.4in 씩 비워 두고
#: 그림만 그만큼 눌려 있었다. 6장(피크특성 2/2)은 이 계산을 안 타서 캡션이
#: 6.68in 에 앉는데 잘리지 않는다 — 그것이 이 값의 상한을 말해 준다.
_BODY_TAIL = 0.1

#: 사양 표의 열 폭 비율 (53세션 2절). **아홉 열을 고르게 나누지 않는다.**
#:
#: 고르게 나누면 한 칸이 1.37in 인데, 「5,220~5,240 kW」 는 11.5pt 로 1.34in 이라
#: 여백을 빼면 **두 줄로 흐른다.** 그렇게 흐른 줄이 표를 1.5배로 키워 아래 그림
#: 자리를 먹었다. **긴 칸(목표·표식)에 폭을 몰고 숫자 칸을 좁힌다.**
#:
#: 차례는 :data:`~kwise.report.frames.ESS_SPEC_HEADER` 와 같다.
SPEC_TABLE_WIDTHS: tuple[float, ...] = (
    0.150,  # 목표
    0.075,  # 저감량
    0.075,  # 출력
    0.085,  # 용량
    0.095,  # 방전시간
    0.110,  # 투자비
    0.115,  # 연간 절감액
    0.095,  # 회수기간
    0.200,  # 표식
)


def _spec_lines(rows: Sequence[Sequence[str]], *, width: float, size: float) -> int:
    """사양 표가 실제로 차지할 **줄 수** (53세션 2절).

    줄 수를 행 수로 세면 안 된다 — 「목표 미달 (실제 264 kW)」 같은 표식이 칸
    폭을 넘어 **두 줄로 흐른다.** 행 수로만 잡아 두면 그만큼 표가 아래를 덮는다.
    열 폭은 고르게 나누므로(``_table`` 이 ``widths`` 를 받지 않는다) 한 칸의
    폭을 한 번만 재면 된다.
    """
    if not rows:
        return 0
    shares = _spec_widths(len(rows[0]))
    total = 0
    for row in rows:
        total += max(
            _fitting_lines([value], span=width * share - 0.12, size=size)
            for value, share in zip(row, shares, strict=True)
        )
    return total


def _spec_widths(columns: int) -> tuple[float, ...]:
    """사양 표의 열 폭 비율. **열 수가 다르면 고르게 나눈다.**"""
    if columns == len(SPEC_TABLE_WIDTHS):
        return SPEC_TABLE_WIDTHS
    return tuple(1.0 / columns for _ in range(columns))


def _spec_block(
    slide: Slide,
    guide: DesignGuide,
    entry: MeasureEntry,
    drawings: tuple[MeasureFigure, ...],
    *,
    top: float,
    height: float,
) -> tuple[float, float, tuple[MeasureFigure, ...]]:
    """목표별 사양 표와 그 아래 참고 한 줄 (46세션 · 53세션 2절에 배분을 다시 잡았다).

    **표가 먼저다.** 「이 목표는 어디서 나왔나」 에 답하는 자리라 그림보다 위에
    온다. 표에 줄 수만큼 주고 **남는 것이 그림 몫**이다 — 46세션의 ``height * 0.55``
    는 줄이 늘어도 같은 비율을 먹어, 격자가 붙은 뒤로 그림을 뭉갰다.

    Returns:
        (그림이 시작할 y, 그림에 남은 높이, 실을 그림들).
        **자리가 :data:`_MIN_FIGURE_BLOCK` 에 못 미치면 그림을 뺀다** — 그만큼
        표가 넓게 앉는다.
    """
    geometry = guide.slide
    gap = geometry.block_gap_in
    # **표식의 뜻도 표 바로 아래가 자리다** (53세션 4-13). 슬라이드 맨 아래
    # 각주로 내리면 :func:`_note_top` 이 그만큼 예약해 **그림이
    # 자리를 잃는다** — 표를 읽는 데 바로 쓰이는 글이라 표에 붙인다.
    lines = [mark_note(line) for line in (entry.spec_caption, entry.spec_note) if line]
    caption_height = _SPEC_CAPTION_HEIGHT * len(lines)
    wanted = _SPEC_ROW_HEIGHT * _spec_lines(
        entry.spec_table,
        width=geometry.content_width_in,
        size=guide.type_scale.caption,
    )
    # **그림이 들어갈 자리가 남는가**로 갈린다. 남지 않으면 그림을 빼고 표가
    # 남은 높이를 다 쓴다 — 상한은 그 자리다.
    gap = gap / 2 if lines else gap
    room = height - wanted - caption_height - _SPEC_CAPTION_GAP - gap
    if room < _MIN_FIGURE_BLOCK:
        # **그림을 빼고 그 자리를 표에 준다.** 아래를 비워 두면 슬라이드가 덜
        # 만들어진 것처럼 보이고, 줄이 벌어진 표가 아홉 열을 훨씬 잘 읽힌다.
        drawings = ()
        roomy = height - caption_height - _SPEC_CAPTION_GAP
        wanted = min(max(wanted, roomy), _SPEC_ROW_MAX * len(entry.spec_table))
    _table(
        slide,
        guide,
        entry.spec_table,
        left=geometry.margin_in,
        top=top,
        width=geometry.content_width_in,
        # ``_table`` 이 줄마다 여유를 더하므로 **그만큼 미리 나눠** 넘긴다 —
        # 그러지 않으면 실제 표가 잡은 자리보다 커져 아래를 덮는다.
        height=wanted / (1.0 + geometry.text_slack),
        widths=_spec_widths(len(entry.spec_table[0])),
    )
    used = wanted
    if lines:
        _text(
            slide,
            guide,
            lines,
            left=geometry.margin_in,
            top=top + used + _SPEC_CAPTION_GAP,
            width=geometry.content_width_in,
            height=caption_height,
            size=guide.type_scale.caption,
            color=guide.colors.muted,
            spacing=1.0,
        )
        used += caption_height + _SPEC_CAPTION_GAP
    # **표와 그림 사이는 반 칸이다.** 참고 한 줄이 이미 둘을 가르므로 온 칸을
    # 두면 그만큼 그림이 눌린다.
    return top + used + gap, height - used - gap, drawings


def _build_measure(
    slide: Slide, guide: DesignGuide, sections: DocumentSections, spec: SlideSpec
) -> None:
    """수단 한 장 — **차트 + 지표**.

    **결론 한 줄이 먼저다.** 그림이 없는 수단은 표 대신 주의사항이 자리를 채운다 —
    빈 그림 자리를 남기지 않는다.

    **그림이 둘이면 좌우로 나눈다** (38세션 3절). 역률·태양광·ESS 가 그렇다.
    """
    geometry = guide.slide
    colors = guide.colors
    assert spec.measure is not None
    entry: MeasureEntry = sections.measures[spec.measure]
    top = _title(slide, guide, measure_slide_title(entry))
    # **결론도 줄 수를 세어 자리를 잡는다** (39세션 7절). 0.4in 로 못박아 두어
    # 두 줄짜리 결론이 아래 지표 라벨과 겹쳤다 — 경제성DR 장이 그랬다.
    size = guide.type_scale.body
    lead = _wrapped_height(entry.conclusion, guide, width=geometry.content_width_in, size=size)
    _text(
        slide,
        guide,
        [entry.conclusion],
        left=geometry.margin_in,
        top=top,
        width=geometry.content_width_in,
        height=lead,
        size=size,
        color=colors.ink,
        bold=True,
    )
    stats_top = top + lead
    # **근거가 결과보다 먼저 오는 장이 있다** (53세션 6-1). 계약전력 조정이
    # 그렇다 — 여유가 얼마나 있는지를 보고 나서야 「얼마로 낮출까」 를 읽는다.
    if entry.facts_first and entry.facts:
        stats_top = (
            _stats(
                slide,
                guide,
                list(entry.facts),
                left=geometry.margin_in,
                top=stats_top,
                width=geometry.content_width_in,
            )
            + geometry.block_gap_in * 0.5
        )
    bottom = _stats(
        slide,
        guide,
        [
            ("절감액", split_reason(entry.slide_saving)[0]),
            ("투자비", slide_investment(entry.investment)),
            ("회수기간", slide_payback(entry.payback)),
        ],
        left=geometry.margin_in,
        top=stats_top,
        width=geometry.content_width_in,
    )
    # **용어 풀이가 먼저다** (39세션 1-2 · 53세션 8-1). 표·지표를 읽는 데 바로
    # 쓰이고, 미산출 사유는 그 아래에서 받는다.
    terms_note = narrative.glossary_note(GLOSSARY_KEYS.get(spec.key, ()))
    note = _measure_note(entry)
    # **장이 따로 적는 줄은 제 줄에 선다** (59세션 14절). 미산출 사유 뒤에
    # 「·」 로 이어 붙이면 「역률 영향 반영 시 279,249,000원」 이 또 하나의
    # 미산출 사유처럼 읽힌다 — 다른 종류의 말이므로 ※ 를 따로 단다.
    extra = entry.slide_note
    body = bottom + geometry.block_gap_in
    height = _note_top(guide, terms_note, note, extra) - body - _BODY_TAIL
    drawings = entry.slide_figures
    crowded = False
    if entry.spec_table:
        body, height, drawings = _spec_block(
            slide, guide, entry, drawings, top=body, height=height
        )
        # **표가 자리를 다 썼으면 그것으로 끝이다** (53세션 2절·4-13). 남은
        # 틈에 주의사항 표를 밀어 넣으면 결론이 이미 한 말을 되풀이하면서
        # 줄이 눌린다 — ESS 가 그랬다.
        crowded = bool(entry.slide_figures) and not drawings
    if drawings and height > _MIN_FIGURE_BLOCK:
        _measure_pictures(slide, guide, drawings, top=body, height=height)
        _note(slide, guide, terms_note, note, extra)
        return
    # **여지가 없는 수단은 그 사실을 숫자로 보인다** (39세션 4-2·4-3). 실행
    # 주의사항 대신 「왜 없는지」 를 세우는 자리다.
    if crowded:
        _note(slide, guide, terms_note, note, extra)
        return
    if not entry.actionable and entry.facts and not entry.facts_first:
        _stats(
            slide,
            guide,
            list(entry.facts),
            left=geometry.margin_in,
            top=body + max(0.0, (height - 0.94) / 2),
            width=geometry.content_width_in,
        )
        _note(slide, guide, terms_note, note, extra)
        return
    rows = [["주의사항"], *[[line] for line in _cautions(entry)]]
    if len(rows) == 1:
        rows.append(["—"])
    table_height = min(height, 0.7 * len(rows))
    # **남는 높이 가운데 앉힌다** — 그림 덩어리와 같은 규약이다 (36세션 3-3).
    # 주의사항이 한 줄뿐인 장(잉여 0)은 표가 짧아, 위로 붙여 두면 아래 절반이
    # 통째로 비어 슬라이드가 덜 만들어진 것처럼 보인다.
    _table(
        slide,
        guide,
        rows,
        left=geometry.margin_in,
        top=body + max(0.0, (height - table_height) / 2),
        width=geometry.content_width_in,
        height=table_height,
    )
    _note(slide, guide, terms_note, note, extra)


def _build_surplus(
    slide: Slide, guide: DesignGuide, sections: DocumentSections, spec: SlideSpec
) -> None:
    """잉여 활용 — **좌우로 그림과 표** (53세션 3절).

    41세션에 잉여가 개선안에서 빠지면서 이 장이 사라졌는데, **잉여가 실제로
    나면 보여 줄 자리가 없어졌다** — 소형 사무빌딩의 연 23,416 kWh 가 태양광
    장의 각주 한 줄로만 나왔다.

    **잉여가 0 이면 이 장을 만들지 않는다** (:func:`slide_specs`). 대형 샘플은
    전량 자가소비라 그렇다 — 빈 장에 「없습니다」 라고 적을 자리가 아니다.
    """
    geometry = guide.slide
    page = sections.surplus
    if page is None:  # pragma: no cover - 자리표가 없으면 부르지 않는다
        return
    top = _title(slide, guide, spec.title)
    top = _lead(slide, guide, page.lead, top=top)
    top = (
        _stats(
            slide,
            guide,
            list(page.facts),
            left=geometry.margin_in,
            top=top,
            width=geometry.content_width_in,
        )
        + geometry.block_gap_in
    )
    gap = geometry.block_gap_in
    half = (geometry.content_width_in - gap) / 2
    right_left = geometry.margin_in + half + gap
    bottom = _note_top(guide, page.note)
    if page.figure is not None:
        _picture_block(
            slide,
            guide,
            page.figure,
            page.figure_caption,
            left=geometry.margin_in,
            top=top,
            width=half,
            height=bottom - top,
        )
    _table(
        slide,
        guide,
        page.scenario_rows,
        left=right_left,
        top=top,
        width=half,
        height=min(bottom - top, 0.46 * len(page.scenario_rows)),
        widths=(0.34, 0.26, 0.40),
    )
    _note(slide, guide, page.note)


def _build_combination(
    slide: Slide, guide: DesignGuide, sections: DocumentSections, spec: SlideSpec
) -> None:
    """조합구성 및 합산효과 — **비교형.** 좌우 2단에 중앙 얇은 구분선 (36세션 3-4).

    왼쪽이 「무엇을 묶었나」, 오른쪽이 「묶으면 얼마인가」다. **조합마다 요금을
    다시 계산한 값이며 수단별 절감액의 단순 합이 아니다.**
    """
    geometry = guide.slide
    colors = guide.colors
    scale = guide.type_scale
    top = _title(slide, guide, spec.title)
    comparison = sections.comparison
    top = _lead(slide, guide, narrative.combination_lead(comparison), top=top)
    gap = geometry.block_gap_in
    half = (geometry.content_width_in - gap) / 2
    right_left = geometry.margin_in + half + gap
    _vrule(
        slide,
        guide,
        left=geometry.margin_in + half + gap / 2,
        top=top,
        length=geometry.height_in - geometry.margin_in - top - 0.1,
        color=colors.rule,
    )

    if comparison is None:
        _text(
            slide,
            guide,
            ["조합 비교는 계약 정보와 수단이 있어야 산출됩니다."],
            left=geometry.margin_in,
            top=top,
            width=half,
            height=0.4,
            size=scale.body,
            color=colors.ink,
        )
        rows = [["조합", "절감액"], ["—", "—"]]
        _table(
            slide,
            guide,
            rows,
            left=right_left,
            top=top,
            width=half,
            height=1.0,
            widths=(0.6, 0.4),
        )
        return

    best = comparison.best
    _text(
        slide,
        guide,
        ["가장 유리한 조합"],
        left=geometry.margin_in,
        top=top,
        width=half,
        height=0.34,
        size=scale.card_title,
        color=colors.ink,
        bold=True,
    )
    # **고른 수단을 순서대로 이어 적는다** (39세션 3-3). 조합 이름은 「+ ESS 목표
    # 5,170 kW」 처럼 **직전 조합에 무엇을 더했는가**를 적는 것이라, 그것만 떼어
    # 놓으면 앞의 수단들이 보이지 않았다. 목록을 여기서 다시 만들지 않는다 —
    # 조합 재계산이 본 것과 갈라진다.
    baseline = comparison.combinations[0].spec.selection if comparison.combinations else None
    _text(
        slide,
        guide,
        [best.spec.composition(baseline)],
        left=geometry.margin_in,
        top=top + 0.42,
        width=half,
        height=0.62,
        size=scale.body,
        color=colors.ink,
    )
    table_top = top + 0.42 + 0.62 + gap
    rows = [["조합", "절감액", "회수기간"]]
    rows.extend(
        [item.name, _won(item.saving_won), _payback(item.payback_years, item.investment_won)]
        for item in comparison.combinations
    )
    _table(
        slide,
        guide,
        rows,
        left=geometry.margin_in,
        top=table_top,
        width=half,
        height=min(
            geometry.height_in - geometry.margin_in - table_top - 0.1,
            0.36 * len(rows),
        ),
        widths=(0.46, 0.3, 0.24),
    )

    _text(
        slide,
        guide,
        ["합산효과"],
        left=right_left,
        top=top,
        width=half,
        height=0.34,
        size=scale.card_title,
        color=colors.ink,
        bold=True,
    )
    stats_bottom = _stats(
        slide,
        guide,
        [
            ("총 절감액", _won(best.saving_won)),
            ("투자비", slide_investment(_won(best.investment_won))),
            ("회수기간", _payback(best.payback_years, best.investment_won)),
        ],
        left=right_left,
        top=top + 0.42,
        width=half,
    )
    chart_top = stats_bottom + gap
    _picture_block(
        slide,
        guide,
        figures.combination_png(comparison),
        "조합별 절감액과 투자비",
        left=right_left,
        top=chart_top,
        width=half,
        height=geometry.height_in - geometry.margin_in - chart_top,
    )


#: 부록 한 장에 담을 근거 줄 수. **넘치면 장을 나눈다** (39세션 5절).
#:
#: 36세션까지는 넘친 줄을 잘라 내고 「자리가 모자라 뺐다」 고 적었는데, 그러면
#: 선택요금 전환 하나만 실리고 **나머지 여섯이 통째로 빠졌다.** 분석한 자료를
#: 감출 이유가 없다 — 장을 늘린다.
#:
#: **열둘에서 열로 줄였다** (39세션 7절). 표 줄 높이는 PowerPoint 에서 최소값일
#: 뿐이라 실제로는 더 벌어진다 — 열둘을 넣으면 마지막 줄이 각주에 닿았다.
#:
#: **53세션에 열셋으로 늘렸다** (8-4). 해석 한 줄이 빠지면서(1-6) 0.5in 이 생겼다.
#: ESS 부록이 열셋이라 두 장으로 갈려 **뒷장에 세 줄만** 있었다.
APPENDIX_ROW_LIMIT = 13

#: 부록 한 장이 가져야 할 **최소 줄 수** (53세션 8-4).
#:
#: 이보다 적게 남으면 앞장에 붙여 한 장으로 만든다 — 상한을 조금 넘겨도 세 줄
#: 짜리 장을 만드는 것보다 낫다. **장이 넘칠 때 나누는 규칙의 하한**이다.
APPENDIX_MIN_ROWS = 4


@dataclass(frozen=True)
class AppendixPage:
    """부록 한 장 — 제목과 표 줄.

    수단 하나가 한 장에 안 들어가면 나눈다. 제목에 ``(1/2)`` 를 달아 이어지는
    장임을 밝힌다.
    """

    title: str
    rows: tuple[tuple[str, str, str], ...]


def appendix_pages(sections: DocumentSections) -> tuple[AppendixPage, ...]:
    """부록 장 목록 (39세션 5절).

    **절감액이 0이거나 미산출인 수단은 뺀다.** 산식과 대입값이 전부 0인 표는
    읽는 사람에게 알려 주는 것이 없고, 그 수단을 검토했다는 사실은 앞의
    「개선안별 요약」 표가 이미 말한다.

    수단이 하나도 남지 않으면 **빈 장 하나**를 돌려준다 — 부록 자리는 늘 있다.
    """
    priced = {entry.kind.key for entry in sections.measures if entry.has_saving}
    labels = {entry.kind.key: measure_slide_title(entry) for entry in sections.measures}
    pages: list[AppendixPage] = []
    for sheet in sections.worksheets:
        if sheet.key not in priced:
            continue
        records = [
            (str(value[0]), str(value[1]), str(value[2])) for value in sheet.frame().to_numpy()
        ]
        chunks = appendix_chunks(records)
        name = labels.get(sheet.key, sheet.title)
        for index, chunk in enumerate(chunks, start=1):
            suffix = f" ({index}/{len(chunks)})" if len(chunks) > 1 else ""
            pages.append(AppendixPage(f"{APPENDIX_SLIDE_TITLE} — {name}{suffix}", tuple(chunk)))
    return tuple(pages) or (AppendixPage(APPENDIX_SLIDE_TITLE, ()),)


def _is_blank(record: tuple[str, str, str]) -> bool:
    """빈 줄 — 근거 표에서 **묶음을 가르는 자리**다 (:func:`tariff_switch_worksheet`)."""
    return not any(value.strip() for value in record)


def appendix_chunks(
    records: Sequence[tuple[str, str, str]],
) -> list[list[tuple[str, str, str]]]:
    """부록 줄을 장으로 나눈다 (53세션 8-2·8-4).

    **먼저 묶음으로 가른다.** 선택요금 전환의 근거 표는 「현행 …」 과 「최적 …」
    사이에 빈 줄을 하나 두는데, 그것을 무시하고 열 줄씩 자르면 **한 장 맨 아래에
    다음 요금제가 걸쳐** 어디까지가 현행인지 알 수 없어진다.

    그 다음에야 :data:`APPENDIX_ROW_LIMIT` 로 자른다. **자투리는 앞장에 붙인다**
    — 세 줄짜리 장을 만드느니 상한을 조금 넘기는 편이 낫다
    (:data:`APPENDIX_MIN_ROWS`).
    """
    sections: list[list[tuple[str, str, str]]] = [[]]
    for record in records:
        if _is_blank(record):
            if sections[-1]:
                sections.append([])
            continue
        sections[-1].append(record)
    chunks: list[list[tuple[str, str, str]]] = []
    for section in sections:
        if not section:
            continue
        for start in range(0, len(section), APPENDIX_ROW_LIMIT):
            chunks.append(list(section[start : start + APPENDIX_ROW_LIMIT]))
        # **꼬리가 짧으면 앞에 붙인다.** 묶음 안에서만 붙인다 — 묶음을 가른
        # 까닭이 그 경계를 지키기 위해서다.
        if len(chunks) > 1 and 0 < len(chunks[-1]) < APPENDIX_MIN_ROWS:
            tail = chunks.pop()
            chunks[-1].extend(tail)
    return chunks or [[]]


def _build_appendix(
    slide: Slide, guide: DesignGuide, sections: DocumentSections, spec: SlideSpec
) -> None:
    """부록 산출근거 상세 — **표형. 수단마다 한 장 이상이다** (39세션 5절).

    **PPT 부록에는 산출 근거만 싣는다** (36세션 6절). ESS 조달 사례 표, 적용
    기준 데이터(Word 부록 B), 알려진 한계와 전제(Word 부록 C)는 여기 오지
    않는다 — 셋 다 **Word 에는 그대로 남아 있다.**

    **절감액이 없는 수단은 빼고, 뺐다는 사실을 각주가 적는다** — 조용히 빼면
    「검토하지 않았다」 로 읽힌다.

    **해석 한 줄을 두지 않는다** (53세션 1-6). 부록은 근거를 그대로 펼치는
    자리이지 읽는 법을 일러 주는 자리가 아니다 — 같은 문장이 부록 장마다
    되풀이돼 표가 밀려 내려가고 있었다. 무엇의 근거인지는 **제목이 적는다.**
    """
    geometry = guide.slide
    top = _title(slide, guide, spec.title)
    pages = appendix_pages(sections)
    page = pages[spec.page or 0]
    rows: list[list[str]] = [list(COLUMNS)]
    rows.extend([list(record) for record in page.rows])
    if len(rows) == 1:
        rows.append(["계산 근거", "절감액이 산출된 수단이 없습니다", "—"])
    appendix_note = _appendix_note(sections)
    _table(
        slide,
        guide,
        rows,
        left=geometry.margin_in,
        top=top,
        width=geometry.content_width_in,
        height=min(
            _table_room_above(guide, top=top, note_top=_note_top(guide, appendix_note)),
            0.4 * len(rows),
        ),
        widths=(0.24, 0.5, 0.26),
    )
    _caption(
        slide,
        guide,
        mark_note(appendix_note),
        left=geometry.margin_in,
        top=geometry.height_in - geometry.margin_in - 0.3,
        width=geometry.content_width_in,
    )


def _appendix_note(sections: DocumentSections) -> str:
    """부록 각주. **제작 사정을 적지 않는다** (39세션 2-4).

    「자리가 모자라 뺐습니다」 는 우리 쪽 사정이다. 뺀 것이 있다면 **무엇을 왜
    뺐는지**와 전문이 어디 있는지를 적는다.

    **규칙을 정확히 적는다** (59세션 11절). 뺀 갈래는 둘이다 — 「값이 0」 과
    「못 냈다」. 「산출되지 않은」 이라고만 적으면 계약전력 조정처럼 **산출은
    됐는데 0 인** 줄이 못 낸 것으로 읽힌다. 8절이 값 자리에서 둘을 가른 것과
    같은 자리다.
    """
    dropped = [entry for entry in sections.measures if not entry.has_saving]
    note = "전문은 Excel 부록 A 에 있습니다."
    if dropped:
        names = " · ".join(measure_slide_title(entry) for entry in dropped)
        note = f"절감액이 0 이거나 산출되지 않은 수단({names})은 근거를 싣지 않았습니다. {note}"
    return note


# ===================================================================== 37세션 · 마무리
#
# **샌드위치의 아랫빵이다** (가이드 3-2). 36세션은 표지만 다크로 두었는데,
# 지시서 목차에 마무리 장이 없었기 때문이다 — 구조가 반쪽이었다.
#
# 여기서 하는 말은 **다음 단계** 하나다. 이 검토는 초기 판단용이고, 채택한
# 개선안은 현장에서 확인해야 확정된다. 앞에서 같은 말을 하지 않는다 — 덱 어디에도
# 없던 문구라 옮겨 올 것이 없었고, 새로 적는 자리도 여기 하나뿐이다.
#
# **Word 에는 이 문구가 없다.** 그쪽에는 마무리 장이 없고, 있지도 않은 자리에
# 문구만 심으면 두 산출물이 같은 말을 다른 무게로 하게 된다. 필요해지면 그때
# 이 상수를 함께 읽으면 된다.

#: 마무리의 헤드라인. **결론이 아니라 성격을 밝힌다.**
NEXT_STEPS_HEADLINE = "이 결과는 초기 판단용입니다."

#: 다음 단계 셋. **무엇을 확인하는지까지 적는다** — 「현장 조사」 만으로는
#: 무엇을 보러 가는지 알 수 없다.
NEXT_STEPS: tuple[tuple[str, str], ...] = (
    ("현장 조사", "설치 공간과 수전 설비, 계통 연계 여건을 봅니다."),
    ("설비 사양 확인", "견적 단가와 실제 사양이 본 산출의 전제와 맞는지 봅니다."),
    ("담당자 인터뷰", "운영 계획과 제약을 듣습니다 — 부하는 자료가 고정입니다."),
)


def _build_closing(
    slide: Slide, guide: DesignGuide, _sections: DocumentSections, spec: SlideSpec
) -> None:
    """마무리 — **전체 배경 밴드가 다크다** (가이드 3-2 · 37세션).

    표지와 짝을 이룬다. 코랄은 라벨 한 줄에만 쓰고(3-2), 슬라이드 폭을 가로지르는
    색 바는 만들지 않는다(3-5) — 항목을 가르는 얇은 선뿐이다.

    **구분선 색이 표형 슬라이드와 다르다.** 밝은 회색을 검은 바탕에 그으면 선이
    튀어 내용보다 먼저 보인다 (:attr:`~kwise.report.design.Palette.on_dark_rule`).
    """
    geometry = guide.slide
    colors = guide.colors
    scale = guide.type_scale
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(colors.closing)

    left = geometry.margin_in + 0.5
    width = geometry.content_width_in - 1.0
    _text(
        slide,
        guide,
        [spec.title],
        left=left,
        top=1.95,
        width=width,
        height=0.3,
        size=scale.body,
        color=colors.coral,
        bold=True,
    )
    _text(
        slide,
        guide,
        [NEXT_STEPS_HEADLINE],
        left=left,
        top=2.42,
        width=width,
        height=0.8,
        size=scale.section,
        color=colors.on_dark,
        bold=True,
    )
    _text(
        slide,
        guide,
        ["채택한 개선안은 아래 셋을 거쳐 확정합니다."],
        left=left,
        top=3.5,
        width=width,
        height=0.32,
        size=scale.body,
        color=colors.on_dark_muted,
    )

    step = 0.72
    top = 4.06
    for index, (title, detail) in enumerate(NEXT_STEPS):
        y = top + step * index
        _text(
            slide,
            guide,
            [title],
            left=left,
            top=y,
            width=3.0,
            height=0.3,
            size=scale.body,
            color=colors.on_dark,
            bold=True,
        )
        _text(
            slide,
            guide,
            [detail],
            left=left + 3.2,
            top=y,
            width=width - 3.2,
            height=0.3,
            size=scale.body,
            color=colors.on_dark_muted,
        )
        # **선이 이 장의 시각 요소다.** 항목을 가르되 배경 위로 튀지 않는다.
        _rule(
            slide,
            guide,
            left=left,
            top=y + step - 0.18,
            length=width,
            color=colors.on_dark_rule,
        )


_BUILDERS: dict[str, Callable[[Slide, DesignGuide, DocumentSections, SlideSpec], None]] = {
    "cover": _build_cover,
    "agenda": _build_agenda,
    "building": _build_building,
    "usage_pattern": _build_usage_pattern,
    "peak_summary": _build_peak_summary,
    "peak_detail": _build_peak_detail,
    "structure": _build_structure,
    "measure_summary": _build_measure_summary,
    "surplus": _build_surplus,
    "combination": _build_combination,
    "closing": _build_closing,
    "appendix": _build_appendix,
}


# ===================================================================== 조립


def build_slides(
    sections: DocumentSections, *, guide: DesignGuide | None = None
) -> PresentationType:
    """덱을 만든다. **차례는 :func:`slide_specs` 하나가 쥔다.**"""
    design = guide if guide is not None else load_design_guide()
    presentation = Presentation()
    presentation.slide_width = Inches(design.slide.width_in)
    presentation.slide_height = Inches(design.slide.height_in)
    blank = presentation.slide_layouts[6]
    colors = design.colors

    for spec in slide_specs(sections):
        slide = presentation.slides.add_slide(blank)
        if spec.layout != "cover":
            # **본문 콘텐츠는 라이트다** (36세션 3-2). 샌드위치의 속이다.
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = _rgb(colors.white)
        builder = _BUILDERS.get(spec.key, _build_measure)
        builder(slide, design, sections, spec)
    return presentation


def slides_path(
    output_dir: Path = DEFAULT_OUTPUT_DIR,
    *,
    prefix: str = "kwise_report",
    now: dt.datetime | None = None,
) -> Path:
    """**날짜·시각 접미사를 붙인다.** PowerPoint 가 열고 있으면 덮어쓰기가 실패한다."""
    stamp = (now if now is not None else dt.datetime.now()).strftime("%Y%m%d_%H%M")
    return output_dir / f"{prefix}_{stamp}.pptx"


def export_slides(
    sections: DocumentSections,
    *,
    output_dir: Path = DEFAULT_OUTPUT_DIR,
    prefix: str = "kwise_report",
    now: dt.datetime | None = None,
) -> Path:
    """파일로 쓴다."""
    path = slides_path(output_dir, prefix=prefix, now=now)
    path.parent.mkdir(parents=True, exist_ok=True)
    build_slides(sections).save(str(path))
    return path


def slides_bytes(
    sections: DocumentSections, *, now: dt.datetime | None = None, prefix: str = "kwise_report"
) -> tuple[bytes, str]:
    """내려받기용 바이트와 파일명. **디스크에 남기지 않는다** (10.2)."""
    buffer = io.BytesIO()
    build_slides(sections).save(buffer)
    return buffer.getvalue(), slides_path(Path(), prefix=prefix, now=now).name
