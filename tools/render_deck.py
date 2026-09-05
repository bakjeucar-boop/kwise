r"""PPT 실물 렌더 — **화면을 그대로 태워 덱을 뽑는다** (59세션).

    .venv\Scripts\python.exe tools\render_deck.py --case large-b --png

53세션 9절부터 「PowerPoint COM 으로 png 를 뽑아 한 장씩 본다」 가 확인 방법인데,
정작 **덱을 만드는 절차가 세션마다 손으로 다시 짜여 있었다.** 같은 자료·같은
계약으로 다시 뽑지 못하면 다음 세션이 「전에 본 그 장」 을 못 찾는다.

**계산을 여기서 다시 하지 않는다.** ``streamlit.testing`` 으로 실제 화면을 띄우고
「PPT 보고서 만들기」 를 눌러 세션에 담긴 바이트를 그대로 꺼낸다 — 사용자가 받는
파일과 한 바이트도 다르지 않다.

png 는 PowerPoint COM 이 있어야 나온다 (Windows + PowerPoint). 없으면 pptx 만
남기고 그 사실을 적는다.
"""

from __future__ import annotations

import argparse
import subprocess
import sys
from collections.abc import Sequence
from dataclasses import dataclass
from pathlib import Path
from typing import TYPE_CHECKING

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT / "src"))

# **`sys.path` 를 세운 뒤에 들여온다.** 위에 두면 저장소를 설치하지 않은
# 환경에서 이 도구가 통째로 못 뜬다.
from kwise.report.figures import FigureFailureCollector  # noqa: E402

if TYPE_CHECKING:  # streamlit 을 도구 시작에 들이지 않는다 — 형에만 쓴다
    from kwise.ui.pipeline import SolarInputs

APP = PROJECT_ROOT / "src" / "kwise" / "ui" / "app.py"
LARGE_CSV = PROJECT_ROOT / "input" / "사용량조회_20240429.csv"
SMALL_CSV = PROJECT_ROOT / "input" / "사용량조회_소형사무빌딩.csv"
#: 용인 소규모 건물 실측 (61세션). **회귀 케이스가 아니라 디버깅용**이다 —
#: 케이스 스터디 목록에는 넣지 않았다 (61세션 8절).
YONGIN_XLSX = PROJECT_ROOT / "input" / "전기사용량_소형건물.xlsx"

#: 통합 시험과 같은 지역을 쓴다 (48세션). 지역을 바꾸면 발전량이 통째로 흔들린다.
#:
#: **벌의 기본값으로 쓰지 않는다** (97세션). 96세션이 낸 결함이 이 상수를
#: `SolarInputs` 에 바로 먹인 것이었다 — 용인 두 벌이 화면에는 「경기도 용인시」
#: 라고 적고 발전량은 강릉 격자로 냈다. 이제 벌마다 `sigungu` 를 반드시 적고,
#: 강릉을 쓰는 벌은 이 이름을 그 자리에 적는다.
REGION = "강원도/강릉시"

ALL_MEASURES = (
    "tariff_switch",
    "contract",
    "demand_response",
    "power_factor",
    "solar",
    "ess",
)


@dataclass(frozen=True)
class Case:
    """뽑을 한 벌. **이름이 곧 재현 조건이다.**"""

    key: str
    title: str
    csv: Path
    contract_type: str
    voltage: str
    option: str
    contract_kw: float
    area_m2: float
    building_name: str
    sigungu: str
    """지역. **기본값을 두지 않는다** — 적지 않으면 벌이 아예 만들어지지 않는다.

    화면과 태양광이 **같은 이 값**을 쓴다 (97세션 1절). 발전량이 통째로 흔들리는
    값이라 「적었겠지」 로 넘길 자리가 아니다.
    """
    surplus_use: str = ""
    """고른 잉여 처리. 비우면 화면 기본값(출력제어)이다 (57세션)."""

    @property
    def province(self) -> str:
        """옆단 1단(시도). **따로 적지 않는다** — 적으면 시군구와 어긋날 수 있다."""
        return self.sigungu.partition("/")[0]


CASES: tuple[Case, ...] = (
    Case(
        key="large-b",
        title="대형 · 일반용(을) 고압A 선택Ⅰ",
        csv=LARGE_CSV,
        contract_type="general_b",
        voltage="high_a",
        option="I",
        contract_kw=6_000.0,
        area_m2=20_000.0,
        building_name="대형 사업장",
        sigungu=REGION,
    ),
    Case(
        key="large-a",
        title="대형 · 일반용(갑)Ⅰ 고압A",
        csv=LARGE_CSV,
        contract_type="general_a_1",
        voltage="high_a",
        option="",
        contract_kw=6_000.0,
        area_m2=20_000.0,
        building_name="대형 사업장(갑)",
        sigungu=REGION,
    ),
    Case(
        key="small-b",
        title="소형 · 일반용(을) 고압A 선택Ⅰ",
        csv=SMALL_CSV,
        contract_type="general_b",
        voltage="high_a",
        option="I",
        contract_kw=300.0,
        area_m2=1_000.0,
        building_name="소형 사무빌딩",
        sigungu=REGION,
    ),
    Case(
        key="small-a",
        title="소형 · 일반용(갑)Ⅰ 고압A",
        csv=SMALL_CSV,
        contract_type="general_a_1",
        voltage="high_a",
        option="",
        contract_kw=300.0,
        area_m2=1_000.0,
        building_name="소형 사무빌딩(갑)",
        sigungu=REGION,
    ),
    Case(
        key="large-b-over",
        title="대형 · 일반용(을) 고압A 선택Ⅰ · 계약전력 과다 20,000 kW",
        csv=LARGE_CSV,
        contract_type="general_b",
        voltage="high_a",
        option="I",
        # **요금적용전력 하한(계약전력의 30%)이 최대수요를 넘는 자리다** —
        # 6,000 kW 하한이 관측 최대 5,293 kW 보다 커서 피크를 낮춰도 기준 전력이
        # 그대로다. ESS 사양 표의 저감량이 전 줄 0 kW 로 선다 (59세션 3절).
        contract_kw=20_000.0,
        area_m2=20_000.0,
        building_name="대형 사업장(계약 과다)",
        sigungu=REGION,
    ),
    Case(
        key="large-b-short",
        title="대형 · 일반용(을) 고압A 선택Ⅰ · 계약전력 부족 4,000 kW (합성 조건)",
        csv=LARGE_CSV,
        contract_type="general_b",
        voltage="high_a",
        option="I",
        # **초과사용부가금이 실제로 서는 유일한 벌이다** (S127 2절 · ②-32).
        # 앞서는 덱 열 + 케이스 일곱이 **전부 초과 구간 0** 이었다 —
        # 계약전력이 언제나 관측 최대 위에 있어 제67조의3 ③ 갈래의 글과 숫자가
        # **한 번도 그려진 적이 없다.** 83세션의 `small-a2-was`, 91세션의
        # `small-edu-a-over` 와 같은 자리다 — **뜨지 않는 갈래는 없는 갈래와 같다.**
        #
        # **`over` 와 반대 방향이다.** `large-b-over` 의 `over` 는 「계약전력이
        # 과다하다」 는 뜻이고(20,000 kW), 여기 `short` 는 「계약전력이 모자라
        # 초과사용부가금이 붙는다」 는 뜻이다.
        #
        # **자료는 실측이고 계약전력 4,000 kW 만 합성이다.** 같은 자료를 쓰는
        # `large-b`(6,000) 에서 계약전력 하나만 갈았다. 일반용(을)은 계약전력
        # 300 kW 이상이라(기본공급약관 제57조 ②) 4,000 kW 는 성립하는 수다.
        #
        #     관측 최대 5,293.44 kW  >  계약 4,000 kW  → 초과 구간 5,050건
        #     초과한 달 13 · 청구되는 달 12 (첫 달은 예고 · 제67조의3 ④)
        #     배수 셋이 다 선다 — 1.5(9달) · 2.0(1달) · 2.5(3달)
        contract_kw=4_000.0,
        area_m2=20_000.0,
        building_name="대형 사업장(계약 부족)",
        sigungu=REGION,
    ),
    Case(
        key="small-b-sell",
        title="소형 · 일반용(을) 고압A 선택Ⅰ · 잉여 외부 판매",
        csv=SMALL_CSV,
        contract_type="general_b",
        voltage="high_a",
        option="I",
        contract_kw=300.0,
        area_m2=1_000.0,
        building_name="소형 사무빌딩(외부 판매)",
        sigungu=REGION,
        # **기본이 아닌 잉여 처리를 골라 본다** (59세션 5절). 기본값(출력제어)은
        # 0원이라 「절감액에 잉여가 얹힌다」 는 사실이 값으로 안 보인다.
        surplus_use="외부 판매",
    ),
    Case(
        key="small-a2",
        title="용인 소규모 · 일반용(갑)Ⅱ 고압A 선택Ⅱ · 계약전력 290 kW",
        csv=YONGIN_XLSX,
        contract_type="general_a_2",
        voltage="high_a",
        option="II",
        # **갑Ⅱ 가 요금적용전력을 쓰는지 실물로 보는 자리다** (61세션).
        # 청구서(8,230원/kW × 118 kW)와 대조한 벌이고, 고치기 전에는 기본요금이
        # 계약전력 290 kW 로 매겨져 12개월 28,640,400원이 나오고 있었다.
        contract_kw=290.0,
        area_m2=1_000.0,
        building_name="용인 소규모 건물(갑Ⅱ)",
        sigungu="경기도/용인시",
    ),
    Case(
        key="small-a2-was",
        title="용인 소규모 · 일반용(을) 고압A 선택Ⅰ · 옛 계약전력 700 kW",
        csv=YONGIN_XLSX,
        contract_type="general_b",
        voltage="high_a",
        option="I",
        # **하한이 이기는 갈래가 뜨는 유일한 실측 벌이다** (83세션 14).
        # 시험 자료 셋(용인 290 · 대형 6,000 · 소형 300)이 전부 하한이 지는
        # 쪽이라 그 갈래가 실물에 한 번도 서지 않았다 — **뜨지 않는 갈래는
        # 없는 갈래와 같다.** 700 은 지어낸 값이 아니라 **이 건물이 실제로
        # 쓰던 계약전력**이고, 290 으로 내리면서 77.7 kW 어치가 사라졌다.
        #
        #     하한 700 × 30% = 210 kW  >  최대수요 132.3 kW  → 하한이 기준
        #     목표 132.3 ÷ 0.3 = 441 kW
        #
        # **종별이 갑Ⅱ 에서 을로 바뀌었다** (96세션). 96세션이 종별 경계를
        # 읽게 하자 **갑Ⅱ 700 kW 는 기준 데이터가 스스로 금지하는 조합**이
        # 됐다 — 갑Ⅱ 는 300 kW 미만이라 계약전력도 목표 441 kW 도 경계 밖이다.
        # 을은 300 kW **이상**이라 700 도 441 도 안쪽이고, 하한이 이기는 갈래는
        # 그대로 선다. 옛 계약 700 kW 는 실측이지만 **그때의 종별이 을이었다는
        # 뜻은 아니다** — 이 벌은 하한 갈래를 세우는 자리이지 청구서 재현이
        # 아니다. 청구서 재현은 `small-a2`(갑Ⅱ 290 kW)가 한다.
        contract_kw=700.0,
        area_m2=1_000.0,
        building_name="용인 소규모 건물(옛 계약 700 · 을)",
        sigungu="경기도/용인시",
    ),
    Case(
        key="small-edu-a",
        title="소형 · 교육용(갑) 고압A 선택Ⅰ · 계약전력 300 kW",
        csv=SMALL_CSV,
        contract_type="education_a",
        voltage="high_a",
        option="I",
        # **89세션이 고친 자리를 실물로 보는 벌이다.** 교육용(갑) 고압은
        # 제38조 ②로 최대수요전력계가 서므로 기본요금이 요금적용전력에 붙는다
        # (제68조 ①) — 고치기 전에는 계약전력 300 kW 로 매겨졌다.
        # **ESS 장이 서는지도 여기서 본다** — 88세션까지는 종별만 보고 막혀
        # 「성립하지 않는다」 장이 나왔다.
        contract_kw=300.0,
        area_m2=1_000.0,
        building_name="소형 학교(교육용 갑)",
        sigungu=REGION,
    ),
    Case(
        key="small-edu-a-over",
        title="소형 · 교육용(갑) 고압A 선택Ⅰ · 계약전력 과다 950 kW",
        csv=SMALL_CSV,
        contract_type="education_a",
        voltage="high_a",
        option="I",
        # **교육용에서 하한이 이기는 갈래를 세우는 벌이다** (91세션 1절).
        # 90세션이 교육용 둘의 하한을 0.3 으로 세웠는데 **이기는 벌이 하나도
        # 없었다** — `small-edu-a` 는 하한 90 kW 가 최대수요 264.7 kW 에 져서
        # 「이미 적정」 으로만 서고, 하한이 이길 때의 글과 그림은 교육용에서
        # 한 번도 그려진 적이 없다. 83세션이 갑Ⅱ 에 `small-a2-was` 를 지은
        # 자리와 같다 — **뜨지 않는 갈래는 없는 갈래와 같다.**
        #
        #     하한 950 × 30% = 285 kW  >  최대수요 264.7 kW  → 하한이 기준
        #     목표 264.7 ÷ 0.3 = 883 kW
        #
        # **950 은 성립하는 수다** — 교육용전력(갑)은 계약전력 4 kW 이상
        # 1,000 kW 미만이고(약관 제58조 ② 1.) 고압A 는 표준전압 3,300~66,000V
        # 고객이라 계약전력에 따로 걸리는 하한이 없다.
        contract_kw=950.0,
        area_m2=1_000.0,
        building_name="소형 학교(교육용 갑 · 계약 과다)",
        sigungu=REGION,
    ),
)

BY_KEY = {case.key: case for case in CASES}


def solar_inputs_for(case: Case) -> SolarInputs:
    """이 벌의 태양광 입력. **좌표는 벌 정의에서 온다** (97세션 1절).

    `build_deck` 안에 있던 한 줄을 여기로 꺼냈다. 안에 있으면 좌표를 확인하려고
    화면을 통째로 띄워야 하고, 그래서 **96세션까지 아무 시험도 이 줄을 안 봤다.**
    """
    from kwise.ui.pipeline import SolarInputs

    return SolarInputs(region_key=case.sigungu, area_m2=case.area_m2)


def _first_option(contract_type: str, voltage: str) -> str:
    from kwise.tariff import load_tariff
    from kwise.ui.pipeline import option_choices

    options = option_choices(load_tariff(), contract_type, voltage)
    return options[0] if options else ""


def build_deck(case: Case, *, timeout: int = 1800) -> bytes:
    """화면을 띄워 덱 바이트를 받는다."""
    from streamlit.testing.v1 import AppTest

    from kwise.ui.artifacts import ARTIFACT_KEY
    from kwise.ui.pipeline import ContractForm

    app = AppTest.from_file(str(APP), default_timeout=timeout)
    state = app.session_state
    state["upload_bytes"] = case.csv.read_bytes()
    state["upload_name"] = case.csv.name
    state["contract_form"] = ContractForm(
        contract_type=case.contract_type,
        voltage=case.voltage,
        option=case.option or _first_option(case.contract_type, case.voltage),
        contract_kw=case.contract_kw,
    )
    state["building_province"] = case.province
    state["building_sigungu"] = case.sigungu
    state["building_name"] = case.building_name
    # **태양광은 「계산」 을 누른 상태로 시작한다.** 위젯에 키가 없어 면적을
    # 세션으로 심을 수 없다 — 눌린 결과(``solar_inputs``)를 바로 넣는다.
    state["solar_inputs"] = solar_inputs_for(case)
    if case.surplus_use:
        state["measure_solar_surplus_use"] = case.surplus_use
    for key in ALL_MEASURES:
        state[f"measure_on_{key}"] = True
    state["combination_pick"] = ALL_MEASURES

    app.run()
    if app.exception:
        raise RuntimeError(f"{case.key} — 화면이 예외로 멈췄다: {app.exception}")
    app.button(key="build_ppt").click().run(timeout=timeout)
    if app.exception:
        raise RuntimeError(f"{case.key} — PPT 만들기가 예외로 멈췄다: {app.exception}")
    try:
        store = dict(app.session_state[ARTIFACT_KEY])
    except (KeyError, AttributeError, TypeError):
        store = {}
    artifact = store.get("ppt")
    if artifact is None:
        raise RuntimeError(f"{case.key} — 덱이 만들어지지 않았다")
    payload: bytes = artifact.payload
    return payload


_EXPORT_PS1 = r"""
$ErrorActionPreference = 'Stop'
$app = New-Object -ComObject PowerPoint.Application
$deck = $app.Presentations.Open('{pptx}', $true, $false, $false)
$deck.Export('{outdir}', 'png', 1600, 900)
$deck.Close()
$app.Quit()
"""


def export_png(pptx: Path, outdir: Path) -> bool:
    """PowerPoint COM 으로 1600×900 png 를 뽑는다. 없으면 ``False``."""
    outdir.mkdir(parents=True, exist_ok=True)
    script = _EXPORT_PS1.format(pptx=str(pptx.resolve()), outdir=str(outdir.resolve()))
    try:
        done = subprocess.run(
            ["powershell", "-NoProfile", "-NonInteractive", "-Command", script],
            capture_output=True,
            text=True,
            timeout=600,
            check=False,
        )
    except (OSError, subprocess.TimeoutExpired) as exc:
        print(f"  png 실패 — {exc}")
        return False
    if done.returncode != 0:
        print(f"  png 실패 — {done.stderr.strip()[:400]}")
        return False
    return True


#: png 폴더에 함께 두는 라벨 파일. **맨 앞에 오게 밑줄로 시작한다.**
LABEL_NAME = "_이_덱은.txt"


def slide_titles(pptx: Path) -> list[str]:
    """덱에서 장 제목을 뽑는다 — **장 수가 다른 까닭을 라벨이 스스로 말하게 한다.**"""
    from pptx import Presentation

    titles: list[str] = []
    for slide in Presentation(str(pptx)).slides:
        heads = sorted(
            (
                shape
                for shape in slide.shapes
                if shape.has_text_frame and shape.top is not None and shape.text_frame.text.strip()
            ),
            key=lambda shape: shape.top,
        )
        titles.append(heads[0].text_frame.text.strip().splitlines()[0] if heads else "?")
    return titles


def label_text(case: Case, titles: Sequence[str] = (), failures: Sequence[str] = ()) -> str:
    """이 덱이 무엇인지 사람이 읽을 한 벌 (60세션 곁가지).

    **이름만 봐서 모르는 산출물은 실물 확인이 안 된다.** ``large-b-over`` 가
    「대형 갑」 으로 읽혀 결론 두 줄 겹침(T1)을 엉뚱한 장에서 찾았다 — 그것은
    **대형 을**이고 갑은 ``large-a`` 다.
    """
    lines = [
        f"{case.key} — {case.title}",
        "",
        f"자료        {case.csv.name}",
        f"계약종별    {case.contract_type}",
        f"전압구분    {case.voltage}",
        f"선택요금    {case.option or '(화면 기본값)'}",
        f"계약전력    {case.contract_kw:,.0f} kW",
        f"면적        {case.area_m2:,.0f} m²",
        f"건물명      {case.building_name}",
        f"지역        {case.sigungu}",
        # **0 이면 0 이라 적는다** (60세션 11절). 줄이 없으면 「기록을 안 남긴
        # 것」 과 「실패가 없던 것」 이 갈리지 않는다.
        f"그림 실패    {len(failures)}개",
    ]
    if case.surplus_use:
        lines.append(f"잉여 처리   {case.surplus_use}")
    lines += [f"    {line}" for line in failures]
    if titles:
        # **장 수는 벌마다 다르다.** 갈리는 것은 수단 장이 아니라 **부록 장**이다 —
        # 「값이 0 이거나 미산출인 수단은 부록에서 뺀다」 (39세션 · ``has_saving``).
        # 그래서 대형 갑은 부록 ESS 가 빠져 20장, 계약 과다는 부록 계약전력
        # 조정이 붙어 22장이다. 소형은 「잉여 활용」 이 하나 붙고 부록 ESS 가
        # 빠져 21장이 된다 — **더한 것과 뺀 것이 상쇄돼 대형 을과 같아 보인다.**
        lines += [
            "",
            f"장 수        {len(titles)}장",
            "",
            "  장 수가 벌마다 다른 까닭 —",
            "    · 부록 산출근거는 절감액이 0 이거나 미산출인 수단을 뺀다",
            "    · 「잉여 활용」 은 잉여가 실제로 날 때만 선다",
            "",
            "  ** 장 수가 같아도 같은 덱이 아니다 **",
            "    소형 21장과 대형 을 21장은 속이 다르다 — 소형은 「잉여 활용」 이",
            "    하나 붙고 부록 ESS 가 하나 빠져 **상쇄된 것**이다. 우연히 같다.",
            "    아래 장 차례를 대조하십시오.",
            "",
            "  장 차례 —",
        ]
        lines += [f"    {index:>2}. {title}" for index, title in enumerate(titles, 1)]
    lines += [
        "",
        "다시 뽑으려면 —",
        f"    .venv\\Scripts\\python.exe tools\\render_deck.py --case {case.key} --png",
    ]
    return "\n".join(lines) + "\n"


def write_label(
    outdir: Path,
    case: Case,
    pptx: Path | None = None,
    failures: Sequence[str] = (),
) -> Path:
    """png 폴더에 라벨을 남긴다. 덱을 주면 **장 차례**도 싣는다."""
    titles = slide_titles(pptx) if pptx is not None else []
    path = outdir / LABEL_NAME
    path.write_text(label_text(case, titles, failures), encoding="utf-8")
    return path


def write_index(outdir: Path, cases: Sequence[Case]) -> Path:
    """산출 폴더에 **이번에 뽑은 벌 목록**을 남긴다."""
    rows = [f"{'벌':<14}{'자료':<26}{'계약종별':<14}{'계약전력':>12}  제목", ""]
    rows += [
        f"{case.key:<14}{case.csv.name:<26}{case.contract_type:<14}"
        f"{case.contract_kw:>10,.0f} kW  {case.title}"
        for case in cases
    ]
    path = outdir / "_뽑은_벌.txt"
    path.write_text("\n".join(rows) + "\n", encoding="utf-8")
    return path


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="PPT 실물 렌더")
    parser.add_argument(
        "--case",
        action="append",
        choices=[case.key for case in CASES],
        help="뽑을 벌. 여러 번 줄 수 있다. 없으면 전부.",
    )
    parser.add_argument("--png", action="store_true", help="PowerPoint COM 으로 png 도 뽑는다")
    parser.add_argument(
        "--out",
        type=Path,
        default=PROJECT_ROOT / "output" / "decks",
        help="산출 폴더",
    )
    args = parser.parse_args(argv)

    picked = [BY_KEY[key] for key in (args.case or [case.key for case in CASES])]
    args.out.mkdir(parents=True, exist_ok=True)
    # **그림 굽기 실패를 벌마다 따로 센다** (60세션 11절). 덱을 다 뽑고 나서
    # 한 번에 세면 어느 벌의 실패인지 갈리지 않는다.
    failed: dict[str, list[str]] = {}
    for case in picked:
        print(f"[{case.key}] {case.title}")
        with FigureFailureCollector() as collector:
            payload = build_deck(case)
        if collector.messages:
            failed[case.key] = list(collector.messages)
            print(f"  ** 그림 {len(collector.messages)}개가 안 구워졌다 **")
            for line in collector.messages:
                print(f"     {line}")
        pptx = args.out / f"{case.key}.pptx"
        pptx.write_bytes(payload)
        print(f"  {pptx} ({len(payload) / 1024:.0f} KB)")
        if args.png and export_png(pptx, args.out / case.key):
            write_label(args.out / case.key, case, pptx, failed.get(case.key, []))
            # Windows 는 glob 이 대소문자를 가리지 않아 `*.PNG` 와 `*.png` 가
            # 같은 파일을 둘 다 문다 — 세는 자리에서 겹치지 않게 한 번만 훑는다.
            count = len(
                [item for item in (args.out / case.key).iterdir() if item.suffix.lower() == ".png"]
            )
            print(f"  png {count}장 → {args.out / case.key}")
    index = write_index(args.out, picked)
    print(f"[목록] {index}")
    # **끝에 다시 낸다.** 여러 벌을 잇달아 뽑으면 앞의 실패는 화면 밖으로 밀려 있다.
    # **수를 적지 않는다** — 「여섯」 이라 적어 두었더니 `small-a2` 가 붙은 뒤
    # 열한 세션 동안 틀린 채 남아 있었다 (72세션).
    if failed:
        total = sum(len(items) for items in failed.values())
        print("")
        print(f"[그림 실패] 벌 {len(failed)}개에서 {total}개 — 그 장은 그림 없이 나갔다")
        for key, items in failed.items():
            print(f"  {key}: {len(items)}개")
    else:
        print("")
        print("[그림 실패] 0개")
    return 0


if __name__ == "__main__":  # pragma: no cover
    raise SystemExit(main())
