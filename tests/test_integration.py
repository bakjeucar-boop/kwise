"""전체 기능 통합 시험 (15세션 5절).

**화면을 실제로 띄워 전 경로를 훑는다.** 순수 함수 시험이 다 통과해도 배선이
틀리면 화면에서만 드러난다 — 8·12·13세션에서 잡은 결함이 모두 그런 종류였다.

여기서 지키는 것 넷.

    5-1 화면 흐름   업로드 → 계약 → 1·2·3단계. 수단 0/1/전부. 켜고 끄기 반복
    5-2 산출물     Excel 시트·Word 장·내려받기 뒤 결과 유지
    5-3 경계       미업로드·미입력·잉여 0·기상 실패
    5-4 회귀       케이스 스터디 값·독립 평가·2↔3단계 일치·표기 규약

**기상은 격리되어 있다** (``conftest.isolated_weather_archive``). 태양광이 실제로
도는 경로는 ``weather`` 표식을 붙여 사전 취득분을 쓰게 한다 — 격리를 풀지 않으면
"API 실패 시 멈추는가" 같은 시험이 조용히 성공해 버린다.
"""

from __future__ import annotations

import functools
import re
from collections.abc import Iterator
from pathlib import Path
from typing import TYPE_CHECKING, Any

import pytest
from streamlit.testing.v1 import AppTest

if TYPE_CHECKING:  # 시험이 무거워지지 않게 형만 빌린다 (70세션 2절)
    from kwise.io import UsageData
    from kwise.measures import DispatchResult, PowerFactorResult, TariffSwitchResult
    from kwise.report.days import RepresentativeDay
    from kwise.tariff import BillingResult, TariffTable

from kwise.notices import texts
from kwise.report import SHEET_ORDER
from kwise.ui.nav import RULES_PAGE, TABS
from kwise.ui.pipeline import ContractForm
from kwise.ui.spec import MEASURES
from tests._appmemo import harvest_ess_memo, seed_ess_memo

APP = (Path("src") / "kwise" / "ui" / "app.py").resolve()
SAMPLE = Path("input") / "사용량조회_20240429.csv"
ALL_MEASURES = (
    "tariff_switch",
    "contract",
    "demand_response",
    "power_factor",
    "solar",
    "ess",
)
# 요금에 영향을 주는 수단만. 태양광은 「계산」 단추를 눌러야 도는 별도 경로다.
BILLED_MEASURES = ("tariff_switch", "contract", "power_factor", "ess")


@pytest.fixture
def real_weather(monkeypatch: pytest.MonkeyPatch) -> Iterator[None]:
    """사전 취득분을 쓴다. **격리를 푸는 유일한 자리다.**"""
    monkeypatch.delenv("KWISE_WEATHER_DIR", raising=False)
    yield


#: 통합 시험이 쓰는 지역. **48세션 전의 기본값 그대로다** — 지역을 바꾸면
#: 발전량이 달라져 태양광 값이 통째로 흔들린다.
TEST_PROVINCE = "강원도"
TEST_REGION = "강원도/강릉시"


def _app(**state: object) -> AppTest:
    running = AppTest.from_file(str(APP), default_timeout=900)
    running.session_state["upload_bytes"] = SAMPLE.read_bytes()
    running.session_state["upload_name"] = SAMPLE.name
    running.session_state["contract_form"] = ContractForm(
        contract_type="general_b", voltage="high_a", option="I", contract_kw=6_000.0
    )
    # **지역을 명시한다** (48세션). 48세션 전에는 옆단 드롭다운이 가나다순 첫
    # 항목을 기본값으로 내고 있어 아무것도 심지 않아도 지역이 잡혔다 — 그 기본값을
    # 없앴으므로 시험이 무엇을 쓰는지 스스로 밝힌다. 옛 기본값 그대로다.
    running.session_state["building_province"] = TEST_PROVINCE
    running.session_state["building_sigungu"] = TEST_REGION
    for key, value in state.items():
        running.session_state[key] = value
    # 3단계는 「합산효과 계산」 을 누른 뒤에 그린다 (33세션 5절). 단추가 하는 일이
    # 이 상태를 심는 것뿐이라, 시험은 눌린 상태로 시작한다.
    if "combination_pick" not in state:
        running.session_state["combination_pick"] = tuple(
            item.key for item in MEASURES if state.get(f"measure_on_{item.key}")
        )
    return harvest_ess_memo(seed_ess_memo(running).run())


def _blank(**state: object) -> AppTest:
    """업로드도 계약 정보도 없는 상태."""
    running = AppTest.from_file(str(APP), default_timeout=900)
    for key, value in state.items():
        running.session_state[key] = value
    return running.run()


def _on(*keys: str) -> dict[str, object]:
    return {f"measure_on_{key}": True for key in keys}


# **탭 구조라 세 화면이 한 번에 그려진다** (16세션 1절). 지표 목록이 셋 이어
# 붙으므로 경계로 갈라 본다 — 1단계의 마지막은 「기본요금 비중」, 3단계의 처음은
# 「단순 합」이다.
_DIAGNOSE_LAST = "기본요금 비중"
_STAGE3_FIRST = "단순 합"


def _stage2_metrics(app: AppTest) -> list[tuple[str, str]]:
    """2단계 카드가 낸 지표만."""
    pairs = [(str(item.label), str(item.value)) for item in app.metric]
    labels = [label for label, _ in pairs]
    start = labels.index(_DIAGNOSE_LAST) + 1 if _DIAGNOSE_LAST in labels else 0
    end = labels.index(_STAGE3_FIRST) if _STAGE3_FIRST in labels else len(pairs)
    return pairs[start:end]


def _stage3_labels(app: AppTest) -> list[str]:
    labels = [str(item.label) for item in app.metric]
    return labels[labels.index(_STAGE3_FIRST) :] if _STAGE3_FIRST in labels else []


def _text(app: AppTest) -> str:
    parts = [
        str(item.value)
        for group in (app.markdown, app.caption, app.warning, app.error, app.info)
        for item in group
    ]
    parts += [f"{item.label} {item.value} {item.delta or ''}" for item in app.metric]
    parts += [str(item.value) for item in app.subheader]
    return "\n".join(parts)


# ===================================================================== 5-1 화면 흐름


def test_세_탭이_한_번에_뜬다() -> None:
    """**탭 구조다** (16세션 1절). 한 실행에서 셋이 모두 그려진다."""
    app = _app()
    assert not app.exception, app.exception
    headers = [str(item.value) for item in app.header]
    for tab in TABS:
        number = tab.split(" ", 1)[0]
        assert any(number in item for item in headers), (number, headers)


def test_기준_데이터_화면도_예외_없이_뜬다() -> None:
    app = _app(nav_page=RULES_PAGE)
    assert not app.exception, app.exception


@pytest.mark.parametrize("count", [0, 1, len(BILLED_MEASURES)])
def test_수단을_0개_1개_전부_켜도_3단계가_돈다(count: int) -> None:
    """**수단 0개가 정상 경로다** — 진단만 보고 받아 가는 경우다."""
    app = _app(**_on(*BILLED_MEASURES[:count]))
    assert not app.exception, app.exception
    body = _text(app)
    assert "개선안별 요약" in body
    if count == 0:
        assert "수단을 하나도 켜지 않았습니다" in body
    else:
        assert "단순 합" in body


def test_모든_수단을_켜도_2단계가_돈다() -> None:
    app = _app(**_on(*ALL_MEASURES))
    assert not app.exception, app.exception
    body = _text(app)
    # 태양광은 계산 전이라 안내만. **잉여 처리는 태양광 안에 있으므로** 계산
    # 전에는 나올 자리가 없다 (41세션 2-2).
    assert "「태양광 계산」 을 누르십시오" in body
    assert "잉여 처리" not in body, body


def test_수단을_켜고_끄기를_반복해도_상태가_꼬이지_않는다() -> None:
    """켜고 끄기를 되풀이해도 카드가 살아 있고 예외가 없어야 한다.

    **선택 오브젝트가 체크박스로 바뀌었다** (27세션 1절). 세션 키는 그대로다 —
    16세션에 잡은 키 충돌·값 유실이 재발하지 않는지 여기서 본다.
    """
    app = _app(**_on("tariff_switch"))
    for _ in range(2):
        app.checkbox(key="measure_on_tariff_switch").set_value(False).run()
        assert not app.exception, app.exception
        app.checkbox(key="measure_on_tariff_switch").set_value(True).run()
        assert not app.exception, app.exception
    assert any(item.label == "가장 유리한 요금제" for item in app.metric)


def test_옆단은_건물_정보다() -> None:
    """**계약 정보는 옆단에 없다** (16세션 2절). 건물이 아니라 계약이다."""
    app = _app()
    labels = [str(item.label) for item in app.sidebar.text_input]
    labels += [str(item.label) for item in app.sidebar.selectbox]
    labels += [str(item.label) for item in app.sidebar.number_input]
    assert set(labels) == {
        "건물명 (선택)",
        "용도 (선택)",
        # 48세션에 「선택 안 함」 이 붙어 선택 항목이 됐다 — 가나다순 첫 항목이
        # 고른 것처럼 보이던 자리다.
        "시도 (선택)",
        "시군구",
        "연면적 (m², 선택)",
        "준공연도 (선택)",
    }
    banned = ("계약종별", "전압구분", "계약전력", "선택요금")
    assert not [item for item in labels if any(word in item for word in banned)], labels


def test_기준_데이터로_갔다_돌아온다() -> None:
    """**단계가 아니라 설정이다** — 옆단 하단의 별도 진입점 (16세션 1절)."""
    app = _app()
    app.sidebar.button(key="nav_rules").click().run()
    assert app.session_state["nav_page"] == RULES_PAGE
    assert not app.exception, app.exception
    app.sidebar.button(key="nav_analysis").click().run()
    assert not app.exception, app.exception
    assert any("1단계" in str(item.value) for item in app.header)


def test_용도를_고르면_계약종별_후보가_좁아진다() -> None:
    """좁히기이지 판정이 아니다 (16세션 2절)."""
    from kwise.tariff import load_tariff
    from kwise.tariff.schema import list_contract_types
    from kwise.ui.building import narrow_contract_types

    every = list_contract_types(load_tariff())
    assert len(narrow_contract_types(every, "")) == len(every)
    factory = narrow_contract_types(every, "factory")
    assert factory and all(key.startswith("industrial") for key, _label in factory)
    school = narrow_contract_types(every, "school")
    assert school and all(key.startswith("education") for key, _label in school)
    # 모르는 용도면 전 종별이다 — 고를 것이 사라지면 입력을 못 한다.
    assert narrow_contract_types(every, "없는용도") == every


def test_연면적을_넣으면_원단위가_나온다() -> None:
    """**없으면 줄 자체가 없다.** 국내 평균과 견주지 않는다 (16세션 2절)."""
    from kwise.ui.building import BuildingInfo, intensity_kwh_per_m2

    assert intensity_kwh_per_m2(1_000.0, None) is None
    assert intensity_kwh_per_m2(1_000.0, BuildingInfo(region_key="서울특별시/강남구")) is None
    value = intensity_kwh_per_m2(
        1_000.0, BuildingInfo(region_key="서울특별시/강남구", floor_area_m2=250.0)
    )
    assert value == pytest.approx(4.0)


def test_원단위가_화면에_한_줄로_나온다() -> None:
    from kwise.ui.building import BuildingInfo

    app = _app(building_info=BuildingInfo(region_key="서울특별시/강남구", floor_area_m2=20_000.0))
    assert not app.exception, app.exception
    assert "원단위" in _text(app)
    assert "국내 평균" not in _text(app)


def test_연면적이_없으면_원단위_줄이_없다() -> None:
    app = _app()
    assert "원단위" not in _text(app)


@pytest.mark.parametrize("key", BILLED_MEASURES)
def test_수단을_함께_켜도_카드_값이_불변이다(key: str) -> None:
    """**독립 평가** (14세션 2절) — 다른 수단을 켜고 끄든 이 카드의 숫자가 같다."""
    alone = _stage2_metrics(_app(**_on(key)))
    both = _stage2_metrics(_app(**_on(*BILLED_MEASURES)))
    span = len(alone)
    assert alone
    assert any(both[start : start + span] == alone for start in range(len(both) - span + 1)), (
        key,
        alone,
    )


def test_대표일을_바꾸면_곡선_차트가_따라_바뀐다() -> None:
    """**세 곡선이 같은 날을 본다** (15세션 2절)."""
    app = _app(**_on("power_factor", "ess"))
    assert not app.exception, app.exception
    assert "연간 최대수요일" in _text(app)

    app.selectbox(key="measure_common_ref_day").set_value("winter").run()
    assert not app.exception, app.exception
    body = _text(app)
    assert "겨울 대표일" in body
    assert "연간 최대수요일" not in body


# ===================================================================== 5-1 태양광 경로


@pytest.mark.usefixtures("real_weather")
def test_태양광_계산_전에는_묵은_결과_경고가_없다() -> None:
    app = _app(**_on("solar"))
    assert not app.exception, app.exception
    assert "입력이 변경되었습니다" not in _text(app)


@pytest.mark.usefixtures("real_weather")
def test_입력을_바꾸면_묵은_결과라고_적는다() -> None:
    """계산 단추를 누르기 전 값은 **이전 계산의 것**이다. 그 사실을 적는다 (13세션)."""
    app = _app(**_on("solar"))
    app.button(key="solar_run").click().run(timeout=900)
    assert not app.exception, app.exception
    assert "입력이 변경되었습니다" not in _text(app)

    app.radio(key="measure_solar_azimuth").set_value("east").run(timeout=900)
    assert not app.exception, app.exception
    body = _text(app)
    assert "입력이 변경되었습니다" in body
    assert "묵은 결과" in body


@pytest.mark.usefixtures("real_weather")
def test_계산한_용량을_먼저_적는다() -> None:
    """**① 무엇을 계산했나** (15세션 1-3 · 52세션 1-4). 곡선은 접어 둔다.

    51세션까지는 「**용량 판정** — 설치 가능 면적 전체(160 kWp)를 쓰는 것이
    회수기간 기준 가장 유리합니다」 였다. **거짓이었다** — 회수기간은 더 작은
    용량이 짧고, 동률 처리로 절감액이 큰 줄을 골랐을 뿐이다.
    """
    app = _app(**_on("solar"))
    app.button(key="solar_run").click().run(timeout=900)
    assert not app.exception, app.exception
    body = _text(app)
    assert "로 계산했습니다" in body, body[:400]
    # **거짓 주장이 사라졌다.**
    assert "가장 유리합니다" not in body
    # **라벨과 줄표를 뗐다** — 본문이 빈 줄표로 시작하는 것처럼 보였다.
    assert "**용량 판정** —" not in body
    assert "20단계 상세는 Excel" in body


@pytest.mark.usefixtures("real_weather")
def test_계산_전에는_방위_라벨이_이름뿐이다() -> None:
    """**옆단에서 시군구를 바꾸는 것만으로 여덟 방위를 돌리지 않는다** (17세션 3-1).

    상대 발전량은 0.85초짜리 시뮬레이션 여덟 번이다. 태양광은 「계산」 단추를
    눌러야 도는 카드이고, 그 규약은 라벨에도 적용된다.
    """
    app = _app(**_on("solar"))
    assert not app.exception, app.exception
    options = list(app.radio(key="measure_solar_azimuth").options)
    assert options[0] == "남"
    assert not [item for item in options if re.search(r"[+−]\d+%$", str(item))], options
    assert "「태양광 계산」 을 누른 뒤 라벨에 붙습니다" in _text(app)


@pytest.mark.usefixtures("real_weather")
def test_계산한_뒤에는_방위_라벨에_상대_발전량이_붙는다() -> None:
    """**하드코딩하지 않는다** — 지역·경사각으로 계산한 값이다 (15세션 1-1)."""
    app = _app(**_on("solar"))
    app.button(key="solar_run").click().run(timeout=900)
    assert not app.exception, app.exception
    options = list(app.radio(key="measure_solar_azimuth").options)
    assert options[0] == "남 (기준)"
    assert any(re.search(r"[+−]\d+%$", str(item)) for item in options), options


# ===================================================================== 5-2 산출물


@pytest.fixture(scope="module")
def compare_screen() -> AppTest:
    running = AppTest.from_file(str(APP), default_timeout=900)
    running.session_state["upload_bytes"] = SAMPLE.read_bytes()
    running.session_state["upload_name"] = SAMPLE.name
    running.session_state["contract_form"] = ContractForm(
        contract_type="general_b", voltage="high_a", option="I", contract_kw=6_000.0
    )
    for key in BILLED_MEASURES:
        running.session_state[f"measure_on_{key}"] = True
    # 3단계는 「합산효과 계산」 을 누른 뒤에 그린다 (33세션 5절).
    running.session_state["combination_pick"] = tuple(BILLED_MEASURES)
    return running.run()


def test_엑셀을_만들고_내려받아도_결과가_남는다(compare_screen: AppTest) -> None:
    """**내려받기 단추는 rerun 을 일으킨다** (12세션). 그 rerun 에서도 결과가 남아야 한다."""
    built = compare_screen.button(key="build_excel").click().run(timeout=900)
    assert not built.exception, built.exception
    assert len(built.download_button) == 1

    after = built.download_button(key="dl_excel").click().run(timeout=900)
    assert not after.exception, after.exception
    assert len(after.download_button) == 1, "내려받기 뒤 단추가 사라졌습니다."
    assert _stage3_labels(after)[:3] == ["단순 합", "합산효과", "차이"], (
        "내려받기 뒤 계산 결과가 사라졌습니다."
    )
    assert _stage2_metrics(after), "내려받기 뒤 2단계 카드가 사라졌습니다."


def test_피피티를_만들고_내려받아도_결과가_남는다(compare_screen: AppTest) -> None:
    built = compare_screen.button(key="build_ppt").click().run(timeout=900)
    assert not built.exception, built.exception
    after = built.download_button(key="dl_ppt").click().run(timeout=900)
    assert not after.exception, after.exception
    assert after.download_button(key="dl_ppt"), "내려받기 뒤 단추가 사라졌습니다."
    assert _stage3_labels(after)[:3] == ["단순 합", "합산효과", "차이"]
    assert _stage2_metrics(after), "내려받기 뒤 2단계 카드가 사라졌습니다."


def test_화면_단추는_엑셀과_피피티_둘뿐이다(compare_screen: AppTest) -> None:
    """**Word 는 화면에서 감췄다** (36세션 1절).

    지운 것이 아니다 — 만드는 코드와 시험은 그대로 있고 단추만 없앴다. 산출물이
    셋이면 받는 사람이 「어느 것을 봐야 하나」 를 먼저 골라야 하는데, Word 와
    PPT 는 같은 재료로 같은 이야기를 하는 두 매체라 고를 근거가 없다.
    """
    keys = {button.key for button in compare_screen.button}
    assert "build_excel" in keys
    assert "build_ppt" in keys
    assert "build_word" not in keys, "Word 단추가 화면에 다시 나왔습니다."

    built = compare_screen.button(key="build_ppt").click().run(timeout=900)
    download_keys = {item.key for item in built.download_button}
    assert "dl_word" not in download_keys

    # 단추 이름·탭 이름·캡션 어디에도 남기지 않는다.
    written = _text(compare_screen) + " ".join(str(item.label) for item in compare_screen.button)
    assert "Word" not in written, "화면 문구에 Word 가 남아 있습니다."


def test_워드_생성_코드와_시험이_남아_있다() -> None:
    """**되살릴 수 있어야 한다** (36세션 1절)."""
    from kwise.report.document import build_document, document_bytes, export_document

    assert callable(build_document) and callable(document_bytes) and callable(export_document)
    assert (Path("tests") / "test_document.py").is_file()
    source = (Path("src") / "kwise" / "ui" / "views" / "compare.py").read_text(encoding="utf-8")
    assert "36세션 1절 · Word" in source, "감춘 이유를 코드 주석에 남기십시오."


def _deck_texts(app: AppTest) -> list[str]:
    """만들어 둔 덱의 슬라이드별 글. **화면이 만든 그 바이트를 연다.**"""
    import io

    from pptx import Presentation

    from kwise.ui.artifacts import ARTIFACT_KEY

    store = dict(app.session_state[ARTIFACT_KEY])
    deck = Presentation(io.BytesIO(store["ppt"].payload))
    return [
        "\n".join(
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        )
        for slide in deck.slides
    ]


def test_갑_종별에서도_ESS_장이_선다() -> None:
    """**장이 통째로 빠지고 있었다** (59세션 2절 · PPT 목록 P2·P10).

    56세션이 갑 종별을 훑기 전에 끊은 뒤로 목표가 없어졌는데, 3단계는 목표가
    있을 때만 개략 곡선을 산출물에 넘기고 있었다 — 「성립 불가」 갈래가 열릴
    조건(``ess_curve is not None``)을 못 채워 **켠 수단이 덱에서 사라졌다.**
    """
    app = _app(
        contract_form=ContractForm(
            contract_type="general_a_1", voltage="high_a", option="I", contract_kw=6_000.0
        ),
        **_on("ess"),
    )
    built = app.button(key="build_ppt").click().run(timeout=900)
    assert not built.exception, built.exception
    ess = [text for text in _deck_texts(built) if text.startswith("ESS\n")]
    assert ess, "갑 종별 덱에 ESS 장이 없습니다."
    assert "기본공급약관 제68조" in ess[0], ess[0][:200]
    assert "계약전력" in ess[0]


def test_수단이_없어도_두_산출물이_만들어진다() -> None:
    """진단만 보고 받아 가는 것이 정상 경로다 (8세션에 Excel 에서 잡았다)."""
    app = _app()
    for build, download in (("build_excel", "dl_excel"), ("build_ppt", "dl_ppt")):
        after = app.button(key=build).click().run(timeout=900)
        assert not after.exception, after.exception
        assert after.download_button(key=download)


def test_엑셀_시트가_규약대로다(
    sample_usage: object, sample_bill: object, sample_diagnosis: object
) -> None:
    """시트 이름은 한글, 순서는 :data:`SHEET_ORDER`, tz 는 해제되어 있다."""
    from kwise.report import ReportSections, build_sheets

    sheets = build_sheets(
        ReportSections(
            usage=sample_usage,  # type: ignore[arg-type]
            bill=sample_bill,  # type: ignore[arg-type]
            diagnosis=sample_diagnosis,  # type: ignore[arg-type]
            include_timeseries=False,
        )
    )
    assert list(sheets) == [name for name in SHEET_ORDER if name in sheets]
    assert all(re.search(r"[가-힣]", name) for name in sheets)


def test_산출물_파일명에_날짜_시각이_붙는다() -> None:
    """Excel 이 파일을 열고 있으면 덮어쓰기가 실패한다."""
    from kwise.report import result_path

    assert re.search(r"_\d{8}_\d{4}\.xlsx$", result_path().name)


# ===================================================================== 5-3 경계


def test_데이터가_없어도_탭을_막지_않는다() -> None:
    """**진행 불가 탭도 막지 않고 안내만 낸다** (16세션 1절)."""
    app = _blank()
    assert not app.exception, app.exception
    headers = [str(item.value) for item in app.header]
    assert any("2단계" in item for item in headers)
    assert any("3단계" in item for item in headers)
    assert "「1단계 · 진단」 탭에서" in _text(app)


def test_계약_정보가_없으면_금액을_내지_않는다() -> None:
    running = AppTest.from_file(str(APP), default_timeout=900)
    running.session_state["upload_bytes"] = SAMPLE.read_bytes()
    running.session_state["upload_name"] = SAMPLE.name
    app = running.run()
    assert not app.exception, app.exception
    assert "계약 정보를 확정하기 전까지는" in _text(app)


def test_잉여_활용_카드가_없다() -> None:
    """**41세션에 7.7 을 없앴다.** 잉여는 개선안이 아니라 태양광의 결과다.

    14세션 2-3 의 독립 평가 원칙은 그대로다 — 남은 여섯 카드는 서로 때문에
    비활성이 되지 않는다. 잉여는 카드가 아니므로 잠글 카드도 없고, 켤 수단
    목록에도 없다.
    """
    from kwise.measures import MEASURE_CATALOG

    assert "surplus" not in [item.key for item in MEASURE_CATALOG]
    # **`_on` 은 열쇠를 받는다** (70세션 2절). 여기가 `_on(*MEASURES)` 라
    # `MeasureSpec` 객체가 그대로 들어가 `measure_on_MeasureSpec(...)` 이라는
    # 없는 열쇠가 됐다 — **수단이 하나도 안 켜진 채로 통과하고 있었다.**
    # mypy 가 잡았다.
    app = _app(**_on(*(item.key for item in MEASURES)))
    assert not app.exception, app.exception
    body = _text(app)
    assert "잉여 활용" not in body, body


def test_여유율_입력이_화면에서_사라졌다() -> None:
    """**움직여도 0% 인 입력칸은 고장으로 보인다** (13세션 → 83세션).

    그 값이 판정에 쓰이지도 않는다는 것이 근본이었다 — 걷어냈다. 판정은
    하한이 하고, 왜 줄 것이 없는지를 본문 한 줄이 말한다.
    """
    from kwise.measures.contract import FLOOR_NOT_BINDING_NOTICE

    running = AppTest.from_file(str(APP), default_timeout=900)
    running.session_state["upload_bytes"] = SAMPLE.read_bytes()
    running.session_state["upload_name"] = SAMPLE.name
    running.session_state["contract_form"] = ContractForm(
        contract_type="general_b", voltage="high_a", option="I", contract_kw=5_400.0
    )
    running.session_state["measure_on_contract"] = True
    app = running.run()
    assert not app.exception, app.exception
    assert FLOOR_NOT_BINDING_NOTICE in _text(app)
    assert not [item for item in app.number_input if "여유율" in str(item.label)]


def test_기상_취득에_실패해도_화면이_살아_있다(monkeypatch: pytest.MonkeyPatch) -> None:
    """**차단 등급으로 알리고 화면은 살아 있어야 한다.**

    실패를 모의로 만든다. 환경변수만 비우면 ``PROJECT_CACHE`` 에 남은 앞선
    시험의 결과나 API 가 대신 성공해 실패 경로를 타지 않는다 — 그러면 이 시험이
    조용히 통과한다.
    """
    import streamlit as st

    from kwise.pv import WeatherUnavailableError
    from kwise.ui import pipeline

    def refuse(*_args: object, **_kwargs: object) -> object:
        raise WeatherUnavailableError("기상 자료를 얻지 못했습니다 (모의 실패).")

    monkeypatch.setattr(pipeline, "load_weather", refuse)
    st.cache_data.clear()

    app = _app(**_on("solar"))
    app.button(key="solar_run").click().run(timeout=900)
    assert not app.exception, app.exception
    assert any("기상 자료를 얻지 못해" in str(item.value) for item in app.error)


def test_저부하_평일이_없으면_사유만_적는다(sample_diagnosis: object) -> None:
    """**차트 대신 사유**를 낸다. 없으면 없다고 적는다 (14세션 4절)."""
    from dataclasses import replace

    import pandas as pd

    from kwise.measures import evaluate_demand_response

    profile = sample_diagnosis.dr  # type: ignore[attr-defined]
    empty = replace(
        profile,
        low_load_days=(),
        daily_reducible_kw=pd.Series(dtype=float),
        daily_hours=pd.Series(dtype=float),
        registered_capacity_kw=0.0,
        period_reducible_kwh=0.0,
        annual_reducible_kwh=0.0,
    )
    result = evaluate_demand_response(empty)
    assert result.annual_reducible_kwh == 0.0
    assert not result.has_low_load_days
    assert any("저부하 평일이 없습니다" in item for item in texts(result.notices))


# ===================================================================== 5-4 회귀


def test_요금적용전력_회귀값이_그대로다(sample_diagnosis: object) -> None:
    """15세션은 표시만 고쳤다. **계산값이 바뀌면 그것은 회귀다.**"""
    peak = sample_diagnosis.peak  # type: ignore[attr-defined]
    assert peak.billing_demand_kw == pytest.approx(5_293.44)


def test_3단계_요약이_2단계_카드와_같다() -> None:
    """**재계산하지 않는다** (14세션 5-1)."""
    card = _app(**_on(*BILLED_MEASURES))
    metrics = [(str(item.label), str(item.value)) for item in card.metric]
    index = next(pos for pos, (label, _) in enumerate(metrics) if label == "출력 / 용량")
    ess_payback = metrics[index + 3][1]

    frame = next(item.value for item in card.dataframe if "수단" in list(item.value.columns))
    rows = {str(row["수단"]): row for _, row in frame.iterrows()}
    # 화면 표는 순번으로 적는다 (27세션 2절). 산출물은 절 번호 그대로다.
    assert str(rows["6. ESS"]["회수기간"]) == ess_payback


def test_단순_합과_합산효과가_다르다() -> None:
    """**상호작용이 실제로 계산되는지** — 같으면 조합 재계산이 도는지 의심해야 한다."""
    app = _app(**_on(*BILLED_MEASURES))
    metrics = {str(item.label): str(item.value) for item in app.metric}
    assert metrics["단순 합"] != metrics["합산효과"]
    gap = next(item for item in app.metric if item.label == "차이")
    assert str(gap.delta).endswith("%")


def test_화면에_코드_식별자가_없다() -> None:
    banned = re.compile(r"\b(general_b|high_a|tariff_switch|prior_peaks|verified=)\b")
    app = _app(**_on(*BILLED_MEASURES))
    offenders = [line for line in _text(app).splitlines() if banned.search(line)]
    assert not offenders, offenders


def test_이스케이프되지_않은_물결표가_없다() -> None:
    """물결표 둘이 한 줄에 있으면 그 사이가 취소선이 된다 (13세션)."""
    app = _app(**_on(*BILLED_MEASURES))
    offenders = [
        line for line in _text(app).splitlines() if len(re.findall(r"(?<!\\)~", line)) >= 2
    ]
    assert not offenders, offenders


def test_원_단위_금액이_모두_천의_배수다() -> None:
    """**표시만 절사하고 내부 계산은 원 단위다** (14세션 1절)."""
    pattern = re.compile(r"(?<![만억\d])(\d[\d,]*)원(?![/\w])")
    app = _app(**_on(*BILLED_MEASURES))
    blobs = [_text(app)] + [frame.value.to_string() for frame in app.dataframe]
    amounts = [match.group(1) for blob in blobs for match in pattern.finditer(blob)]
    assert amounts, "원 단위 금액을 찾지 못했습니다."
    offenders = [item for item in amounts if int(item.replace(",", "")) % 1_000]
    assert not offenders, offenders


def test_건물명이_없으면_산출물_제목이_미입력이다() -> None:
    """**빈 표지는 만들다 만 문서로 보인다** (16세션 2절)."""
    from kwise.ui.building import NAME_MISSING, BuildingInfo

    blank = BuildingInfo(region_key="서울특별시/강남구")
    assert blank.title == NAME_MISSING
    assert not blank.named
    named = BuildingInfo(region_key="서울특별시/강남구", name=" 본사 ")
    assert named.title == "본사"

    app = _app(**_on(*BILLED_MEASURES))
    assert not app.exception, app.exception
    body = _text(app)
    assert f"표지 이름 — **{NAME_MISSING}**" in body
    # 같은 것을 두 곳에서 묻지 않는다 — 건물명 입력칸은 옆단 하나뿐이다.
    fields = [str(item.label) for item in app.text_input if "건물명" in str(item.label)]
    assert fields == ["건물명 (선택)"], fields


# ===================================================================== 17세션 · 그래프 표시
#
# **표시만 본다.** 계산값은 5-4 회귀 시험이 지킨다 — 여기서 보는 것은 "축이
# 잘렸는가 · 색이 있는가 · 자리를 옮겼는가" 다.


@functools.cache
def _material() -> tuple[
    UsageData, TariffSwitchResult, PowerFactorResult, RepresentativeDay, BillingResult
]:
    """차트 재료 한 벌. **한 번만 만든다** — 요금 계산이 매번 돌면 시험이 느려진다."""
    from kwise.io import load_usage
    from kwise.measures import evaluate_power_factor, evaluate_tariff_switch
    from kwise.report.days import find_day
    from kwise.tariff import calculate_bill, load_tariff
    from kwise.ui import pipeline

    usage = load_usage(SAMPLE)
    form = ContractForm(
        contract_type="general_b", voltage="high_a", option="I", contract_kw=6_000.0
    )
    table = load_tariff()
    quality = pipeline.load_quality(usage)
    baseline = calculate_bill(
        usage, table, form.selection, options=form.billing_options(), quality=quality
    )
    switch = evaluate_tariff_switch(
        usage, table, form.selection, options=form.billing_options(), quality=quality
    )
    power_factor = evaluate_power_factor(
        usage,
        table,
        form.selection,
        baseline=baseline,
        target_pct=97.0,
        options=form.billing_options(),
        quality=quality,
    )
    day = find_day(usage, "peak")
    return usage, switch, power_factor, day, baseline


def _rows(spec: dict[str, Any], layer: dict[str, Any]) -> list[dict[str, Any]]:
    """altair 는 자료를 ``datasets`` 에 이름으로 담는다. 그것을 풀어 준다.

    **`object` 가 아니라 `Any` 다** (70세션 2절). 이것은 altair 가 낸 JSON 이라
    속이 진짜로 정해져 있지 않다 — `object` 로 적으면 사실과 다르고, 꺼낼
    때마다 억제를 붙이게 된다.
    """
    data = layer.get("data", {})
    values = data.get("values")
    if values is not None:
        return list(values)
    datasets = spec.get("datasets", {})
    return list(datasets.get(data.get("name"), []))


def _fake_generation(usage: object) -> object:
    """정오에 솟는 가짜 발전 프로파일. **기상 없이도 그림을 볼 수 있다.**"""
    import numpy as np
    import pandas as pd

    index = pd.DatetimeIndex(usage.kw.index)  # type: ignore[attr-defined]
    hours = index.hour + index.minute / 60.0
    shape = np.clip(np.cos((hours - 12.5) / 6.0 * np.pi / 2.0), 0.0, None) ** 2
    peak = float(usage.kw.max()) * 0.12  # type: ignore[attr-defined]
    return pd.Series(shape * peak, index=index, name="발전량(kW)")


def test_선택요금이_제도_순서로_나온다() -> None:
    """**Ⅰ·Ⅱ·Ⅲ 순이다** (17세션 1-1). 절감액 순으로 정렬하면 자료마다 뒤섞인다."""
    from kwise.tariff import option_sort_key
    from kwise.ui.charts import tariff_option_frame

    assert sorted(["III", "I", "II"], key=option_sort_key) == ["I", "II", "III"]
    _usage, switch, _pf, _day, _base = _material()
    assert list(tariff_option_frame(switch)["요금제"]) == ["선택Ⅰ", "선택Ⅱ", "선택Ⅲ"]


def test_선택Ⅲ도_기본_전력량으로_갈라진다() -> None:
    """**「상세 미산출」 은 값이 없어서가 아니었다** (17세션 1-4).

    현행·최적 둘만 상세를 내던 최적화가 원인이다 — 갈아탈 수 있는 조합은 같은
    계약종별·전압 안의 선택요금뿐이라 많아야 셋이고, 늘어나는 계산은 한 벌이다.
    """
    from kwise.ui.charts import tariff_option_frame, tariff_option_long_frame

    _usage, switch, _pf, _day, _base = _material()
    frame = tariff_option_frame(switch)
    assert frame["기본요금(원)"].notna().all(), frame.to_string()
    assert frame["전력량요금(원)"].notna().all(), frame.to_string()
    assert set(tariff_option_long_frame(switch)["구분"]) == {"기본요금", "전력량요금", "합계"}


def test_요금제_그래프가_그룹_막대이고_차액_차트가_따로_있다() -> None:
    """쌓으면 항목별 비교가 안 된다 (17세션 1-2·1-3)."""
    from kwise.ui.charts import tariff_delta_chart, tariff_option_chart

    _usage, switch, _pf, _day, _base = _material()
    grouped = tariff_option_chart(switch).to_dict()
    assert "xOffset" in grouped["encoding"], list(grouped["encoding"])
    assert grouped["encoding"]["y"]["scale"]["zero"] is False
    # **축 제목에 그 사실을 적지 않는다** (21세션 5절). 눈금을 보면 안다.
    assert grouped["encoding"]["y"]["title"] == "요금 (원)"

    delta = tariff_delta_chart(switch).to_dict()
    fields = {layer["encoding"]["x"].get("field") for layer in delta["layer"]}
    assert "현행 대비(원)" in fields


def test_화면에_상세_미산출이_없다() -> None:
    app = _app(**_on("tariff_switch"))
    assert not app.exception, app.exception
    body = _text(app)
    assert "상세 미산출" not in body
    # **읽는 법은 툴팁 하나로 족하다** (27세션 4-2). 본문 캡션이 같은 말을 했다.
    assert "왼쪽(초록)이 절감" not in body


def test_역률_범례가_그림_아래에_있다() -> None:
    """**꼭짓점 옆 글자와 자리를 다투지 않는다** (23세션 1절 → 27세션 6절).

    바깥 오른쪽 범례는 그림 오른쪽 **위**에서부터 쌓인다. 이 그림은 도형 옆에
    설명을 직접 적으므로 (17세션 2절) 그 글자가 같은 자리로 뻗었다. 배경 없음·
    안쪽 금지는 그대로다.
    """
    from kwise.ui.charts import power_triangle_chart

    _usage, _switch, power_factor, _day, _base = _material()
    spec = power_triangle_chart(power_factor).to_dict()
    legends = [
        layer["encoding"]["color"]["legend"]
        for layer in spec["layer"]
        if "legend" in layer["encoding"].get("color", {})
    ]
    assert legends, "범례가 없습니다."
    assert {item["orient"] for item in legends} == {"bottom"}
    assert all(item.get("fillColor") is None for item in legends), "범례에 배경이 있습니다."


def test_전력삼각형에_각도와_역률을_직접_적는다() -> None:
    """**범례 의존을 줄인다** (17세션 2절)."""
    from kwise.ui.charts import power_triangle_chart

    _usage, _switch, power_factor, _day, _base = _material()
    spec = power_triangle_chart(power_factor).to_dict()
    texts = [layer for layer in spec["layer"] if layer.get("mark", {}).get("type") == "text"]
    values = [row for layer in texts for row in _rows(spec, layer)]
    assert any(
        "역률" in str(row.get("설명", "")) and "°" in str(row.get("설명", "")) for row in values
    )
    assert any(str(row.get("각도라벨", ""))[-1:] == "°" for row in values)


@pytest.mark.usefixtures("real_weather")
def test_참고는_선정_용량과_다를_때만_적는다() -> None:
    """**② 참고** (17세션 3-2 → 52세션 1-4).

    이 자료는 권장이 곧 면적 상한이고 잉여도 나지 않는다 — **적을 참고가 없다.**
    51세션까지는 그 자리에 「회수기간을 기준으로 고른 용량입니다. 그 용량에서
    멈춘 것은 **설치 가능 면적 상한** 입니다」 를 늘 적었고, 배지가 문장 중간에
    박혀 읽기 어려웠다.
    """
    from kwise.measures import RECOMMENDED

    app = _app(**_on("solar"))
    app.button(key="solar_run").click().run(timeout=900)
    assert not app.exception, app.exception
    body = _text(app)
    assert "참고 —" not in body, body[:400]
    # **배지가 문장 중간에 박히지 않는다.**
    assert "그 용량에서 멈춘 것은" not in body
    assert "회수기간을 기준으로 고른 용량입니다" not in body
    assert "최적" not in body, "개선안 맥락의 「최적」 이 화면에 남아 있다"
    # 권장은 표의 표식이 낸다 — 단가를 안 넣은 자료라 「최대 절감액」 이다.
    assert RECOMMENDED not in body


@pytest.mark.usefixtures("real_weather")
def test_태양광_용량_표가_잉여_지점을_세운다() -> None:
    """20단계는 Excel 로, 화면에는 의미 있는 지점만 (17세션 3-3 · 31세션 4-1).

    **잉여 지점 둘이 표에 있어야 한다.** 26세션에 판단을 잉여로 옮겼는데 표는
    선정 용량 아래에서만 뽑아, 정작 「어디서부터 잉여가 생기나」 가 없었다.
    """
    from kwise.measures import (
        AREA_EXCEEDED,
        LARGEST_SAVING,
        RECOMMENDED,
        SELECTED_CAPACITY,
        SURPLUS_HEAVY,
        SURPLUS_ONSET,
    )
    from kwise.ui.charts import CAPACITY_ROWS

    app = _app(**_on("solar"))
    app.button(key="solar_run").click().run(timeout=900)
    assert not app.exception, app.exception
    frame = next(item.value for item in app.dataframe if "자가소비율" in list(item.value.columns))
    assert len(frame) == CAPACITY_ROWS == 5
    assert set(frame.columns) >= {
        "용량(kWp)",
        "필요 면적",  # 31세션 4-1 — 면적이 판단 기준이다
        "발전량",  # 26세션 3-3 — MWh/년 로 낸다
        "자가소비율",
        "절감액",  # 51세션 2절 — 절감 열은 하나다. 줄을 가르는 것이 이 값이다
        "투자비",
        "회수기간",
    }
    # **가른 열 둘은 사라졌다** (51세션 2절). 분해는 카드 툴팁·계산 근거·Excel 에 있다.
    assert "기본요금 절감" not in frame.columns
    assert "전력량요금 절감" not in frame.columns
    marks = " ".join(str(value) for value in frame["표식"])
    assert SELECTED_CAPACITY in marks
    assert SURPLUS_ONSET in marks
    assert SURPLUS_HEAVY in marks
    # 설치 가능 면적을 넘는 지점은 값을 지우지 않고 **그 사실만** 적는다.
    assert AREA_EXCEEDED in marks
    # **각주가 말하는 표식이 표에 있어야 한다** (51세션 1절). 50세션은 권장이
    # 면적 상한과 같으면 「선정 용량」 에 먹혀 「권장」 이 화면에서 사라졌고,
    # 각주는 조건 없이 붙어 「최대 절감액」 인 화면에서도 「권장」 을 말했다.
    #
    # **이 시험은 단가를 넣지 않는다** — 그래서 고른 자리의 이름은 「최대 절감액」
    # 이고 동률 각주는 나오지 않아야 한다.
    assert LARGEST_SAVING in marks, marks
    assert RECOMMENDED not in marks, marks
    caption = " ".join(str(item.value) for item in app.caption)
    assert RECOMMENDED not in caption, "쓰이지 않은 규칙을 각주로 적으면 안 된다"


def test_태양광_연간_차트가_발전량만_그린다() -> None:
    """**한 그림은 한 가지만 말한다** (17세션 3-4 → 23세션 5절).

    17세션에는 「사용량이 줄어드는 모습」이 주인공이었다. 뜻은 옳았지만 사용량이
    일 60 MWh 대인데 발전량은 3 MWh 대라 **저감분 띠가 선 굵기만큼도 서지
    않았다.** 사용량이 얼마나 줄었는지는 카드의 절감액과 대표일 곡선이 낸다.
    """
    from kwise.ui.charts import solar_annual_chart, solar_saving_ratio

    usage, _switch, _pf, _day, _base = _material()
    generation = _fake_generation(usage)
    spec = solar_annual_chart(usage, generation).to_dict()
    assert spec["encoding"]["y"]["field"] == "발전량(kWh)"
    assert [item["field"] for item in spec["encoding"]["tooltip"]] == ["날짜", "발전량(kWh)"]
    # 절감 비율은 여전히 화면 문구가 쓴다 — 그림에서 뺐을 뿐이다.
    ratio = solar_saving_ratio(usage, generation)
    assert ratio is not None and 0.0 < ratio < 1.0


def test_일일_곡선의_축이_0부터가_아니고_저감분을_채운다() -> None:
    """5,000 kW 부하에 수백 kW 를 얹으면 0 부터 그린 축에서 두 선이 붙는다 (3-5)."""
    from kwise.ui.charts import PEAK_ZOOM_HOURS, solar_day_chart

    usage, _switch, _pf, day, _base = _material()
    generation = _fake_generation(usage)
    spec = solar_day_chart(usage, generation, day).to_dict()
    scales = [
        layer["encoding"]["y"].get("scale", {})
        for layer in spec["layer"]
        if "y" in layer["encoding"]
    ]
    assert any(scale.get("zero") is False for scale in scales), scales
    filled = [
        layer
        for layer in spec["layer"]
        if layer.get("mark", {}).get("type") == "area" and "y2" in layer["encoding"]
    ]
    assert filled, "원부하와 순부하 사이를 채우지 않았습니다."
    texts = [layer for layer in spec["layer"] if layer.get("mark", {}).get("type") == "text"]
    values = [row for layer in texts for row in _rows(spec, layer)]
    assert any("피크 −" in str(row.get("설명", "")) for row in values), values

    zoomed = solar_day_chart(usage, generation, day, zoom=True).to_dict()
    assert f"피크 앞뒤 {PEAK_ZOOM_HOURS}시간" in str(zoomed)


def test_ESS_그래프가_한_칸이고_충방전은_문구다() -> None:
    """**그림이 둘일 필요가 없다** (17세션 4-1 → 23세션 6절).

    아래 칸(충전＋·방전−)은 위 칸의 결과라 종속이다. 시각과 양은 글이 더
    정확하다 — 막대에서 시각을 눈으로 읽어 내야 했다.
    """
    from kwise.report.frames import dispatch_schedule, ess_day_frame
    from kwise.ui import charts as charts_module
    from kwise.ui.charts import ess_day_chart

    usage, _switch, _pf, day, dispatch = _dispatch()
    spec = ess_day_chart(usage, dispatch, day).to_dict()
    assert "vconcat" not in spec, "2단 그림이 남아 있습니다."
    assert any(
        layer["encoding"]["y"].get("scale", {}).get("zero") is False
        for layer in spec["layer"]
        if "y" in layer["encoding"]
    )
    # **배경 띠를 걷어냈다** (26세션 2-1·2-2). 확대한 창이 한 시간대 안에 들어가
    # 그림 전체가 주황 한 색이 되고, 범례에는 셋이 남아 그림에 없는 것을 가리켰다.
    assert "시간대" not in str(spec), "계시별 시간대 배경 띠가 남아 있습니다."
    assert not hasattr(charts_module, "_BAND_COLORS"), "배경 띠 색이 남아 있습니다."
    # 충·방전 시각은 문구가 낸다.
    charge, discharge = dispatch_schedule(ess_day_frame(usage, dispatch, day.date))
    assert charge or discharge, "충·방전 구간을 못 읽었습니다."


def test_ESS_본문에_투자비_상세와_성립_조건이_없다() -> None:
    """본문에는 투자비 합계·절감액·회수기간만 남긴다 (17세션 4-2 · 20세션).

    화면이 다시 적던 넉 줄은 지웠다. 계산 쪽이 같은 사실을 근거로 이미 내고
    있어 사실 ID 로 견주니 전부 중복이었다.
    """
    source = (Path("src") / "kwise" / "ui" / "views" / "measures.py").read_text(encoding="utf-8")
    body = source[source.index("def _ess(") : source.index("def _ess_cost_inputs(")]
    # 주석은 화면에 나가지 않는다.
    body = " ".join(line for line in body.splitlines() if not line.lstrip().startswith("#"))
    for banned in ("설비 **", "성립 조건", "kW당 배터리비"):
        assert banned not in body, banned
    # **정밀화가 낸 안내도 같은 자리로 간다** (40세션 1-2). 창 가장자리 경고가
    # 여기 실려야 「그 자리가 최적」 으로 읽히지 않는다.
    assert "_notices((*result.notices, *optimum.notices))" in source
    # **곡선 대비 확인사항은 46세션에 뺐다.** 곡선이 없어져 설명할 차이가 없다.
    assert "_ess_basis_note" not in source
    assert "_ess_details" not in source, "화면이 계산 쪽 근거를 다시 적고 있습니다."


def test_화면에_조달_사례_표가_없다() -> None:
    """계수 산출 근거는 데이터와 재적합 스크립트에 남긴다 (17세션 4-3)."""
    source = (Path("src") / "kwise" / "ui" / "views" / "measures.py").read_text(encoding="utf-8")
    assert "case_table()" not in source
    app = _app(**_on("ess"))
    assert not app.exception, app.exception
    columns = {str(name) for item in app.dataframe for name in item.value.columns}
    assert "설치 형태" not in columns, columns


def test_ESS_투자비_상세가_근거_툴팁에_있다() -> None:
    """**확인사항 상자에서 근거 툴팁으로 옮겼다** (19세션 1절).

    17세션에 본문에서 확인사항으로 내렸던 넉 줄이 이제 툴팁으로 간다. 스물이
    넘는 줄을 접힌 상자에 쌓으면 정작 위험한 두 건이 묻히기 때문이다. 지운
    것이 아니라 자리를 옮긴 것이므로 **툴팁에는 그대로 있어야 한다.**

    20세션에 문구가 하나로 합쳐졌다 — 화면이 다시 적던 「투자비 내역」 대신
    계산 쪽 「투자비 = 설비 … + 전기공사 …」 하나가 남는다. 같은 사실이다.
    """
    app = _app(**_on("ess"))
    assert not app.exception, app.exception
    tips = "\n".join(str(item.help or "") for item in app.caption)
    body = "\n".join(str(item.value) for item in app.markdown)

    # 산식·계수·내역은 **근거**다 — 툴팁으로 접힌다.
    assert "투자비 = 설비" in tips
    assert "설비비 = " in tips
    assert "투자비 = 설비" not in body

    # **성립 조건은 다르다.** 이 자료에서는 성립하지 않아 주의로 올라오므로
    # 본문에 남는다 — 등급이 자리를 정한다는 것이 이 한 줄로 드러난다.
    assert "회수에 필요한 저감량" in body


@functools.cache
def _dispatch() -> tuple[
    UsageData, TariffSwitchResult, PowerFactorResult, RepresentativeDay, DispatchResult
]:
    """ESS 디스패치 한 벌. 그림만 보므로 목표는 관측 최대의 97% 로 잡는다."""
    from kwise.measures import EssCostInput, evaluate_ess

    usage, switch, power_factor, day, baseline = _material()
    result = evaluate_ess(
        usage,
        _table(),
        _contract().selection,
        target_kw=float(usage.kw.max()) * 0.97,
        cost=EssCostInput(),
        baseline=baseline,
        options=_contract().billing_options(),
    )
    return usage, switch, power_factor, day, result.dispatch


def _contract() -> ContractForm:
    return ContractForm(
        contract_type="general_b", voltage="high_a", option="I", contract_kw=6_000.0
    )


@functools.cache
def _table() -> TariffTable:
    from kwise.tariff import load_tariff

    return load_tariff()
