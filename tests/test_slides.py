"""PPT 보고서 시험 (36세션).

**여기서 지키는 것 여덟**

    ① 슬라이드 차례가 지시서의 목차 그대로다
    ② 수단 0개여도 만들어진다 — **수단별 장만 빠진다**
    ③ 색·크기가 **가이드 한 곳**에서 온다 (하드코딩 검사)
    ④ 슬라이드마다 **시각 요소가 하나 이상** 있다
    ⑤ 레이아웃이 반복되지 않는다 (형태별 개수)
    ⑥ 부록에서 뺀 셋이 PPT 에 없고 **Word 에는 있다**
    ⑦ 한글이 깨지지 않고 **OS 전용 글꼴 이름이 코드에 없다**
    ⑧ 파일명에 날짜·시각 접미사가 붙는다
"""

from __future__ import annotations

import datetime as dt
import re
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from kwise.compare import ComparisonResult
from kwise.diagnose import Diagnosis
from kwise.io import UsageData
from kwise.measures import MEASURE_CATALOG, TariffSwitchResult
from kwise.report import figures
from kwise.report.appendix import APPENDIX_TITLES
from kwise.report.design import DesignGuideError, design_path, load_design_guide
from kwise.report.document import (
    DocumentSections,
    MeasureEntry,
    build_document,
    measure_entries,
)
from kwise.report.slides import (
    ANNUAL_BASIS_NOTE,
    APPENDIX_ROW_LIMIT,
    LAYOUTS,
    MEASURE_AGENDA_ITEM,
    NEXT_STEPS,
    NEXT_STEPS_HEADLINE,
    SLIDE_TITLES,
    _cautions,
    agenda_items,
    appendix_pages,
    build_slides,
    export_slides,
    measure_slide_title,
    plain_text,
    season_pairs,
    slide_specs,
    slides_bytes,
    slides_path,
    split_reason,
)
from kwise.tariff import BillingResult, TariffTable

#: DrawingML 이름공간. 글꼴 지정이 이 안에 있다.
_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

PROJECT_ROOT = Path(__file__).resolve().parent.parent
SRC_ROOT = PROJECT_ROOT / "src" / "kwise"

#: 지시서 2절의 괄호 — **하나가 슬라이드 한 장이다.**
EXPECTED_ORDER: tuple[str, ...] = (
    SLIDE_TITLES["cover"],
    SLIDE_TITLES["agenda"],
    SLIDE_TITLES["building"],
    SLIDE_TITLES["usage_pattern"],
    SLIDE_TITLES["peak_summary"],
    SLIDE_TITLES["peak_detail"],
    SLIDE_TITLES["structure"],
    SLIDE_TITLES["measure_summary"],
    SLIDE_TITLES["combination"],
    SLIDE_TITLES["appendix"],
    SLIDE_TITLES["closing"],
)

#: 부록 앞의 뼈대 장 수 (표지·목차·진단 다섯·요약·조합).
FRAME_BEFORE_APPENDIX = 9


# ===================================================================== 도우미


@pytest.fixture(scope="module")
def entries(
    sample_switch: TariffSwitchResult,
    sample_bill: BillingResult,
    sample_usage: UsageData,
) -> tuple[object, ...]:
    from kwise.measures import evaluate_contract_adjustment

    contract = evaluate_contract_adjustment(sample_usage, sample_bill, contract_kw=5_800.0)
    return measure_entries(switch=sample_switch, contract=contract)


@pytest.fixture(scope="module")
def full_sections(
    sample_usage: UsageData,
    sample_bill: BillingResult,
    sample_diagnosis: Diagnosis,
    sample_comparison: ComparisonResult,
    sample_switch: TariffSwitchResult,
    entries: tuple[object, ...],
) -> DocumentSections:
    from kwise.report.worksheet import tariff_switch_worksheet

    return DocumentSections(
        usage=sample_usage,
        bill=sample_bill,
        diagnosis=sample_diagnosis,
        comparison=sample_comparison,
        measures=entries,  # type: ignore[arg-type]
        worksheets=(tariff_switch_worksheet(sample_switch),),
        building_name="본사 사옥",
        prepared_on=dt.date(2026, 8, 11),
    )


@pytest.fixture(scope="module")
def diagnosis_only(
    sample_usage: UsageData, sample_bill: BillingResult, sample_diagnosis: Diagnosis
) -> DocumentSections:
    """**수단을 하나도 켜지 않은 덱.** 진단만 보고 받아 가는 정상 경로다."""
    return DocumentSections(usage=sample_usage, bill=sample_bill, diagnosis=sample_diagnosis)


def _slide_text(slide: object) -> str:
    parts: list[str] = []
    for shape in slide.shapes:  # type: ignore[attr-defined]
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
        if shape.has_table:
            parts.extend(
                cell.text
                for row in shape.table.rows
                for cell in row.cells  # type: ignore[union-attr]
            )
    return "\n".join(parts)


def _deck_text(deck: object) -> str:
    return "\n".join(_slide_text(slide) for slide in deck.slides)  # type: ignore[attr-defined]


def _all_measures(sections: DocumentSections) -> tuple[MeasureEntry, ...]:
    """수단을 **모두 켠** 항목 목록. 41세션부터 여섯이다.

    실제로 여섯을 계산하면 시험이 몇 분 늘어난다. 여기서 보는 것은 **장 수가
    수단 수를 따라가는가** 하나이므로 항목만 채워 만든다 — 카탈로그 차례를
    그대로 쓰므로 실제 경로와 같은 순서다.
    """
    known = {entry.kind.key: entry for entry in sections.measures}
    return tuple(
        known.get(
            kind.key,
            MeasureEntry(
                kind=kind,
                conclusion=f"{kind.title} 검토 결과입니다.",
                saving="0원",
                investment="0원",
                payback="즉시",
                certainty="높음",
            ),
        )
        for kind in MEASURE_CATALOG
    )


def _slide_by_key(deck: object, sections: DocumentSections, key: str) -> object:
    """자리표의 열쇠로 실물 슬라이드를 찾는다.

    **차례에서 세지 않는다.** 「맨 뒤」 같은 자리 표현은 장이 하나 붙는 순간
    조용히 다른 슬라이드를 가리킨다 — 37세션에 마무리를 붙이며 겪었다.
    """
    index = [spec.key for spec in slide_specs(sections)].index(key)
    return list(deck.slides)[index]  # type: ignore[attr-defined]


def _visual_shapes(slide: object) -> list[object]:
    """**글자 상자가 아닌 것.** 그림·표·선이 시각 요소다 (36세션 5절)."""
    visual = {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.LINE}
    return [shape for shape in slide.shapes if shape.shape_type in visual]  # type: ignore[attr-defined]


# ===================================================================== ① 차례


def test_슬라이드_구성이_목차대로다(full_sections: DocumentSections) -> None:
    """**괄호 하나가 슬라이드 한 장이다** (36세션 2절)."""
    titles = [spec.title for spec in slide_specs(full_sections)]
    measures = [measure_slide_title(entry) for entry in full_sections.measures]
    assert measures, "이 시험은 수단이 켜진 덱을 본다."
    appendix = [page.title for page in appendix_pages(full_sections)]
    expected = [
        *EXPECTED_ORDER[:8],
        *measures,  # 검토한 수단별 1장씩 — **켠 것만, 차례대로**
        SLIDE_TITLES["combination"],
        *appendix,  # 부록은 수단마다 한 장 이상 (39세션 5절)
        SLIDE_TITLES["closing"],
    ]
    assert expected[-1] == SLIDE_TITLES["closing"], "마무리는 맨 뒤다."
    assert titles == expected


def test_실물_덱이_차례와_같다(full_sections: DocumentSections) -> None:
    """자리표와 실물이 갈라지지 않는다 — 덱은 :func:`slide_specs` 대로 만들어진다."""
    deck = build_slides(full_sections)
    specs = slide_specs(full_sections)
    assert len(deck.slides) == len(specs)
    for slide, spec in zip(deck.slides, specs, strict=True):
        assert spec.title in _slide_text(slide), f"「{spec.title}」 장이 제목을 잃었습니다."


def test_목차가_수단을_한_줄로_묶는다(full_sections: DocumentSections) -> None:
    """슬라이드마다 한 줄을 적으면 수단이 일곱일 때 목차가 열두 줄이 된다.

    **수단 이름을 나열하지도 않는다** (38세션 1-1). 어느 수단을 보았는지는
    바로 다음 장인 「개선안별 요약」 표가 낸다.
    """
    items = agenda_items(full_sections)
    assert len(items) == len(EXPECTED_ORDER) - 2 + 1  # 표지·목차를 빼고 수단 묶음 한 줄
    assert MEASURE_AGENDA_ITEM in items
    for entry in full_sections.measures:
        assert measure_slide_title(entry) not in " ".join(items), (
            "목차에 수단을 나열하면 줄이 넘쳐 다음 항목을 덮는다."
        )


def test_목차_항목이_한_줄을_넘지_않는다(full_sections: DocumentSections) -> None:
    """**07번이 08번을 덮고 있었다** (38세션 1-1).

    목차는 번호 칸(0.8in)을 뺀 폭에 한 줄로 앉는다. 두 줄이 되면 아래 항목의
    자리를 침범하는데, 줄 간격이 항목 수로 정해져 있어 **넘친 만큼 겹친다.**

    수단이 일곱일 때가 가장 긴 경우다 — 그때도 넘치지 않아야 한다.
    """
    from dataclasses import replace

    from kwise.report.design import load_design_guide
    from kwise.report.slides import _text_width_in

    guide = load_design_guide()
    span = guide.slide.content_width_in - 0.8
    every = replace(full_sections, measures=_all_measures(full_sections))
    for sections in (full_sections, every, diagnosis_only_of(full_sections)):
        for line in agenda_items(sections):
            width = _text_width_in(line, guide.type_scale.card_title)
            assert width <= span, (
                f"목차 「{line}」 가 한 줄을 넘습니다 ({width:.2f}in > {span:.2f}in)"
            )


def diagnosis_only_of(sections: DocumentSections) -> DocumentSections:
    """수단을 뺀 같은 자료. 목차가 가장 짧은 경우다."""
    from dataclasses import replace

    return replace(sections, measures=())


# ===================================================================== 38세션 · 편집 결함


#: PPT 가 해석하지 못하는 마크다운 표식 (38세션 1-2).
#: 화면(Streamlit)은 이것을 굵게 그리지만 PowerPoint 는 글자 그대로 찍는다.
MARKDOWN_MARKS: tuple[str, ...] = ("**", "__", "`")


def test_마크다운_기호가_슬라이드에_남지_않는다(loaded_sections: DocumentSections) -> None:
    """**PPT 는 마크다운을 해석하지 않는다** (38세션 1-2).

    「자격요건은 판정하지 않았습니다」 를 감싼 별표가 슬라이드에 그대로
    찍히고 있었다. 문구는 화면·Word 로도 가므로 낳는 자리에서 고칠 수 없어,
    **적는 순간 벗긴다** — :func:`~kwise.report.slides.plain_text` 한 곳이다.
    """
    text = _deck_text(build_slides(loaded_sections))
    for mark in MARKDOWN_MARKS:
        assert mark not in text, f"슬라이드에 마크다운 표식 「{mark}」 이 남아 있습니다."


def test_마크다운을_벗기는_자리가_하나다() -> None:
    """**자리마다 벗기면 새 문구를 붙일 때 빠뜨린다.**

    글자를 쓰는 곳은 :func:`_text` 와 :func:`_table` 둘뿐이고, 둘 다
    ``plain_text`` 를 지나야 한다.
    """
    assert plain_text("**굵게** 와 `코드` 와 __밑줄__") == "굵게 와 코드 와 밑줄"
    source = (SRC_ROOT / "report" / "slides.py").read_text(encoding="utf-8")
    writes = [line for line in source.splitlines() if "run.text =" in line]
    assert writes, "글자를 쓰는 자리를 찾지 못했습니다."
    for line in writes:
        assert "plain_text(" in line, f"마크다운을 벗기지 않는 자리가 있습니다: {line.strip()}"


def test_절_번호가_슬라이드에_없다(loaded_sections: DocumentSections) -> None:
    """**7.1~7.7 을 뺀다** (38세션 1-3).

    「7.」 이 무엇인지 덱 어디에도 적혀 있지 않다 — 처음 받아 보는 사람에게는
    없는 7장을 찾게 만드는 표시다. 화면은 27세션에 이미 뗐다.
    """
    sections = loaded_sections
    text = _deck_text(build_slides(sections))
    for kind in MEASURE_CATALOG:
        assert kind.number not in text, f"절 번호 「{kind.number}」 이 슬라이드에 남아 있습니다."
    # **이름은 남는다.** 번호만 뗀 것이지 수단을 감춘 것이 아니다.
    for entry in sections.measures:
        assert entry.kind.label in text


def test_Word_에는_절_번호가_남아_있다(loaded_sections: DocumentSections) -> None:
    """**PPT 에서만 뗐다.** Word·Excel 은 요구사항서와 맞물린 번호를 쓴다."""
    document = build_document(loaded_sections)
    text = "\n".join(item.text for item in document.paragraphs)
    for entry in loaded_sections.measures:
        assert entry.kind.number in text, f"Word 에서 「{entry.kind.number}」 이 사라졌습니다."


# ===================================================================== 38세션 2절 · 4~6장


def test_4장이_전력사용현황과_부하패턴을_합친다(full_sections: DocumentSections) -> None:
    """**화면 1단계의 차례를 그대로 따른다** (38세션 2-1).

    화면은 머릿수 지표 자리에 그림이 없고 첫 그림이 「부하 패턴」 절에 나온다.
    36세션의 덱은 그 구조와 어긋나 두 장이 같은 이야기를 지표만 바꿔 되풀이했다.
    """
    specs = slide_specs(full_sections)
    frame = [spec for spec in specs if spec.measure is None]
    assert [spec.key for spec in frame[2:6]] == [
        "building",
        "usage_pattern",
        "peak_summary",
        "peak_detail",
    ]
    slide = _slide_by_key(build_slides(full_sections), full_sections, "usage_pattern")
    text = _slide_text(slide)
    assert SLIDE_TITLES["usage_pattern"] == "전력사용현황 및 부하패턴"
    for label in ("부하율", "기저부하 비율", "운영시간 외 부하 비중"):
        assert label in text, f"「{label}」 지표가 4장에 없습니다."
    assert "사용량" in text
    # **겹치는 것은 뺐다** — 최대수요·요금적용전력은 다음 장이 낸다.
    assert "최대수요" not in text and "요금적용전력" not in text
    assert "분석 기간" not in text, "3장 표에 이미 있는 값이다."


def test_4장_그림이_기온을_겹쳐_그린다(full_sections: DocumentSections) -> None:
    """**냉난방이 부하의 얼마를 차지하는지**가 태양광·ESS 판단을 가른다.

    기온이 없으면 사용량만 그리고 캡션이 그 사실을 적는다 — 화면과 같은 규칙이다.
    """
    from dataclasses import replace

    import pandas as pd

    slide = _slide_by_key(build_slides(full_sections), full_sections, "usage_pattern")
    assert "일평균 기온" not in _slide_text(slide), "기온이 없으면 사용량만 그린다."

    index = pd.date_range(full_sections.usage.meta.start, periods=48, freq="h")
    warm = replace(full_sections, temperature=pd.Series(range(48), index=index, dtype=float))
    assert "일평균 기온" in _slide_text(_slide_by_key(build_slides(warm), warm, "usage_pattern"))


def test_5장이_정오_비중을_낸다(full_sections: DocumentSections) -> None:
    """**정오 비중이 태양광 판정의 근거다** (38세션 2-2).

    36세션의 덱에는 이 숫자가 없어 상위 구간 그래프만 덩그러니 서 있었다.
    """
    slide = _slide_by_key(build_slides(full_sections), full_sections, "peak_summary")
    text = _slide_text(slide)
    assert "상위 구간 정오 비중" in text
    assert "상위 구간 주말 비중" in text
    assert "요금적용전력" in text


def test_5장_지표가_화면과_같은_갈림이다(full_sections: DocumentSections) -> None:
    """관측 최대와 요금적용 대상 최대가 같으면 **한 칸으로 접는다** (13세션).

    갈릴 때는 넷이 된다 — 화면은 칸이 셋이라 정오 비중을 밀어냈지만 슬라이드는
    밀어낼 것이 없다.
    """
    from dataclasses import replace

    from kwise.report.slides import _peak_stats

    diagnosis = full_sections.diagnosis
    assert diagnosis is not None
    joined = [label for label, _value in _peak_stats(full_sections)]
    assert joined[0] == "최대수요 = 요금적용전력"
    assert len(joined) == 3

    night = replace(diagnosis, peak=replace(diagnosis.peak, billing_demand_kw=1_000.0))
    split = [label for label, _value in _peak_stats(replace(full_sections, diagnosis=night))]
    assert split[:2] == ["관측 최대수요", "요금적용전력"]
    assert "상위 구간 정오 비중" in split
    assert len(split) == 4


def test_6장이_그림_둘을_좌우로_놓는다(full_sections: DocumentSections) -> None:
    """지금 구성 그대로다 — 시간대별 평균 부하와 상위 구간 발생 시각."""
    sections = full_sections
    spec = next(item for item in slide_specs(sections) if item.key == "peak_detail")
    assert spec.layout == "chart_pair"
    slide = _slide_by_key(build_slides(sections), sections, "peak_detail")
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 2
    text = _slide_text(slide)
    assert "평균 부하" in text and "발생한 시각" in text


# ===================================================================== ② 수단 0개


def test_수단이_없어도_만들어진다(diagnosis_only: DocumentSections) -> None:
    """**수단별 장만 빠진다.** 나머지는 그대로다 (Word 와 같은 규약)."""
    specs = slide_specs(diagnosis_only)
    assert [spec.title for spec in specs] == list(EXPECTED_ORDER)
    assert all(spec.measure is None for spec in specs)

    deck = build_slides(diagnosis_only)
    assert len(deck.slides) == len(EXPECTED_ORDER)
    payload, _name = slides_bytes(diagnosis_only)
    assert payload[:2] == b"PK", "pptx 가 아닙니다."


def test_수단_0개면_개선안별_요약이_비어_있음을_적는다(
    diagnosis_only: DocumentSections,
) -> None:
    """**빈 표를 남기지 않는다.** 빈칸은 「효과가 없다」 로 읽힌다."""
    slide = _slide_by_key(build_slides(diagnosis_only), diagnosis_only, "measure_summary")
    assert "검토한 수단이 없습니다" in _slide_text(slide)


# ===================================================================== ③ 값이 오는 자리


def test_가이드가_색과_크기를_쥔다() -> None:
    """**데이터 파일 한 곳에서 온다.** 코드에 기본값을 두지 않는다."""
    guide = load_design_guide()
    assert design_path().is_file()
    assert guide.slide.width_in == pytest.approx(13.333)
    assert guide.slide.height_in == pytest.approx(7.5)
    assert guide.slide.margin_in == pytest.approx(0.5)
    assert 0.3 <= guide.slide.block_gap_in <= 0.4
    assert 0.25 <= guide.slide.title_gap_in <= 0.3
    assert guide.slide.text_slack == pytest.approx(0.1), "텍스트 상자 여유 10% 규약."
    scale = guide.type_scale
    assert 44 <= scale.cover <= 54
    assert 36 <= scale.section <= 40
    assert 28 <= scale.slide_title <= 32
    assert 20 <= scale.card_title <= 22
    assert 14 <= scale.body <= 16
    assert 11 <= scale.caption <= 12


def test_가이드가_없으면_멈춘다(tmp_path: Path) -> None:
    """**조용히 기본값으로 넘어가지 않는다.** 고쳐도 반영되지 않는 사고가 난다."""
    load_design_guide.cache_clear()
    try:
        with pytest.raises(DesignGuideError, match="디자인 가이드가 없습니다"):
            load_design_guide(str(tmp_path / "없는파일.json"))
    finally:
        load_design_guide.cache_clear()


#: 색을 코드에 박았는지 보는 자리. **슬라이드와 png 를 그리는 두 모듈이다.**
_COLOR_SCANNED = ("report/slides.py", "report/figures.py")
_HEX_COLOR = re.compile(r"[\"']#[0-9a-fA-F]{3,8}[\"']")


def test_색을_코드에_박지_않는다() -> None:
    """가이드가 바뀌면 **고칠 자리가 한 곳**이어야 한다 (36세션 3절)."""
    offenders = {
        name: _HEX_COLOR.findall((SRC_ROOT / name).read_text(encoding="utf-8"))
        for name in _COLOR_SCANNED
    }
    found = {name: hits for name, hits in offenders.items() if hits}
    assert not found, (
        f"색을 코드에 박았습니다: {found}. data\\ppt_design.json 에 넣고 "
        "kwise.report.design.load_design_guide() 로 읽으십시오."
    )


def test_png_팔레트가_가이드에서_온다() -> None:
    """**화면은 그대로 두고 png 만 맞춘다** (36세션 4절)."""
    guide = load_design_guide()
    assert figures.chart_palette() is guide.chart
    assert guide.chart.canvas == "#FFFFFF", "png 는 흰 캔버스가 확정이다."
    assert guide.chart.series[0] == guide.colors.blue
    assert guide.chart.increase == guide.colors.coral
    # 계시 색은 **단가가 높을수록 짙다** — 화면과 같은 규칙이다.
    assert set(guide.chart.band) == {"경부하", "중간부하", "최대부하"}


def test_png_은_배경을_투명으로_굽지_않는다(sample_diagnosis: Diagnosis) -> None:
    """**png 글자색 규약은 「흰 바탕이 확정」 위에 선다** (36세션 4절).

    화면(altair)은 다크 모드를 타서 글자에 중립색을 박지 않지만, png 는 배경이
    하나뿐이라 그 규약이 성립하지 않는다. 대신 그 배경을 여기서 확정한다.
    """
    source = (SRC_ROOT / "report" / "figures.py").read_text(encoding="utf-8")
    assert "transparent=True" not in source, "png 를 투명 배경으로 구우면 글자가 사라진다."
    assert "facecolor=canvas" in source
    png = figures.hourly_profile_png(sample_diagnosis.peak)
    assert png[:8] == b"\x89PNG\r\n\x1a\n"


# ===================================================================== ④⑤ 형태


def test_슬라이드마다_시각_요소가_하나_이상이다(full_sections: DocumentSections) -> None:
    """**표형 슬라이드도 표 자체가 시각 요소다** (36세션 5절)."""
    deck = build_slides(full_sections)
    specs = slide_specs(full_sections)
    for slide, spec in zip(deck.slides, specs, strict=True):
        assert _visual_shapes(slide), f"「{spec.title}」 장에 시각 요소가 없습니다."


def test_레이아웃이_반복되지_않는다(full_sections: DocumentSections) -> None:
    """**같은 형태를 되풀이하지 않는다** (36세션 3-4). 가이드가 금지한다."""
    specs = slide_specs(full_sections)
    used = [spec.layout for spec in specs]
    assert set(used) <= set(LAYOUTS), f"모르는 레이아웃이 있습니다: {set(used) - set(LAYOUTS)}"

    # 수단별 장은 가이드가 「차트 + 지표」 로 못박은 자리라 형태가 같은 것이 옳다.
    # **나머지 장**을 본다 — 여기가 되풀이되면 덱이 한 형태로 보인다.
    # **부록은 세지 않는다** (39세션 5절). 수단마다 한 장씩 붙는 자리라 형태가
    # 같은 것이 옳다 — 세면 「표형이 잦다」 는 결론만 나온다.
    frame = [spec.layout for spec in specs if spec.measure is None and spec.page is None]
    counts = {kind: frame.count(kind) for kind in set(frame)}
    assert len(counts) >= 5, f"형태가 {len(counts)}가지뿐입니다: {counts}"
    assert max(counts.values()) <= 3, f"한 형태가 너무 잦습니다: {counts}"


def test_카드로만_만들지_않는다(full_sections: DocumentSections) -> None:
    """**표·리스트·여백만인 슬라이드를 섞는다** (36세션 3-4).

    가이드가 금지하는 것은 카드가 아니라 **카드만인 덱**이다. 그림 없이 표와
    선만으로 서는 장이 있어야 한다.
    """
    deck = build_slides(full_sections)
    picture_free = [
        slide
        for slide in deck.slides
        if not any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)
    ]
    assert len(picture_free) >= 3, "그림 없이 서는 슬라이드가 셋에 못 미칩니다."


def test_슬라이드당_큰_타이틀은_하나다(full_sections: DocumentSections) -> None:
    """**크기 위계** (36세션 3-3). 제목 크기의 글자가 한 장에 둘 있으면 안 된다."""
    guide = load_design_guide()
    threshold = guide.type_scale.card_title
    deck = build_slides(full_sections)
    for slide in deck.slides:
        big = [
            run
            for shape in slide.shapes
            if shape.has_text_frame
            for paragraph in shape.text_frame.paragraphs
            for run in paragraph.runs
            if run.font.size is not None and run.font.size.pt > threshold
        ]
        assert len(big) <= 1, f"큰 타이틀이 {len(big)}개입니다: {[run.text for run in big]}"


def test_본문은_좌측_정렬이고_그림_캡션만_가운데다(full_sections: DocumentSections) -> None:
    """**본문은 좌측 정렬이다** (36세션 3-3). **그림 캡션만 가운데다** (53세션 1-2).

    그림은 칸 가운데 앉는데(``_picture_block``) 캡션만 칸 왼쪽 귀퉁이에 붙어
    있어 그림 아래가 아니라 그림 옆에 매달린 것처럼 보였다.
    """
    from pptx.enum.text import PP_ALIGN

    guide = load_design_guide()
    deck = build_slides(full_sections)
    centered = 0
    for slide in deck.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.alignment is PP_ALIGN.CENTER:
                    centered += 1
                    sizes = {run.font.size.pt for run in paragraph.runs if run.font.size}
                    assert sizes == {guide.type_scale.caption}, (
                        f"캡션이 아닌 글이 가운데 정렬입니다: {paragraph.text}"
                    )
                    continue
                assert paragraph.alignment in (None, PP_ALIGN.LEFT)
    assert centered, "가운데 정렬된 캡션이 하나도 없습니다."


def test_슬라이드_밖으로_나가는_것이_없다(full_sections: DocumentSections) -> None:
    """**텍스트 상자에 여유를 두었어도 자리를 벗어나면 안 된다** (36세션 3-1)."""
    guide = load_design_guide()
    width = Inches(guide.slide.width_in)
    height = Inches(guide.slide.height_in)
    deck = build_slides(full_sections)
    for index, slide in enumerate(deck.slides, start=1):
        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            right = shape.left + (shape.width or 0)
            bottom = shape.top + (shape.height or 0)
            assert shape.left >= 0 and shape.top >= 0, f"{index}장이 왼쪽·위로 넘칩니다."
            assert right <= width + Inches(0.01), f"{index}장이 오른쪽으로 넘칩니다."
            assert bottom <= height + Inches(0.01), f"{index}장이 아래로 넘칩니다."


def test_표지와_마무리가_다크다(full_sections: DocumentSections) -> None:
    """**샌드위치 구조** (36세션 3-2 · 37세션).

    36세션은 표지만 다크로 두어 구조가 반쪽이었다 — 지시서 목차에 마무리 장이
    없었기 때문이다. 아랫빵을 붙여 앞뒤가 짝을 이룬다.
    """
    guide = load_design_guide()
    deck = build_slides(full_sections)
    fills = [str(slide.background.fill.fore_color.rgb) for slide in deck.slides]
    assert fills[0] == guide.colors.cover.lstrip("#").upper()
    assert fills[-1] == guide.colors.closing.lstrip("#").upper()
    assert set(fills[1:-1]) == {guide.colors.white.lstrip("#").upper()}, (
        "본문 콘텐츠는 라이트다 — 다크는 바깥 둘뿐이다."
    )


def test_전체_배경_밴드에_포인트색을_쓸_수_없다() -> None:
    """**딥그린·네이비만 전체 배경이다** (36세션 3-2). 코랄은 작은 포인트뿐이다."""
    from dataclasses import replace

    from kwise.report.design import DesignGuideError

    colors = load_design_guide().colors
    for value in (colors.cover, colors.closing):
        assert value in (colors.deep_green, colors.dark_primary)
    with pytest.raises(DesignGuideError, match="전체 배경 밴드"):
        _ = replace(colors, cover_background="coral").cover
    with pytest.raises(DesignGuideError, match="전체 배경 밴드"):
        _ = replace(colors, closing_background="coral").closing


# ===================================================================== 37세션 · 마무리


def test_마무리가_맨_뒤에_있다(full_sections: DocumentSections) -> None:
    """**표지·목차와 같은 취급이다** — 자료가 무엇이든 늘 나온다 (37세션)."""
    for sections in (
        full_sections,
        DocumentSections(usage=full_sections.usage, bill=full_sections.bill),
    ):
        specs = slide_specs(sections)
        assert specs[-1].key == "closing"
        assert specs[-1].layout == "closing"
        assert specs[-1].title == SLIDE_TITLES["closing"]


def test_마무리가_다음_단계를_적는다(full_sections: DocumentSections) -> None:
    """**결론이 아니라 성격을 밝힌다.** 채택한 안은 현장에서 확정된다."""
    deck = build_slides(full_sections)
    text = _slide_text(list(deck.slides)[-1])
    assert SLIDE_TITLES["closing"] in text
    assert NEXT_STEPS_HEADLINE in text
    for title, detail in NEXT_STEPS:
        assert title in text and detail in text
    assert _visual_shapes(list(deck.slides)[-1]), "마무리에 시각 요소가 없습니다."


def test_다음_단계_문구가_앞뒤로_겹치지_않는다(loaded_sections: DocumentSections) -> None:
    """**같은 말을 두 번 하지 않는다** (37세션).

    덱 어디에도 없던 문구라 옮겨 올 것이 없었다 — 새로 적는 자리도 마무리
    하나뿐이어야 한다. 앞 장이 같은 말을 하기 시작하면 여기서 걸린다.
    """
    deck = build_slides(loaded_sections)
    slides = list(deck.slides)
    assert _deck_text(deck).count(NEXT_STEPS_HEADLINE) == 1
    front = "\n".join(_slide_text(slide) for slide in slides[:-1])
    for _title, detail in NEXT_STEPS:
        assert detail not in front, f"「{detail}」 가 앞 장에도 있습니다."
    assert "초기 판단용" not in front


def test_장수가_수단_개수를_따라간다(
    full_sections: DocumentSections, diagnosis_only: DocumentSections
) -> None:
    """뼈대 열한 장 + 켠 수단 수 (37세션). 수단 0개면 11장, 여섯이면 17장이다.

    **41세션에 수단이 여섯이 됐다** — 7.7 잉여 활용을 태양광 안으로 넣었다.
    """
    from dataclasses import replace

    # 뼈대 = 표지·목차·진단 다섯·요약·조합·마무리 + **부록 장 수** (39세션 5절).
    frame = FRAME_BEFORE_APPENDIX + 1  # 마무리
    assert len(slide_specs(diagnosis_only)) == frame + len(appendix_pages(diagnosis_only))
    assert len(slide_specs(full_sections)) == (
        frame + len(full_sections.measures) + len(appendix_pages(full_sections))
    )

    every = replace(full_sections, measures=_all_measures(full_sections))
    assert len(every.measures) == 6
    assert len(slide_specs(every)) == frame + 6 + len(appendix_pages(every))


# ===================================================================== 39세션 · 해석과 문구


#: 해석 한 줄과 용어 각주가 있어야 할 장 (39세션 1절).
#: 표지·목차·마무리는 뼈대라 빼고, 수단 장은 결론 한 줄이 그 몫을 진다.
LEAD_SLIDES: tuple[str, ...] = (
    "building",
    "usage_pattern",
    "peak_summary",
    "peak_detail",
    "structure",
    "measure_summary",
    "combination",
    "appendix",
)


def test_슬라이드마다_해석_한_줄이_있다(full_sections: DocumentSections) -> None:
    """**그래프만 나열하면 무엇을 보아야 하는지 알 수 없다** (39세션 1-1).

    제목 아래 문장이 그림보다 먼저 읽혀야 한다. 진단이 이미 내린 판정을 옮기거나,
    값의 크기로 고르거나, 고정 문장이다.
    """
    deck = build_slides(full_sections)
    specs = slide_specs(full_sections)
    for slide, spec in zip(deck.slides, specs, strict=True):
        if spec.key not in LEAD_SLIDES:
            continue
        body = _slide_text(slide).replace(spec.title, "", 1)
        assert any(line.strip().endswith("니다.") for line in body.splitlines()), (
            f"「{spec.title}」 장에 해석 한 줄이 없습니다."
        )


def test_수단_장은_결론_한_줄이_해석을_진다(full_sections: DocumentSections) -> None:
    """수단 장에는 별도 해석을 얹지 않는다 — **같은 말을 두 번 하지 않는다.**"""
    deck = build_slides(full_sections)
    specs = slide_specs(full_sections)
    for slide, spec in zip(deck.slides, specs, strict=True):
        if spec.measure is None:
            continue
        entry = full_sections.measures[spec.measure]
        assert entry.conclusion in _slide_text(slide)


def test_판정을_문장으로_옮기고_근거_숫자를_붙인다(full_sections: DocumentSections) -> None:
    """**진단이 이미 판정한 것은 그 판정을 쓴다** (39세션 1-1).

    화면 「개선 여지 요약」 의 태양광 등급이고, 근거 숫자가 바로 아래 지표의
    정오 비중이다 — 판정과 근거가 함께 읽혀야 한다.
    """
    from kwise.diagnose.summary import PvPotential
    from kwise.report import narrative

    diagnosis = full_sections.diagnosis
    assert diagnosis is not None
    lead = narrative.peak_summary_lead(diagnosis)
    share = f"{diagnosis.summary.pv_midday_share * 100:,.0f}%"
    assert share in lead, "판정만 적고 근거 숫자를 빠뜨렸습니다."
    assert set(PvPotential) == {PvPotential.HIGH, PvPotential.MEDIUM, PvPotential.LOW}
    # 등급마다 다른 문장을 낸다 — 하나로 뭉뚱그리면 판정을 옮긴 것이 아니다.
    assert len({text for text in narrative._PV_LEAD.values()}) == len(PvPotential)


def test_용어_풀이가_화면_툴팁에서_온다() -> None:
    """**두 벌로 적지 않는다** (39세션 1-2).

    화면은 「산식 한 줄 + 의미 한 줄」 로 툴팁을 적고, 슬라이드는 같은 산식을
    각주로 깐다 — 낳는 자리가 하나여야 한쪽만 고쳐지지 않는다.
    """
    from kwise.report import narrative

    table = narrative.terms()
    for key in ("load_factor", "base_load_ratio", "weekend_ratio", "off_hours_energy_share"):
        term = table[key]
        assert term.formula and term.meaning
        assert term.tooltip.startswith(term.formula)
        assert term.meaning in term.tooltip
        assert term.name in term.line

    source = (SRC_ROOT / "ui" / "views" / "diagnose.py").read_text(encoding="utf-8")
    assert "narrative.terms(" in source, "화면이 용어를 따로 들고 있으면 두 벌이 된다."


def test_용어_각주가_그_장의_용어만_깐다(full_sections: DocumentSections) -> None:
    """**전부 깔면 각주가 본문이 된다.** 그 장에 나오는 것만 고른다."""
    from kwise.report import narrative

    deck = build_slides(full_sections)
    table = narrative.terms()
    present = {spec.key for spec in slide_specs(full_sections)}
    for key, chosen in narrative.GLOSSARY_KEYS.items():
        # 수단 장은 켠 것만 선다 (53세션 8-1 에 수단 장에도 용어를 깔았다).
        if not chosen or key not in present:
            continue
        text = _slide_text(_slide_by_key(deck, full_sections, key))
        for term in chosen:
            assert f"{table[term].name} = " in text, f"「{key}」 장에 「{term}」 풀이가 없습니다."
        loose = [name for name in set(table) - set(chosen) if f"{table[name].name} = " in text]
        assert not loose, f"「{key}」 장에 그 장의 용어가 아닌 풀이가 있습니다: {loose}"


def test_상위_구간을_왜_보는지_적고_판정까지_한다(full_sections: DocumentSections) -> None:
    """화면 툴팁 문장 뒤에 **판정이 온다** (39세션 1-3 · 53세션 4-4).

    39세션까지는 「무엇을 보는 그림인가」 만 적고 이 건물이 어느 쪽인지는
    적지 않았다 — 그림 둘을 보고 읽는 사람이 스스로 세어야 했다.
    """
    from kwise.report import narrative

    deck = build_slides(full_sections)
    text = _slide_text(_slide_by_key(deck, full_sections, "peak_detail"))
    assert "낮에 몰리면 태양광이" in text
    assert full_sections.diagnosis is not None
    assert narrative.peak_summary_lead(full_sections.diagnosis) in text
    # **5장은 그 판정을 더 이상 적지 않는다** — 그림이 월별 최대수요다.
    summary = _slide_text(_slide_by_key(deck, full_sections, "peak_summary"))
    assert "낮 시간에 발생해" not in summary, summary
    assert "최대수요가" in summary


def test_6장_그림이_문장_차례와_같다(full_sections: DocumentSections) -> None:
    """**문장이 말하는 그림이 왼쪽이다** (53세션 4-5)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu

    slide = _slide_by_key(build_slides(full_sections), full_sections, "peak_detail")
    captions = sorted(

            (Emu(shape.left).inches, shape.text_frame.text)
            for shape in slide.shapes  # type: ignore[attr-defined]
            if (shape.has_text_frame and "구간이 발생한 시각" in shape.text_frame.text)
            or (shape.has_text_frame and "평균 부하 모양" in shape.text_frame.text)

    )
    assert next(text for _left, text in captions).startswith("최대수요 상위")
    pictures = sorted(
        Emu(shape.left).inches
        for shape in slide.shapes  # type: ignore[attr-defined]
        if shape.shape_type is MSO_SHAPE_TYPE.PICTURE
    )
    assert len(pictures) == 2


# ===================================================================== 39세션 2절 · 어투


#: 개발자 어투 (39세션 2-1·2-3·2-4). **고객에게 우리 사정을 말하지 않는다.**
DEVELOPER_VOICE: tuple[str, ...] = (
    "모릅니다",
    "넣지 않았습니다",
    "자리가 모자라",
    "그렸습니다",
    "채우지 않습니다",
    "다시 잰 값",
)


def test_개발자_어투가_없다(loaded_sections: DocumentSections) -> None:
    """**「모릅니다」 는 우리 사정이다** (39세션 2-1).

    「미산출 — 투자비를 모릅니다」 와 「미산출 — 계통 연계·설비 조건에 따릅니다」 가
    나란히 있었다. 값 자리에는 「미산출」 만 두고, 무엇을 넣으면 값이 나오는지는
    각주가 받는다.
    """
    text = _deck_text(build_slides(loaded_sections))
    for phrase in DEVELOPER_VOICE:
        assert phrase not in text, f"개발자 어투가 남아 있습니다: 「{phrase}」"


def test_미산출은_값과_사유를_가른다() -> None:
    """지표 칸에는 값만, 사유는 각주로 (39세션 2-1)."""
    assert split_reason("미산출 — 투자비 미입력") == ("미산출", "투자비 미입력")
    assert split_reason("미산출") == ("미산출", "")
    assert split_reason("53,575,000원") == ("53,575,000원", "")


def test_12개월_환산을_값마다_되풀이하지_않는다(full_sections: DocumentSections) -> None:
    """**같은 값을 두 번 적지 않는다** (39세션 2-2).

    「9,050,000원 (12개월 환산 9,050,000원)」 이 두 줄로 흐르고, 읽는 사람은 어느
    쪽이 답인지 되묻는다. 기준은 각주가 **한 번** 적는다.
    """
    text = _deck_text(build_slides(full_sections))
    assert "12개월 환산" not in text.replace(ANNUAL_BASIS_NOTE, "")
    assert text.count(ANNUAL_BASIS_NOTE) == 1
    for entry in full_sections.measures:
        assert "12개월 환산" not in entry.slide_saving


def test_Word_에도_마크다운_표식이_없다(loaded_sections: DocumentSections) -> None:
    """**Word 도 마크다운을 해석하지 않는다** (39세션 2-6).

    38세션에 PPT 만 막았고 Word 에는 별표가 남아 있었다. 벗기는 자리는
    :func:`kwise.report.notices.plain_text` 하나다.
    """
    document = build_document(loaded_sections)
    parts = [item.text for item in document.paragraphs]
    for table in document.tables:
        parts.extend(cell.text for row in table.rows for cell in row.cells)
    text = "\n".join(parts)
    for mark in MARKDOWN_MARKS:
        assert mark not in text, f"Word 에 마크다운 표식 「{mark}」 이 남아 있습니다."


def test_Word_가_글자를_쓰는_자리를_모두_지난다() -> None:
    """**자리마다 벗기면 새 문구를 붙일 때 빠뜨린다** (39세션 2-6)."""
    source = (SRC_ROOT / "report" / "document.py").read_text(encoding="utf-8")
    calls = ("document.add_paragraph(", "document.add_heading(", "cell.text =")
    for line in source.splitlines():
        stripped = line.strip()
        if stripped.startswith("#") or "plain_text(" in stripped:
            continue
        if "add_paragraph()" in stripped:
            continue
        for call in calls:
            assert call not in stripped, f"마크다운을 벗기지 않는 자리가 있습니다: {stripped}"


def test_단위가_빠진_칸이_없다() -> None:
    """「역률 개선 (97)」 이 그렇게 나가고 있었다 (39세션 2-5)."""
    from kwise.measures import AppliedMeasure

    assert AppliedMeasure("power_factor", (("power_factor_pct", 97.0),)).label == "역률 개선 (97%)"
    assert AppliedMeasure("solar", (("capacity_kwp", 240.0),)).label == "태양광 (240 kWp)"


def test_조합_구성을_한_줄로_이어_적는다(full_sections: DocumentSections) -> None:
    """**고른 수단을 순서대로** (39세션 3-3).

    조합 이름은 「+ ESS 목표 5,170 kW」 처럼 직전 조합에 무엇을 더했는가를 적는
    것이라, 그것만 떼어 놓으면 앞의 수단들이 보이지 않았다.
    """
    comparison = full_sections.comparison
    assert comparison is not None
    baseline = comparison.combinations[0].spec.selection
    line = comparison.best.spec.composition(baseline)
    assert " + " in line
    assert "선택요금 전환" in line
    assert "(" not in line, "괄호를 겹쳐 적지 않는다."
    text = _slide_text(_slide_by_key(build_slides(full_sections), full_sections, "combination"))
    assert line in text


# ===================================================================== 39세션 4절 · 값이 0인 수단


def test_잉여가_0인_까닭을_숫자로_적는다() -> None:
    """**0이라는 것 자체가 정보다. 다만 왜 0인지가 있어야 한다** (39세션 4-3)."""
    from kwise.report import narrative

    line = narrative.surplus_lead(
        total_kwh=0.0,
        generation_kwh=281_293.0,
        self_consumed_kwh=281_293.0,
        surplus_free_kwp=2_048.0,
    )
    assert "자가소비" in line and "2,048 kWp" in line
    assert "281" in line, "발전량을 숫자로 보인다."


def test_DR_이_0인_까닭을_적는다() -> None:
    """**거래 가능일과 저부하 평일만 적으면 왜 그런지 알 수 없다** (39세션 4-1)."""
    from kwise.report import narrative

    assert narrative.dr_lead(None) == ""


def test_결론과_무관한_주의사항을_싣지_않는다(full_sections: DocumentSections) -> None:
    """**하지도 않을 일을 조심하라는 글이 결론보다 길어진다** (39세션 4-2)."""
    from dataclasses import replace

    base = full_sections.measures[0]
    idle = replace(
        base,
        actionable=False,
        cautions=("계약전력 하향은 되돌리기 어렵습니다.",),
        facts=(("하향 여지", "0 kW"),),
        figure=None,
        figures=(),
    )
    assert _cautions(idle) == (), "여지가 없으면 실행 주의사항을 싣지 않는다."
    assert _cautions(replace(idle, actionable=True)), "여지가 있으면 그대로 싣는다."


# ===================================================================== 38세션 3절 · 그래프 대조


#: **화면 그림 → PPT 그림.** 38세션 3-5 에 하나씩 대조한 결과다.
#:
#: 화면에 있는데 PPT 에 없던 것이 셋이었다 — 역률 일일 곡선 · 태양광 연간
#: 발전량 · ESS 회수기간 곡선. 셋 다 「무엇을 보고 그렇게 판단했나」 에 답하는
#: 그림이라, 빠지면 슬라이드에 결론만 남는다.
#:
#: **요금제 전환 둘은 png 한 장이 함께 진다** — ``tariff_option_png`` 가 위 칸에
#: 그룹 막대, 아래 칸에 현행 대비 차액을 그린다 (17세션 1-2·1-3).
SCREEN_TO_DECK: dict[str, str] = {
    "daily_temperature_chart": "daily_temperature_png",
    "monthly_peak_chart": "monthly_peak_png",
    "hourly_profile_chart": "hourly_profile_png",
    "top_hour_chart": "top_hour_png",
    "monthly_charge_chart": "monthly_charge_png",
    "band_donut_chart": "band_donut_grid_png",
    "tariff_delta_chart": "tariff_option_png",
    "tariff_option_chart": "tariff_option_png",
    "dr_daily_chart": "dr_daily_png",
    "power_triangle_chart": "power_triangle_png",
    "power_factor_day_chart": "power_factor_day_png",
    "solar_annual_chart": "solar_annual_png",
    "solar_day_chart": "solar_day_png",
    "ess_day_chart": "ess_day_png",
    "surplus_daily_chart": "surplus_daily_png",
}


def _screen_charts() -> set[str]:
    """화면이 **실제로 그리는** 그림 이름. 정의만 있고 안 그리는 것은 세지 않는다."""
    pattern = re.compile(r"charts\.(\w+_chart)\(")
    names: set[str] = set()
    for path in (SRC_ROOT / "ui" / "views").glob("*.py"):
        names |= set(pattern.findall(path.read_text(encoding="utf-8")))
    return names


def test_화면에_있는_그림이_PPT_에도_있다() -> None:
    """**하나씩 세어 대조한다** (38세션 3-5).

    화면에 그림을 더하고 슬라이드를 잊으면, 받는 사람은 「분석은 했는데 근거는
    없는」 보고서를 받는다. 그 어긋남은 둘을 나란히 놓고서야 드러나므로 시험이
    센다 — 짝을 못 지으면 :data:`SCREEN_TO_DECK` 에 한 줄을 더하거나 png 를
    새로 만들어야 한다.
    """
    drawn = _screen_charts()
    assert drawn, "화면 그림을 찾지 못했습니다 — 훑는 규칙이 낡았습니다."
    missing = sorted(drawn - set(SCREEN_TO_DECK))
    assert not missing, (
        f"화면에만 있고 PPT 에 없는 그림입니다: {missing}. "
        "png 를 만들어 슬라이드에 싣거나, 싣지 않는 이유를 SCREEN_TO_DECK 옆에 적으십시오."
    )
    stale = sorted(set(SCREEN_TO_DECK) - drawn)
    assert not stale, f"화면이 더는 그리지 않는 그림이 표에 남아 있습니다: {stale}"
    for screen, deck in SCREEN_TO_DECK.items():
        assert deck in figures.__all__, f"「{screen}」 의 짝 「{deck}」 이 figures 에 없습니다."


def test_역률_태양광_ESS_가_그림_둘을_싣는다(full_sections: DocumentSections) -> None:
    """**화면은 셋 다 그림이 둘이다** (38세션 3-1·3-2·3-3).

    36세션의 덱은 하나씩만 실어, 전력삼각형만으로 「어느 시간대가 요금 대상인가」
    를, 대표일 곡선만으로 「한 해에 얼마나 만드나」 를, 충·방전 곡선만으로
    「목표 5,170 kW 는 어디서 나왔나」 를 답하지 못했다.
    """
    from kwise.report.document import MeasureFigure

    entry = MeasureEntry(
        kind=MEASURE_CATALOG[0],
        conclusion="본문",
        saving="0원",
        investment="0원",
        payback="즉시",
        certainty="높음",
        figures=(MeasureFigure(b"a", "왼쪽"), MeasureFigure(b"b", "오른쪽")),
    )
    assert len(entry.slide_figures) == 2
    spec = slide_specs(_with_measures(full_sections, (entry,)))[8]
    assert spec.layout == "chart_pair", "그림이 둘이면 좌우로 나눈다."

    # 하나만 주면 폭 전체를 쓴다 — 빈 칸을 남기지 않는다.
    single = MeasureEntry(
        kind=MEASURE_CATALOG[0],
        conclusion="본문",
        saving="0원",
        investment="0원",
        payback="즉시",
        certainty="높음",
        figure=b"a",
        figure_caption="하나",
    )
    assert len(single.slide_figures) == 1
    assert slide_specs(_with_measures(full_sections, (single,)))[8].layout == "stat_chart"


def _with_measures(
    sections: DocumentSections, measures: tuple[MeasureEntry, ...]
) -> DocumentSections:
    from dataclasses import replace

    return replace(sections, measures=measures)


# ===================================================================== ⑥ 부록


#: **PPT 부록에서 뺀 셋** (36세션 6절). Word 에는 그대로 있다.
DROPPED_FROM_DECK: tuple[str, ...] = (
    "ESS 조달 사례",
    APPENDIX_TITLES["B"],
    APPENDIX_TITLES["C"],
)


@pytest.fixture(scope="module")
def loaded_sections(full_sections: DocumentSections) -> DocumentSections:
    """**부록 재료를 다 넣은 덱.** 뺐다는 것을 보이려면 넣어 놓고 봐야 한다."""
    from dataclasses import replace

    from kwise.measures import load_ess_cost_model
    from kwise.tariff import load_tariff

    return replace(
        full_sections,
        ess_cases=load_ess_cost_model().case_table(),
        tariff_table=load_tariff(),
    )


def test_부록에서_뺀_셋이_PPT_에_없다(loaded_sections: DocumentSections) -> None:
    """슬라이드에서 읽히지 않는 분량이다. **지운 것이 아니라 매체를 가린 것이다.**"""
    sections = loaded_sections
    text = _deck_text(build_slides(sections))
    for title in DROPPED_FROM_DECK:
        assert title not in text, f"「{title}」 이 PPT 에 남아 있습니다."


def test_뺀_셋이_Word_에는_남아_있다(loaded_sections: DocumentSections) -> None:
    """**PPT 에서만 뺐다.** Word 는 그대로다 (36세션 6절)."""
    document = build_document(loaded_sections)
    parts = [item.text for item in document.paragraphs]
    for table in document.tables:
        parts.extend(cell.text for row in table.rows for cell in row.cells)
    text = "\n".join(parts)
    for title in DROPPED_FROM_DECK:
        assert title in text, f"「{title}」 이 Word 에서 사라졌습니다."


def test_부록이_넘치면_장을_나눈다(full_sections: DocumentSections) -> None:
    """**자르지 않고 나눈다** (39세션 5절).

    36세션까지는 한 장에 눌러 담고 「자리가 모자라 뺐다」 고 적었는데, 그러면
    선택요금 전환 하나만 실리고 나머지 여섯이 통째로 빠졌다.
    """
    from dataclasses import replace

    from kwise.report.worksheet import Worksheet

    sheet = full_sections.worksheets[0]
    many = Worksheet(key=sheet.key, title=sheet.title, rows=sheet.rows * 6)
    sections = replace(full_sections, worksheets=(many,))
    pages = appendix_pages(sections)
    assert len(many.frame()) > APPENDIX_ROW_LIMIT
    assert len(pages) > 1, "넘치면 장을 나눈다."
    # **빈 줄은 묶음을 가르는 자리다** (53세션 8-2) — 내용이 아니라 경계라
    # 장에는 싣지 않는다. 그 밖의 줄은 한 줄도 버리지 않는다.
    records = [
        tuple(str(value) for value in row) for row in many.frame().to_numpy()
    ]
    kept = [record for record in records if any(value.strip() for value in record)]
    assert sum(len(page.rows) for page in pages) == len(kept), "한 줄도 버리지 않는다."
    assert all("(" in page.title for page in pages), "이어지는 장임을 제목이 밝힌다."
    text = _deck_text(build_slides(sections))
    assert "자리가 모자라" not in text, "제작 사정을 고객에게 말하지 않는다."


def test_부록이_수단마다_한_장_이상이다(full_sections: DocumentSections) -> None:
    """**분석한 자료를 감추지 않는다** (39세션 5절).

    절감액이 산출된 수단은 모두 근거가 실리고, 0이거나 미산출인 수단은 빠지되
    **뺐다는 사실을 각주가 적는다** — 조용히 빼면 「검토하지 않았다」 로 읽힌다.
    """
    sections = full_sections
    priced = [entry for entry in sections.measures if entry.has_saving]
    dropped = [entry for entry in sections.measures if not entry.has_saving]
    assert priced and dropped, "이 시험은 값이 있는 수단과 없는 수단을 함께 본다."
    titles = [page.title for page in appendix_pages(sections)]
    sheets = {sheet.key for sheet in sections.worksheets}
    for entry in priced:
        if entry.kind.key not in sheets:
            continue
        assert any(measure_slide_title(entry) in title for title in titles), (
            f"「{entry.kind.label}」 의 근거가 부록에 없습니다."
        )
    for entry in dropped:
        assert not any(measure_slide_title(entry) in title for title in titles)
    text = _deck_text(build_slides(sections))
    for entry in dropped:
        assert entry.kind.label in text, "뺐다는 사실을 각주가 적는다."
    assert "전문은 Excel 부록 A 에 있습니다" in text


# ===================================================================== ⑦ 글꼴·한글


def test_OS_전용_글꼴_이름이_슬라이드_코드에_없다() -> None:
    """**설정으로 뺐다** (36세션 3-1). 배포 시험이 ``src\\`` 를 훑는다."""
    from tests.test_deployment import FONT_NAME_ALLOWED

    assert "src/kwise/report/slides.py" not in FONT_NAME_ALLOWED
    assert "src/kwise/report/design.py" not in FONT_NAME_ALLOWED
    source = (SRC_ROOT / "report" / "slides.py").read_text(encoding="utf-8")
    assert "guide.typography.primary" in source, "글꼴 이름은 가이드에서 와야 합니다."


def test_글꼴_후보가_없으면_물러선다() -> None:
    """**없을 때 물러설 곳을 정해 둔다** (36세션 3-1)."""
    from dataclasses import replace

    typography = load_design_guide().typography
    assert typography.primary == typography.candidates[0]
    assert replace(typography, candidates=()).primary == typography.fallback


def test_한글이_깨지지_않는다(full_sections: DocumentSections, tmp_path: Path) -> None:
    """**동아시아 글꼴을 따로 지정해야** PowerPoint 가 한글에 쓴다."""
    path = export_slides(full_sections, output_dir=tmp_path)
    deck = Presentation(str(path))
    text = _deck_text(deck)
    assert "전력 비용 진단 보고서" in text
    assert "건물현황 및 계약정보" in text
    assert "�" not in text and "?" * 3 not in text

    # **동아시아 글꼴(``a:ea``)은 ``a:rPr`` 안에 있다.** ``a:latin`` 만 있으면
    # PowerPoint 가 한글에 그 글꼴을 쓰지 않는다 (Word 의 ``w:eastAsia`` 와 같다).
    name = load_design_guide().typography.primary
    eastasian = [
        run._r.find(f"{_A}rPr/{_A}ea")
        for slide in deck.slides
        for shape in slide.shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    ]
    assert eastasian and all(item is not None for item in eastasian)
    assert all(item.get("typeface") == name for item in eastasian if item is not None)


# ===================================================================== ⑧ 파일명


def test_파일명에_날짜_시각이_붙는다() -> None:
    """PowerPoint 가 파일을 열고 있으면 덮어쓰기가 실패한다 (Excel·Word 와 같다)."""
    path = slides_path(Path("out"), now=dt.datetime(2026, 8, 11, 9, 5))
    assert path.name == "kwise_report_20260811_0905.pptx"


def test_바이트를_돌려줄_때도_같은_이름이다(full_sections: DocumentSections) -> None:
    _payload, name = slides_bytes(full_sections, now=dt.datetime(2026, 8, 11, 9, 5))
    assert name == "kwise_report_20260811_0905.pptx"


def test_저장_경로에_폴더가_없으면_만든다(full_sections: DocumentSections, tmp_path: Path) -> None:
    target = tmp_path / "없던" / "폴더"
    path = export_slides(full_sections, output_dir=target)
    assert path.is_file()


# ===================================================================== 도넛 넷


def test_계시_도넛이_2x2_로_한_장에_들어간다(sample_diagnosis: Diagnosis) -> None:
    """**넷의 조각 순서·색·시작각을 맞춘다** (36세션 4절).

    도넛에는 y 축이 없으므로 축 대신 이 셋이 「같은 잣대로 그렸다」 를 진다.
    """
    structure = sample_diagnosis.structure
    assert structure is not None
    pairs = season_pairs(structure)
    assert [label for label, _key in pairs] == ["전체", "봄·가을", "여름", "겨울"]
    assert figures.DONUT_GRID == (2, 2)
    png = figures.band_donut_grid_png(structure, pairs)
    assert png[:8] == b"\x89PNG\r\n\x1a\n"

    source = (SRC_ROOT / "report" / "figures.py").read_text(encoding="utf-8")
    assert "startangle=90" in source and "counterclock=False" in source


def test_없는_계절은_도넛을_그리지_않는다(sample_diagnosis: Diagnosis) -> None:
    """**빈 원은 「그 계절에 안 썼다」 로 읽힌다** (화면과 같은 규칙)."""
    from dataclasses import replace

    structure = sample_diagnosis.structure
    assert structure is not None
    half = structure.band_season_kwh.drop(index="summer")
    pairs = season_pairs(replace(structure, band_season_kwh=half))
    assert [label for label, _key in pairs] == ["전체", "봄·가을", "겨울"]


# ===================================================================== png 날짜 축


def test_새_png_도_한국식_날짜_축이다(sample_usage: UsageData) -> None:
    """**화면만 고치면 어긋난다** (33세션 1절 · 13세션). matplotlib 은 기본이 영어다."""
    source = (SRC_ROOT / "report" / "figures.py").read_text(encoding="utf-8")
    body = source.split("def daily_usage_png")[1].split("\ndef ")[0]
    assert "date_axis(axes)" in body, "연간 사용량 png 가 날짜 축 규약을 타지 않습니다."
    png = figures.daily_usage_png(sample_usage)
    assert png[:8] == b"\x89PNG\r\n\x1a\n"


def test_월별_요금_png_의_달_축이_한국식이다(sample_diagnosis: Diagnosis) -> None:
    """**해는 바뀔 때만 적는다** (33세션 1절)."""
    source = (SRC_ROOT / "report" / "figures.py").read_text(encoding="utf-8")
    body = source.split("def monthly_charge_png")[1].split("\ndef ")[0]
    assert "month_labels(months)" in body
    structure = sample_diagnosis.structure
    assert structure is not None
    assert figures.monthly_charge_png(structure)[:8] == b"\x89PNG\r\n\x1a\n"


# ==================================================== 41세션 · 잉여를 태양광 안으로


def test_ppt_에_잉여_장이_없고_태양광_장이_진다(
    sample_usage: UsageData, tariff: TariffTable, sample_bill: BillingResult
) -> None:
    """**7.7 장을 없애고 7.5 에 녹였다** (41세션 2-1·2-6).

    잉여는 태양광을 얼마나 크게 지을지에 따라 나오는 결과이지 개선안이 아니다.
    장을 따로 세우면 같은 발전량 이야기를 두 번 한다 — 그래서 **장이 하나 준다.**

    일별 잉여 그림은 슬라이드로 오지 않는다. 수단 장은 그림 둘까지이고(38세션
    3절) 태양광 장은 연간 발전량과 대표일 곡선이 그 둘을 이미 쓴다. 그림은
    화면(잉여 처리 접힘)에 남겼다.
    """
    import numpy as np
    import pandas as pd

    from kwise.measures import apply_generation, solar_point
    from kwise.measures.surplus import evaluate_surplus
    from kwise.report.document import measure_entries
    from kwise.tariff import TariffSelection

    selection = TariffSelection("general_b", "high_a", "I")
    index = pd.DatetimeIndex(sample_usage.kw.index)
    hours = index.hour + index.minute / 60.0
    unit = pd.Series(
        np.clip(np.sin(np.pi * (hours - 7.0) / 12.0), 0.0, None) * 0.8,
        index=index,
        name="kw",
    )
    capacity = 4_000.0
    point = solar_point(sample_usage, tariff, selection, unit, capacity, baseline=sample_bill)
    net = apply_generation(sample_usage, unit * capacity)
    surplus = evaluate_surplus(
        sample_usage,
        tariff,
        selection,
        net.surplus_kw,
        generation_kwh=net.generated_kwh,
        net_usage=net.usage,
        capacity_kwp=capacity,
    )
    assert surplus.total_kwh > 0

    entries = measure_entries(solar=point, surplus=surplus, surplus_free_kwp=2_048.0)
    # **잉여 장이 없다.** 태양광 하나뿐이다.
    assert [entry.kind.key for entry in entries] == ["solar"]

    solar = next(entry for entry in entries if entry.kind.key == "solar")
    facts = dict(solar.facts)
    assert facts["잉여"].endswith("MWh"), facts
    assert facts["자가소비"].endswith("MWh"), facts
    assert facts["잉여 없는 최대 용량"] == "2,048 kWp", facts
    # **시나리오 나열은 태양광 각주에 없다** (53세션 3절). 잉여가 나면 「잉여
    # 활용」 장이 다음에 붙어 표로 낸다 — 각주는 기준선 0원까지 이어 적고 있었다.
    # 고른 시나리오가 없는 점(``solar_point`` 원본)이라 잉여 조각은 없다 —
    # 남는 것은 역률 조정 한 줄뿐이다 (59세션 12절).
    assert "자가소비로 줄인 요금" not in solar.slide_note, solar.slide_note
    assert solar.slide_note.startswith("역률 영향 반영 시"), solar.slide_note

    # **고르면 절감액이 무엇의 합인지 한 줄 적는다** (59세션 5절). 화면은 절감액
    # 물음표가 늘 이 줄을 낸다 (57세션) — PPT 에만 없었다.
    from kwise.measures import with_surplus_revenue

    chosen = with_surplus_revenue(
        point, revenue_won=350_000.0, scenario="외부 판매", base_fee_months=12.0
    )
    noted = next(
        entry
        for entry in measure_entries(solar=chosen, surplus=surplus, surplus_free_kwp=2_048.0)
        if entry.kind.key == "solar"
    )
    assert noted.slide_note.startswith("절감액 = 자가소비로 줄인 요금"), noted.slide_note
    assert "잉여 외부 판매" in noted.slide_note
    assert any("자격요건은 판정하지 않았습니다" in line for line in solar.cautions)

    from kwise.report.document import surplus_page

    page = surplus_page(surplus, capacity_kwp=capacity, surplus_free_kwp=2_048.0)
    assert page is not None
    names = [row[0] for row in page.scenario_rows[1:]]
    # **셋을 다 싣는다** (59세션 4·13절). 출력제어가 기본 선택이라 그 줄이 없으면
    # 지금 절감액에 든 것이 표에 없다.
    assert "외부 판매" in names
    assert "출력제어" in names, names
    # 그림은 여전히 둘이다 — 셋을 넣으면 축 눈금이 뭉개진다 (38세션 3절).
    assert len(solar.slide_figures) <= 2


def test_잉여가_0이면_태양광_장이_늘지_않는다(
    sample_usage: UsageData, tariff: TariffTable, sample_bill: BillingResult
) -> None:
    """**팔 것이 없는데 파는 절차를 조심하라고 적지 않는다** (39세션 4-2)."""
    import pandas as pd

    from kwise.measures import apply_generation, solar_point
    from kwise.measures.surplus import evaluate_surplus
    from kwise.report.document import measure_entries
    from kwise.tariff import TariffSelection

    selection = TariffSelection("general_b", "high_a", "I")
    index = pd.DatetimeIndex(sample_usage.kw.index)
    unit = pd.Series(0.0, index=index, name="kw")
    point = solar_point(sample_usage, tariff, selection, unit, 100.0, baseline=sample_bill)
    net = apply_generation(sample_usage, unit * 100.0)
    surplus = evaluate_surplus(
        sample_usage,
        tariff,
        selection,
        net.surplus_kw,
        generation_kwh=net.generated_kwh,
        net_usage=net.usage,
        capacity_kwp=100.0,
    )
    assert surplus.total_kwh == 0.0

    solar = next(
        entry
        for entry in measure_entries(solar=point, surplus=surplus)
        if entry.kind.key == "solar"
    )
    assert "잉여" not in dict(solar.facts)
    # **발전이 0 이면 역률도 안 떨어진다** — 조정 줄도 없다 (59세션 12절).
    assert solar.slide_note == ""
    assert not any("자격요건" in line for line in solar.cautions), solar.cautions


# ===================================================================== 53세션 · 공통 규약


def _note_lines(slide: object, guide: object) -> list[str]:
    """슬라이드 아래쪽 작은 회색 글씨. **캡션과 각주가 여기 든다.**"""
    size = guide.type_scale.caption  # type: ignore[attr-defined]
    lines: list[str] = []
    for shape in slide.shapes:  # type: ignore[attr-defined]
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            sizes = {run.font.size.pt for run in paragraph.runs if run.font.size}
            if sizes == {size} and paragraph.text.strip():
                lines.append(paragraph.text.strip())
    return lines


def test_참고용_작은_글씨에_표식이_붙는다(full_sections: DocumentSections) -> None:
    """**※ 하나가 참고와 본문을 가른다** (53세션 1-1).

    작은 회색 글씨라는 것만으로는 본문의 끝인지 참고인지 갈리지 않는다.
    **그림 캡션은 제외한다** — 캡션은 그림의 이름이지 참고가 아니다.
    """
    from kwise.report.slides import NOTE_MARK, mark_note

    guide = load_design_guide()
    deck = build_slides(full_sections)
    captions = {figure.caption for entry in full_sections.measures for figure in entry.figures}
    marked = 0
    for slide in deck.slides:
        for line in _note_lines(slide, guide):
            if line.startswith(NOTE_MARK.strip()):
                marked += 1
                continue
            assert line in captions or "—" in line or "습니다" not in line, (
                f"※ 없는 참고 줄: {line}"
            )
    assert marked >= 5, f"※ 가 붙은 줄이 {marked}개뿐입니다."
    # 두 번 붙지 않는다.
    assert mark_note(mark_note("가나다")) == "※ 가나다"


def test_각주가_산식을_쉼표로_가른다() -> None:
    """**「·」 가 산식 자리에서 수식 기호로 읽힌다** (53세션 1-3)."""
    from kwise.report.narrative import GLOSSARY_KEYS, glossary_note

    note = glossary_note(GLOSSARY_KEYS["usage_pattern"])
    assert "÷" in note
    assert " · " not in note, note
    assert note.count(", ") == 2, note


def test_투자가_없으면_줄표와_즉시로_적는다(full_sections: DocumentSections) -> None:
    """**「0 원」 은 값이 아니라 없음이다** (53세션 1-5)."""
    from kwise.report.slides import IMMEDIATE, NO_INVESTMENT, slide_investment, slide_payback

    assert slide_investment("0원") == NO_INVESTMENT
    assert slide_investment("261,893,000원") == "261,893,000원"
    assert slide_payback("즉시 (투자 없음)") == IMMEDIATE
    assert slide_payback("31.7년") == "31.7년"
    assert slide_payback("미산출 — 투자비 미입력") == "미산출"

    text = _deck_text(build_slides(full_sections))
    assert "즉시 (투자 없음)" not in text
    # **투자비 칸에만 걸린다.** 절감액이 실제로 0원인 수단은 그대로 0원이다 —
    # 계산해 보니 0 인 값과 애초에 살 것이 없는 값은 다른 사실이다.
    investments = {
        row.cells[2].text
        for slide in build_slides(full_sections).slides
        for shape in slide.shapes
        if shape.has_table and shape.table.rows[0].cells[0].text == "개선 수단"
        for row in list(shape.table.rows)[1:]
    }
    assert "0원" not in investments, investments


def test_부록에_해석_문구가_없다(full_sections: DocumentSections) -> None:
    """**부록은 근거를 펼치는 자리다** (53세션 1-6). 읽는 법을 일러 주지 않는다."""
    from kwise.report.narrative import APPENDIX_LEAD

    deck = build_slides(full_sections)
    slide = _slide_by_key(deck, full_sections, "appendix")
    assert APPENDIX_LEAD not in _slide_text(slide)
    assert APPENDIX_LEAD not in _deck_text(deck)


def test_확실성이_슬라이드에_없다(full_sections: DocumentSections) -> None:
    """**등급을 산출물에서 뺐다** (53세션 1-4). 계산은 그대로다."""
    assert "확실성" not in _deck_text(build_slides(full_sections))
    assert all(entry.certainty for entry in full_sections.measures), "계산값은 남는다"


# ===================================================================== 53세션 · 2절 ESS 배치


def _ess_entry_with_spec() -> MeasureEntry:
    """목표별 사양 표가 달린 ESS 항목. **표 모양만 실물과 같으면 된다.**"""
    from kwise.measures import measure_kind
    from kwise.report import figures, frames

    png = figures.chart_palette  # 존재 확인용 — 아래에서 실제 png 를 굽지 않는다
    assert png is not None
    rows = [frames.ESS_SPEC_HEADER]
    rows.extend(
        (
            f"{5_240 - 40 * index:,.0f}~{5_260 - 40 * index:,.0f} kW",
            f"{73 + 40 * index:,.0f} kW",
            "150 kW",
            "100 kWh",
            "0.67h",
            "261,893,000",
            "8,264,000",
            "31.7년",
            "목표 미달 (실제 5,264 kW)" if index else "최단 회수기간",
        )
        for index in range(5)
    )
    return MeasureEntry(
        kind=measure_kind("ess"),
        conclusion="피크를 5,180 kW 까지 낮추려면 ESS 150 kW / 100 kWh 가 필요합니다.",
        saving="8,264,000원",
        investment="261,893,000원",
        payback="31.7년",
        certainty="중간~낮음",
        spec_table=tuple(rows),
        spec_caption=frames.ESS_SPEC_CAPTION,
    )


def test_사양_표가_아래_그림을_덮지_않는다(full_sections: DocumentSections) -> None:
    """**두 세션 연속 지적된 자리다** (53세션 2절).

    50세션에 격자가 붙어 표가 여섯 줄이 되자, 46세션의 ``height * 0.55`` 배분이
    그대로 표를 키워 그림이 **0.85 × 0.41in** 로 우겨 넣어졌다. 줄 수로 잡는다.
    """
    import dataclasses

    from pptx.util import Emu

    entry = _ess_entry_with_spec()
    sections = dataclasses.replace(full_sections, measures=(entry,))
    slide = _slide_by_key(build_slides(sections), sections, f"measure_{entry.kind.key}")
    boxes = [
        (Emu(shape.top).inches, Emu(shape.top + shape.height).inches, shape)
        for shape in slide.shapes  # type: ignore[attr-defined]
        if shape.top is not None and shape.height
    ]
    from kwise.report.frames import ESS_SPEC_HEADER

    tables = [
        box
        for box in boxes
        if box[2].has_table
        and tuple(cell.text for cell in box[2].table.rows[0].cells) == ESS_SPEC_HEADER
    ]
    assert len(tables) == 1, "사양 표가 하나여야 한다"
    table_top, table_bottom, shape = tables[0]
    assert table_bottom < 7.0, f"표가 슬라이드 밖으로 나갑니다: {table_bottom:.2f}in"
    for start, _end, other in boxes:
        if other is shape or start < table_top:
            continue
        assert start >= table_bottom - 1e-6, (
            f"표(밑변 {table_bottom:.2f}in) 아래 요소가 {start:.2f}in 에서 시작합니다."
        )


def test_사양_표_열이_고르게_나뉘지_않는다() -> None:
    """**긴 칸에 폭을 몰아 준다** (53세션 2절).

    아홉 열을 고르게 나누면 「5,220~5,240 kW」 가 두 줄로 흘러 표가 1.5배가 된다.
    """
    from kwise.report.design import load_design_guide
    from kwise.report.frames import ESS_SPEC_HEADER
    from kwise.report.slides import SPEC_TABLE_WIDTHS, _spec_lines

    assert len(SPEC_TABLE_WIDTHS) == len(ESS_SPEC_HEADER)
    assert abs(sum(SPEC_TABLE_WIDTHS) - 1.0) < 1e-9
    guide = load_design_guide()
    rows = _ess_entry_with_spec().spec_table
    width = guide.slide.content_width_in
    size = guide.type_scale.caption
    # **몰아 준 폭에서는 한 줄에 앉는다.**
    assert _spec_lines(rows, width=width, size=size) == len(rows)
    # **고르게 나누면 흐른다** — 이것이 46세션 배치가 무너진 까닭이다.
    even = tuple(1.0 / len(ESS_SPEC_HEADER) for _ in ESS_SPEC_HEADER)
    monkey = width * even[0] - 0.12
    from kwise.report.slides import _fitting_lines

    assert _fitting_lines(["5,240~5,260 kW"], span=monkey, size=size) == 2


# ===================================================================== 53세션 · 3절 잉여 장


def _surplus_result(*, weekday: float, weekend: float, holiday: float) -> object:
    """잉여 결과를 **값만 채워** 만든다. 갈래 시험은 계산이 아니라 문장을 본다."""
    import pandas as pd

    from kwise.measures.surplus import (
        CURTAIL_SCENARIO,
        EXTERNAL_SCENARIO,
        OFFSET_SCENARIO,
        SurplusResult,
        SurplusScenario,
    )

    total = weekday + weekend + holiday
    return SurplusResult(
        total_kwh=total,
        generation_kwh=total * 8,
        share_of_generation=0.125 if total else None,
        hour_distribution=pd.Series(dtype=float),
        weekday_kwh=weekday,
        weekend_kwh=weekend,
        holiday_kwh=holiday,
        scenarios=(
            SurplusScenario(OFFSET_SCENARIO, 2_545_000.0, "상계", "계약 변경"),
            SurplusScenario(EXTERNAL_SCENARIO, None, "단가 미입력", "구매자 발굴"),
            SurplusScenario(CURTAIL_SCENARIO, 0.0, "출력을 낮춘다", "없다"),
        ),
    )


def test_잉여가_있으면_장이_생기고_0이면_안_생긴다(full_sections: DocumentSections) -> None:
    """**41세션에 잉여 장이 사라지면서 보여 줄 자리가 없어졌다** (53세션 3-1)."""
    import dataclasses

    from kwise.report.document import surplus_page

    assert surplus_page(None, capacity_kwp=160.0) is None
    empty = _surplus_result(weekday=0.0, weekend=0.0, holiday=0.0)
    assert surplus_page(empty, capacity_kwp=160.0) is None  # type: ignore[arg-type]

    page = surplus_page(
        _surplus_result(weekday=105.0, weekend=20_000.0, holiday=3_311.0),  # type: ignore[arg-type]
        capacity_kwp=160.0,
        surplus_free_kwp=40.0,
    )
    assert page is not None
    # 태양광 장이 있어야 그 **다음**이 정해진다 — 수단을 모두 켠 덱으로 본다.
    base = dataclasses.replace(full_sections, measures=_all_measures(full_sections))
    with_page = dataclasses.replace(base, surplus=page)
    without = dataclasses.replace(base, surplus=None)
    keys = [item.key for item in slide_specs(with_page)]
    assert keys.count("surplus") == 1
    assert keys.index("surplus") == keys.index("measure_solar") + 1
    assert "surplus" not in [item.key for item in slide_specs(without)]
    assert len(slide_specs(with_page)) == len(slide_specs(without)) + 1
    # 목차도 한 줄 는다.
    assert SLIDE_TITLES["surplus"] in agenda_items(with_page)
    assert SLIDE_TITLES["surplus"] not in agenda_items(without)


def test_잉여_문장이_평일_휴일_비중으로_갈린다() -> None:
    """**이 건물 값에 맞춘 문장을 고정으로 박지 않는다** (53세션 3-2 · 4절 공통).

    소형 사무빌딩은 휴일이 99.6% 지만, 평일 낮에 문을 닫지 않는 건물은 반대다.
    """
    from kwise.report.document import surplus_page

    def lead(**parts: float) -> str:
        page = surplus_page(_surplus_result(**parts), capacity_kwp=160.0)  # type: ignore[arg-type]
        assert page is not None
        return page.lead

    assert "대부분 토·일·공휴일에 발생합니다" in lead(
        weekday=105.0, weekend=20_000.0, holiday=3_311.0
    )
    assert "대부분 평일에 발생합니다" in lead(weekday=9_000.0, weekend=500.0, holiday=500.0)
    assert "고르게 발생합니다" in lead(weekday=5_000.0, weekend=3_000.0, holiday=2_000.0)


def test_잉여_장이_셋을_다_싣고_금액_순으로_세운다(full_sections: DocumentSections) -> None:
    """**출력제어 줄이 없었다** (59세션 4·10·13절 · 목록 P4·P7·P16).

    53세션까지는 출력제어를 「기준선」 으로 보고 뺐다 (27세션이 화면에서 뺀
    것을 따랐다). **57세션에 그것이 기본 선택이 됐다** — 표에서 빼 두면 지금
    절감액에 든 것이 어디에도 없고, 「절감액에 반영」 표식이 어느 줄에도 안
    붙어 상계거래가 골라진 것처럼 읽힌다.

    **금액이 큰 순서다.** 58세션에 단가 기본값이 생겨 순위가 뒤집혔는데
    (외부 판매 > 상계거래), 정의 순서가 고정이라 표에서 그 사실이 안 보였다.
    """
    import dataclasses

    from kwise.measures.surplus import (
        CURTAIL_SCENARIO,
        EXTERNAL_SCENARIO,
        OFFSET_SCENARIO,
        SurplusScenario,
    )
    from kwise.report.document import (
        SURPLUS_CHOSEN_MARK,
        SURPLUS_SCENARIO_HEADER,
        surplus_page,
    )

    base = _surplus_result(weekday=105.0, weekend=20_000.0, holiday=3_311.0)
    # 58세션 뒤의 자리 — 두 줄 다 금액이 서고 **외부 판매가 더 크다.**
    priced = dataclasses.replace(
        base,  # type: ignore[type-var]
        scenarios=(
            SurplusScenario(OFFSET_SCENARIO, 2_283_732.0, "상계", "계약 변경"),
            SurplusScenario(EXTERNAL_SCENARIO, 2_947_148.0, "외부", "구매자 발굴"),
            SurplusScenario(CURTAIL_SCENARIO, 0.0, "출력을 낮춘다", "없다"),
        ),
    )
    page = surplus_page(
        priced,
        capacity_kwp=160.0,
        surplus_free_kwp=40.0,
        chosen_scenario=CURTAIL_SCENARIO,
    )
    assert page is not None
    assert page.scenario_rows[0] == SURPLUS_SCENARIO_HEADER
    assert [row[0] for row in page.scenario_rows[1:]] == [
        EXTERNAL_SCENARIO,
        OFFSET_SCENARIO,
        CURTAIL_SCENARIO,
    ]
    assert len(page.facts) == 4
    # **고른 것에만 붙는다.** 기본 선택은 출력제어다 (57세션).
    marked = [row[0] for row in page.scenario_rows[1:] if SURPLUS_CHOSEN_MARK in row[2]]
    assert marked == [CURTAIL_SCENARIO]

    sections = dataclasses.replace(
        full_sections, measures=_all_measures(full_sections), surplus=page
    )
    slide = _slide_by_key(build_slides(sections), sections, "surplus")
    text = _slide_text(slide)
    assert "출력제어" in text
    assert SURPLUS_CHOSEN_MARK in text
    assert "상계거래는 계약 변경과 역송 계량기가 필요합니다" in text
    assert "자격요건은 판정하지 않았습니다" in text


def test_금액을_못_낸_줄은_맨_뒤다() -> None:
    """**견줄 수 없는 것을 견주는 자리에 세우지 않는다** (59세션 10절).

    단가를 비우면 「미산출」 로 돌아간다 (58세션). 그 줄을 금액 순 어딘가에
    끼워 넣으면 순위가 값이 아니라 정의 순서로 정해진다.
    """
    from kwise.measures.surplus import CURTAIL_SCENARIO, EXTERNAL_SCENARIO, OFFSET_SCENARIO
    from kwise.report.document import surplus_page

    # ``_surplus_result`` 의 외부 판매는 단가 미입력(``None``)이다.
    page = surplus_page(
        _surplus_result(weekday=105.0, weekend=20_000.0, holiday=3_311.0),  # type: ignore[arg-type]
        capacity_kwp=160.0,
        surplus_free_kwp=40.0,
    )
    assert page is not None
    assert [row[0] for row in page.scenario_rows[1:]] == [
        OFFSET_SCENARIO,
        CURTAIL_SCENARIO,
        EXTERNAL_SCENARIO,
    ]
    assert page.scenario_rows[-1][1] == "미산출"


def test_상계거래_비고가_기간말_잔여를_늘_적는다() -> None:
    """**「120원인데 왜 이 금액인가」 를 표가 닫는다** (59세션 10절).

    표 아래 각주가 「상계거래 SMP 120원/kWh」 라고 적는데 금액은 그 단가와
    무관할 수 있다 — 잉여가 사용량에 다 잠기면 SMP 가 곱해질 몫이 없다.
    잔여 0 을 적지 않으면 그 사실이 표 어디에도 없다.
    """
    import dataclasses

    from kwise.measures.surplus import OFFSET_SCENARIO, OffsetSettlement, SurplusScenario
    from kwise.report.document import surplus_page

    base = _surplus_result(weekday=105.0, weekend=20_000.0, holiday=3_311.0)
    settled = dataclasses.replace(
        base,  # type: ignore[type-var]
        offset=OffsetSettlement(
            months=(),
            deducted_kwh=23_416.0,
            deducted_won=2_283_732.0,
            remaining_kwh=0.0,
            settles_cash=True,
            smp_price_won_per_kwh=120.0,
            smp_won=0.0,
        ),
        scenarios=(
            SurplusScenario(OFFSET_SCENARIO, 2_283_732.0, "상계", "계약 변경"),
            *base.scenarios[1:],  # type: ignore[attr-defined]
        ),
    )
    page = surplus_page(settled, capacity_kwp=160.0, surplus_free_kwp=40.0)
    assert page is not None
    remark = next(row[2] for row in page.scenario_rows[1:] if row[0] == OFFSET_SCENARIO)
    assert "당월 차감 23,416 kWh" in remark
    assert "기간 말 잔여 0 kWh" in remark


def test_잉여_장이_적용_단가를_밝힌다(full_sections: DocumentSections) -> None:
    """**표 아래에 무슨 단가로 산출했는지 적는다** (58세션 2절).

    단가에 기본값이 생겨 「미산출」 이 사라졌다. 금액이 늘 나오면 어느 단가로
    나온 것인지를 함께 적지 않고는 읽는 사람이 참고값을 확정값으로 본다.
    **숫자를 박지 않는다** — 사용자가 고친 값이 그대로 실린다.
    """
    import dataclasses

    from kwise.measures.surplus import (
        EXTERNAL_SCENARIO,
        OFFSET_SCENARIO,
        SurplusScenario,
    )
    from kwise.report.document import surplus_page

    base = _surplus_result(weekday=105.0, weekend=20_000.0, holiday=3_311.0)
    priced = dataclasses.replace(
        base,  # type: ignore[type-var]
        external_price_won_per_kwh=140.0,
        scenarios=(
            SurplusScenario(OFFSET_SCENARIO, 2_545_000.0, "상계", "계약 변경"),
            SurplusScenario(EXTERNAL_SCENARIO, 3_267_824.0, "외부", "구매자 발굴"),
            base.scenarios[2],  # type: ignore[attr-defined]
        ),
    )
    page = surplus_page(priced, capacity_kwp=160.0, surplus_free_kwp=40.0)
    assert page is not None
    assert "외부 판매 140원/kWh" in " ".join(page.notes)
    assert "참고용입니다" in " ".join(page.notes)
    # **단가가 없으면 그 문장도 없다** — 지어낸 단가를 밝히지 않는다.
    plain = surplus_page(base, capacity_kwp=160.0, surplus_free_kwp=40.0)  # type: ignore[arg-type]
    assert plain is not None
    assert "원/kWh" not in " ".join(plain.notes)

    sections = dataclasses.replace(
        full_sections, measures=_all_measures(full_sections), surplus=page
    )
    text = _slide_text(_slide_by_key(build_slides(sections), sections, "surplus"))
    assert "외부 판매 140원/kWh" in text


# ===================================================================== 53세션 · 4절 갈래


def test_결측_갈래가_셋이다(sample_report: object) -> None:
    """**결측이 없으면 뒷문장이 없다** (53세션 4-1).

    39세션은 「가장 심한 달」 하나만 짚었다. 여러 달이 높으면 이름을 다 이어
    붙이는 대신 개수로 적는다 — 해석 한 줄이 세 줄로 흐르지 않게.
    """
    from dataclasses import replace

    from kwise.report.narrative import building_lead

    assert "결측" in building_lead(None)
    clean = replace(sample_report, missing_slots=0, missing_ratio=0.0, monthly=())  # type: ignore[type-var]
    text = building_lead(clean)
    assert "결측이 없어" in text
    assert "낮게 잡혔을" not in text, text

    one = building_lead(sample_report)  # type: ignore[arg-type]
    assert "낮게 잡혔을 수 있습니다" in one, one
    assert "개 달이" not in one, one

    heavy = replace(
        sample_report,  # type: ignore[type-var]
        monthly=tuple(
            replace(month, ratio=0.4, missing_slots=100)
            for month in sample_report.monthly  # type: ignore[attr-defined]
        ),
    )
    many = building_lead(heavy)
    assert "개 달이 크게 비어" in many, many


def test_부하패턴_문장_둘이_각각_갈린다(sample_diagnosis: Diagnosis) -> None:
    """**조합이 넷이다** (53세션 4-2). 부하율이 높으면 「짧은 피크」 가 아니다."""
    from dataclasses import replace

    from kwise.report.narrative import base_load_high, load_factor_flat, pattern_lead

    flat = load_factor_flat() + 0.05
    peaky = load_factor_flat() - 0.05
    high = base_load_high() + 0.05
    low = base_load_high() - 0.05
    base = replace(sample_diagnosis.pattern, load_factor=peaky, base_load_ratio=high)
    assert "짧은 피크" in pattern_lead(base)
    # **「좁습니다」 를 쓰지 않는다** (59세션 6절). 폭이 좁다는 말로 읽혀 무엇이
    # 좁은지 되묻게 된다 — 뜻은 「여력이 제한적이다」 이다.
    assert "충전 여력이 제한적입니다" in pattern_lead(base)
    assert "좁습니다" not in pattern_lead(base)
    both = pattern_lead(replace(base, load_factor=flat, base_load_ratio=low))
    assert "하루 내내 고르게" in both
    assert "충전 여력이 있습니다" in both
    # 값이 없으면 그 문장이 통째로 빠진다.
    assert "기저부하" not in pattern_lead(replace(base, base_load_ratio=None))
    assert "산출하지 못했습니다" in pattern_lead(
        replace(base, load_factor=None, base_load_ratio=None)
    )


def test_요금구조가_세_갈래다() -> None:
    """**가운데를 가운데라고 적는다** (53세션 4-6)."""
    from types import SimpleNamespace

    from kwise.report.narrative import base_fee_share_high, base_fee_share_low, structure_lead

    def lead(share: float) -> str:
        return structure_lead(
            SimpleNamespace(  # type: ignore[arg-type]
                base_won=share * 100.0,
                energy_won=(1 - share) * 100.0,
                total_won=100.0,
                bill=SimpleNamespace(total_power_factor_won=0.0),
            )
        )

    low = lead(base_fee_share_low() - 0.05)
    assert "나머지는 전력량요금입니다" in low
    # **뒷문장이 지침에서 결과로 바뀌었다** (59세션 7절). 「단가가 낮은 시간대로
    # 부하를 옮기거나 사용량을 줄이는 방안」 은 앞말에서 곧바로 따라오는 말이라
    # 읽는 사람이 얻는 것이 없었다 — 높음 갈래의 거울 문장으로 바꿨다.
    assert "피크를 낮춰도 줄어드는 몫이 작습니다" in low
    assert "단가가 낮은 시간대로 부하를 옮기거나" not in low
    # **갈래는 그대로 셋이다** (54세션). 뒷문장이 셋을 가르는 자리다.
    assert "함께 큽니다" in lead((base_fee_share_low() + base_fee_share_high()) / 2)
    assert "최대수요를 낮추는 방안을 먼저" in lead(base_fee_share_high() + 0.05)
    tails = {
        lead(base_fee_share_low() - 0.05).split(" — ", 1)[-1],
        lead((base_fee_share_low() + base_fee_share_high()) / 2).split(" — ", 1)[-1],
        lead(base_fee_share_high() + 0.05).split(" — ", 1)[-1],
    }
    assert len(tails) == 3, tails
    assert "산출하지 못했습니다" in structure_lead(
        SimpleNamespace(  # type: ignore[arg-type]
            base_won=0.0,
            energy_won=0.0,
            total_won=0.0,
            bill=SimpleNamespace(total_power_factor_won=0.0),
        )
    )


def test_요약_근거에_0원인_수단이_안_든다() -> None:
    """**계약전력 조정은 0원인데 근거로 들고 있었다** (53세션 4-7)."""
    from types import SimpleNamespace

    from kwise.report.narrative import measure_summary_lead

    def lead(switch: float | None, contract: float | None) -> str:
        summary = SimpleNamespace(
            tariff_switch_saving_won=switch,
            contract_saving_won=contract,
            no_investment_saving_won=(switch or 0.0) + (contract or 0.0),
        )
        return measure_summary_lead(SimpleNamespace(summary=summary), "5,358만원")  # type: ignore[arg-type]

    only = lead(53_575_000.0, 0.0)
    assert "요금제 전환입니다" in only, only
    assert "계약전력" not in only, only
    assert "계약전력 조정" in lead(53_575_000.0, 1_000_000.0)
    assert "미산출" not in lead(53_575_000.0, None)
    assert "줄일 수 있는 몫은 없습니다" in lead(0.0, 0.0)


def test_갑_종별_ESS_장이_결론_한_줄로_선다(
    sample_usage: UsageData, tariff: TariffTable, sample_bill: BillingResult
) -> None:
    """**장이 통째로 사라지고 있었다** (59세션 2절 · 목록 P2·P10).

    56세션이 갑 종별을 훑기 전에 끊으면서 ``EssOptimum.points`` 가 비었고, 화면
    경로는 목표가 있을 때만 곡선을 담아 넘겨(``ess_curve``) **성립 불가 갈래
    둘이 열리지 않았다.** 켠 수단이 산출물에서 사라지면 「검토하지 않았다」 와
    구분되지 않는다 (48세션).

    문장은 계산이 낸 것을 그대로 옮긴다 — 슬라이드가 새로 짓지 않는다.
    """
    from kwise.measures.ess import (
        BASE_FEE_ON_CONTRACT_CONCLUSION,
        ess_target_curve,
        refine_ess_target,
    )
    from kwise.report.document import measure_entries
    from kwise.tariff import TariffSelection

    selection = TariffSelection("general_a_1", "high_a", "I")
    curve = ess_target_curve(
        sample_usage.kw,
        15,
        baseline_demand_kw=5_293.44,
        base_fee_won_per_kw=float(tariff.rates(selection).base_won_per_kw),
    )
    optimum = refine_ess_target(
        sample_usage, tariff, selection, curve=curve, baseline=sample_bill
    )
    assert not optimum.viable and not optimum.points

    entry = next(
        item
        for item in measure_entries(ess_optimum=optimum, ess_curve=curve)
        if item.kind.key == "ess"
    )
    assert entry.conclusion == BASE_FEE_ON_CONTRACT_CONCLUSION.format(
        label=tariff.contract("general_a_1").label
    )
    assert "제68조" in entry.conclusion
    # **잰 점이 없으면 표를 그리지 않는다** — 화면 카드와 같은 규약이다.
    assert entry.spec_table == ()
    assert not entry.spec_caption
    # 까닭이 지표 둘에 보인다 — 기본요금이 붙는 자리가 피크가 아니다.
    assert ("기본요금 기준", "계약전력") in entry.facts
    # 마진 갈래의 지표는 쓰지 않는다 — 잰 것이 없어 값이 없다.
    assert not any(label == "성립 한계 방전시간" for label, _value in entry.facts)


def test_마진_미달_ESS_장은_표를_남긴다(
    sample_usage: UsageData, tariff: TariffTable, sample_bill: BillingResult
) -> None:
    """**성립 불가 갈래 둘이 다르다** (59세션 2절). 이쪽은 잰 점이 있다."""
    from kwise.measures.ess import ess_target_curve, refine_ess_target
    from kwise.report.document import measure_entries
    from kwise.tariff import TariffSelection

    curve = ess_target_curve(
        sample_usage.kw, 15, baseline_demand_kw=5_293.44, base_fee_won_per_kw=70.0
    )
    optimum = refine_ess_target(
        sample_usage,
        tariff,
        TariffSelection("general_b", "high_a", "I"),
        curve=curve,
        baseline=sample_bill,
    )
    assert not optimum.viable and optimum.points

    entry = next(
        item
        for item in measure_entries(ess_optimum=optimum, ess_curve=curve)
        if item.kind.key == "ess"
    )
    assert entry.spec_table
    assert any(label == "성립 한계 방전시간" for label, _value in entry.facts)
    assert "성립하는 목표가 없어" in entry.saving


def test_계약전력이_세_갈래다(sample_usage: UsageData, sample_bill: BillingResult) -> None:
    """**초과 위약 갈래가 없었다** (53세션 4-9). 계산은 이미 알고 있었다."""
    from kwise.measures import evaluate_contract_adjustment
    from kwise.report.document import measure_entries

    def conclusion(contract_kw: float) -> str:
        contract = evaluate_contract_adjustment(sample_usage, sample_bill, contract_kw=contract_kw)
        entry = next(
            item for item in measure_entries(contract=contract) if item.kind.key == "contract"
        )
        return entry.conclusion

    # 계약전력이 실측 최대보다 낮으면 초과 위약이다.
    over = conclusion(3_000.0)
    assert "초과 위약 검토 대상입니다" in over, over
    assert "상향을 검토해야 합니다" in over
    assert "낮출 여지가" in conclusion(6_500.0)


def test_장이_따로_적는_줄은_제_줄에_선다(full_sections: DocumentSections) -> None:
    """**다른 종류의 말을 「·」 로 잇지 않는다** (59세션 14절).

    미산출 사유 뒤에 이어 붙이니 「역률 영향 반영 시 279,249,000원」 이 또 하나의
    미산출 사유처럼 읽혔다. ``slide_note`` 는 제 ※ 를 단다.
    """
    import dataclasses

    from kwise.report.slides import NOTE_MARK, _measure_note

    measures = _all_measures(full_sections)
    entry = next(item for item in measures if item.kind.key == "solar")
    noted = dataclasses.replace(entry, slide_note="역률 영향 반영 시 279,249,000원")
    # 각주 본문에는 들어가지 않는다 — 슬라이드가 따로 한 줄로 깐다.
    assert "역률 영향" not in _measure_note(noted)

    sections = dataclasses.replace(
        full_sections, measures=tuple(_replace_solar(measures, noted))
    )
    slide = _slide_by_key(build_slides(sections), sections, "measure_solar")
    lines = [
        line
        for shape in slide.shapes
        if shape.has_text_frame
        for line in shape.text_frame.text.split("\n")
        if line.startswith(NOTE_MARK.strip())
    ]
    assert any(line.strip().startswith("※ 역률 영향 반영 시") for line in lines), lines


def _replace_solar(
    measures: tuple[MeasureEntry, ...], swapped: MeasureEntry
) -> list[MeasureEntry]:
    return [swapped if item.kind.key == "solar" else item for item in measures]


def test_부록이_효과_있는_수단을_빠뜨리지_않는다(full_sections: DocumentSections) -> None:
    """**사용자가 물었다 — 「선택요금 조정이 빠진 이유는 효과가 0 이라서인가요?」**
    (59세션 11절 · 목록 P11).

    규칙은 「값이 0 이거나 미산출인 수단은 뺀다」 이고, 뒤집으면 **효과가 있으면
    반드시 실린다** 는 뜻이다. 그 역이 실제로 참인지를 못박는다.

    **판정은 ``MeasureEntry.has_saving`` 이 쥔다** — 39세션이 세웠고, 「0 이거나」
    는 59세션 11절에 붙었다. (60세션 7절에 **55세션** 참조를 여기서 걷어냈다.
    55세션 자리는 비어 있다.)

    각주는 규칙을 그대로 적는다 — 「산출되지 않은」 이라고만 적으면 계약전력
    조정처럼 **산출은 됐는데 0 인** 줄이 못 낸 것으로 읽힌다.
    """
    import dataclasses
    import inspect

    from kwise.report.slides import _appendix_note, appendix_pages, measure_slide_title

    sections = dataclasses.replace(full_sections, measures=_all_measures(full_sections))
    titles = " ".join(page.title for page in appendix_pages(sections))
    sheets = {sheet.key for sheet in sections.worksheets}
    for entry in sections.measures:
        name = measure_slide_title(entry)
        if entry.has_saving and entry.kind.key in sheets:
            assert name in titles, f"효과가 있는데 부록에 없습니다: {name}"
        if not entry.has_saving:
            assert name not in titles, f"효과가 0·미산출인데 부록에 있습니다: {name}"

    # **근거 표가 없으면 효과가 있어도 못 싣는다.** 화면 경로는 여섯 수단 모두
    # 근거 표를 만든다 — 그 배선이 끊기면 부록에서 조용히 사라지므로 못박는다.
    from kwise.ui.views.compare import _MeasureResults

    source = inspect.getsource(_MeasureResults.worksheets)
    for name in (
        "tariff_switch_worksheet",
        "contract_worksheet",
        "demand_response_worksheet",
        "power_factor_worksheet",
        "solar_worksheet",
        "ess_worksheet",
    ):
        assert name in source, f"근거 표를 만들지 않는 수단이 있습니다: {name}"

    note = _appendix_note(sections)
    dropped = [item for item in sections.measures if not item.has_saving]
    if dropped:
        assert "절감액이 0 이거나 산출되지 않은 수단" in note, note
        for entry in dropped:
            assert measure_slide_title(entry) in note, note
    assert "자리가 모자라" not in note


def test_역률_영향을_큰_글자에_녹이지_않는다(
    sample_usage: UsageData, tariff: TariffTable, sample_bill: BillingResult
) -> None:
    """**둘을 나눠 보인다** (59세션 12절 · 목록 P6 · 31세션 독립 평가).

    태양광이 유효전력만 상쇄해 역률이 떨어지고 역률요금이 는다 — 사실이고
    계산이 이미 내던 값인데(``power_factor_extra_won``) **어느 산출물에도 금액
    으로 서 있지 않았다.** 큰 글자는 조정 전 값이고 조정값은 곁에 적는다.
    """
    import dataclasses

    import pandas as pd

    from kwise.measures import solar_point
    from kwise.report.document import measure_entries
    from kwise.report.excel import measure_summary_frame
    from kwise.report.narrative import power_factor_adjusted_saving
    from kwise.tariff import TariffSelection

    selection = TariffSelection("general_b", "high_a", "I")
    index = pd.DatetimeIndex(sample_usage.kw.index)
    unit = pd.Series(0.0, index=index, name="kw")
    point = solar_point(sample_usage, tariff, selection, unit, 100.0, baseline=sample_bill)
    # 역률이 떨어진 자리를 만든다 — 발전이 0 인 점에는 조정이 없다.
    dropped = dataclasses.replace(
        point,
        generation_kwh=1_940_781.0,
        self_consumed_kwh=1_940_781.0,
        self_consumption_ratio=1.0,
        total_saving_won=29_564_000.0,
        annual_saving_won=29_564_000.0,
        power_factor_extra_won=156_000.0,
        power_factor_after_pct=91.8,
    )
    line = power_factor_adjusted_saving(saving_won=29_564_000.0, extra_won=156_000.0)
    assert line == "역률 영향 반영 시 29,408,000원"

    entry = next(
        item
        for item in measure_entries(solar=dropped, base_fee_months=12.0)
        if item.kind.key == "solar"
    )
    # **큰 글자는 조정 전이다.**
    assert entry.saving_annual.startswith("29,564,000원"), entry.saving_annual
    assert "역률" not in entry.saving_annual
    # 곁에 적는다 — PPT 는 각주, Word 는 주의사항 목록. **같은 문장이다.**
    assert line in entry.slide_note
    assert line in entry.cautions

    # Excel 도 같은 문장을 쓴다.
    frame = measure_summary_frame(solar=dropped, base_fee_months=12.0)
    key = next(name for name in frame.index if str(name).startswith("태양광"))
    row = frame.loc[key]
    assert line in str(row["비고"]), row["비고"]
    assert str(row["절감액(원)"]).startswith("29,564,000")

    # **영향이 0 이면 어디에도 줄이 없다.**
    flat = dataclasses.replace(dropped, power_factor_extra_won=0.0)
    quiet = next(
        item
        for item in measure_entries(solar=flat, base_fee_months=12.0)
        if item.kind.key == "solar"
    )
    assert "역률 영향 반영 시" not in quiet.slide_note
    assert not [line for line in quiet.cautions if "역률 영향 반영 시" in line]


def test_합산효과는_태양광_역률_영향을_반영하지_않는다() -> None:
    """**이중 반영이 아니다** (59세션 12절). 오히려 아예 들어가지 않는다.

    조합 요금은 :class:`BillingOptions` 의 ``power_factor_pct`` 하나로만 역률을
    본다 — 켠 「역률 개선」 의 목표값이거나 기준선 값이다. 태양광이 떨어뜨리는
    역률은 그 옵션에 들어가지 않으므로 2단계 카드가 참고로 내는 값이다.
    **고치려면 계산을 바꿔야 하고, 그것은 이 세션의 범위가 아니다.**
    """
    import inspect

    from kwise.compare import combination

    source = inspect.getsource(combination)
    assert "power_factor_after_pct" not in source
    assert "power_factor_extra_won" not in source


def test_각주에_미산출이_두_번_서지_않는다() -> None:
    """**앞에도 뒤에도 붙는다** (53세션 9절 · 59세션 8절).

    태양광 투자비 사유는 「미산출 — …」 라 **앞**에 겹쳤고, 계약전력 사유는
    「하한 규정 미확인 — 금액 미산출」 이라 **뒤**에 겹쳤다.
    """
    from kwise.report.slides import _measure_note, _trim_repeat

    assert _trim_repeat("하한 규정 미확인 — 금액 미산출", "미산출") == "하한 규정 미확인"
    # 꼬리가 머리말이 아니면 그대로 둔다.
    assert _trim_repeat("정산 단가 미입력", "미산출") == "정산 단가 미입력"
    assert (
        _trim_repeat("요금적용전력 하한 30% 적용, 12.00개월분 재계산", "기본요금 변화없음")
        == "요금적용전력 하한 30% 적용, 12.00개월분 재계산"
    )

    from kwise.measures import measure_kind
    from kwise.report.document import MeasureEntry

    entry = MeasureEntry(
        kind=measure_kind("contract"),
        conclusion="",
        saving="미산출 — 하한 규정 미확인 — 금액 미산출",
        investment="—",
        payback="즉시",
        certainty="높음",
    )
    note = _measure_note(entry)
    assert note.count("미산출") == 1, note
    assert "하한 규정 미확인" in note


def test_계약종별을_바꾸면_계약전력_결과가_따라온다() -> None:
    """**캐시 열쇠에 계약종별이 없었다** (59세션 8절).

    하한 비율은 종별 속성이라 ``BillingResult`` 가 들고 오는데 그 인자는
    ``_bill`` 이라 열쇠에서 빠진다. 다른 캐시 함수는 모두 ``form`` 을 열쇠에
    두는데 여기만 없어, **한 세션에서 종별을 바꾸면 앞 결과가 다시 나왔다.**
    """
    import inspect

    from kwise.ui.cache import cached_contract_adjustment

    params = list(inspect.signature(cached_contract_adjustment).parameters)
    keys = [name for name in params if not name.startswith("_")]
    assert "form" in keys, params


def test_계약전력_각주가_하한_미결착을_말한다(
    sample_usage: UsageData, sample_bill: BillingResult
) -> None:
    """**「하향 여지 8 kW」 옆의 0원이 설명 없이 서 있었다** (59세션 9절).

    여지는 「낮출 수 있는가」 이고 절감액은 「낮추면 돈이 주는가」 다. 둘이
    갈리는 까닭을 계산이 이미 안내로 내고 있었는데(화면 산출 근거에 있다)
    슬라이드만 그것을 안 읽었다 — **문장을 새로 짓지 않는다.**
    """
    from kwise.measures import evaluate_contract_adjustment
    from kwise.report.document import CONTRACT_FLOOR_NOT_BINDING_FACT, measure_entries

    # 요금적용전력 5,293 kW 가 하한(계약전력의 30%)보다 훨씬 커서 안 걸린다.
    contract = evaluate_contract_adjustment(sample_usage, sample_bill, contract_kw=6_000.0)
    assert not contract.saving_won
    entry = next(
        item for item in measure_entries(contract=contract) if item.kind.key == "contract"
    )
    assert entry.slide_note == (
        "하한이 요금적용전력에 걸리지 않아 계약전력을 낮춰도 기본요금이 줄지 않습니다."
    )
    assert any(
        item.fact == CONTRACT_FLOOR_NOT_BINDING_FACT for item in contract.notices
    ), "화면 산출 근거에 있던 안내다 — 문구를 새로 짓지 않았다."


def test_역률이_세_갈래다(
    sample_usage: UsageData, tariff: TariffTable, sample_bill: BillingResult
) -> None:
    """**기준 미달이면 감액이 아니라 추가요금 회피다** (53세션 4-11)."""
    from kwise.measures import evaluate_power_factor
    from kwise.report.document import measure_entries
    from kwise.tariff import TariffSelection

    selection = TariffSelection("general_b", "high_a", "I")

    def conclusion(current: float, target: float) -> str:
        result = evaluate_power_factor(
            sample_usage,
            tariff,
            selection,
            current_pct=current,
            target_pct=target,
            baseline=sample_bill,
        )
        entry = next(
            item
            for item in measure_entries(power_factor=result)
            if item.kind.key == "power_factor"
        )
        return entry.conclusion

    penalty = conclusion(85.0, 97.0)
    assert "추가요금이 붙습니다" in penalty, penalty
    assert "추가요금이 없어지고 감액을 받아" in penalty
    assert "올리면" in conclusion(92.0, 97.0)
    assert "올릴 여지가 없습니다" in conclusion(97.0, 97.0)


def test_태양광_괄호가_면적과_한계를_나란히_놓지_않는다() -> None:
    """**2,038 kWp 를 지을 수 있는 것처럼 읽혔다** (53세션 4-12)."""
    from kwise.report.document import _solar_conclusion

    class _Point:
        capacity_kwp = 160.0
        generation_kwh = 194_078.0
        total_saving_won = 30_647_000.0
        surplus_kwh = 23_416.0

    point = _Point()
    ample = _solar_conclusion(point, 2_038.0, 2_000.0)  # type: ignore[arg-type]
    assert "설치 가능 면적 2,000 m² 가 허용하는 상한입니다" in ample
    assert "전부 자가소비됩니다" in ample
    assert "2,038" not in ample, "참고값을 실제 상한과 나란히 놓지 않는다"

    tight = _solar_conclusion(point, 40.0, 2_000.0)  # type: ignore[arg-type]
    assert "잉여 없이 지을 수 있는 것은 40 kWp 까지입니다" in tight
    assert "23,416 kWh 가 남습니다" in tight
    # 잉여 장이 뒤따르면 그 수는 다음 장이 맡는다.
    follows = _solar_conclusion(point, 40.0, 2_000.0, surplus_page_follows=True)  # type: ignore[arg-type]
    assert "23,416" not in follows, follows


def test_ESS_표식의_뜻이_표_아래에_있다() -> None:
    """**대형 표에 「마진 미달」 이 셋인데 뜻이 어디에도 없었다** (53세션 4-13)."""
    from kwise.measures import MARGIN_SHORT, SHORTEST_PAYBACK, TARGET_MISSED, spec_mark_note

    note = spec_mark_note([SHORTEST_PAYBACK, MARGIN_SHORT, MARGIN_SHORT, ""])
    assert note.startswith("표식 — ")
    assert SHORTEST_PAYBACK in note and MARGIN_SHORT in note
    # **표에 없는 표식은 설명하지 않는다.**
    assert TARGET_MISSED not in note, note
    # 값이 붙은 표식도 알아본다.
    assert TARGET_MISSED in spec_mark_note([f"{TARGET_MISSED} (실제 5,264 kW)"])
    assert spec_mark_note([]) == ""


def test_조합_문장이_수단_수로_갈린다() -> None:
    """**겹칠 것이 없으면 겹침을 설명하지 않는다** (53세션 4-14)."""
    from types import SimpleNamespace

    from kwise.report.narrative import COMBINATION_LEAD, SINGLE_MEASURE_LEAD, combination_lead

    assert combination_lead(None) == COMBINATION_LEAD
    assert combination_lead(SimpleNamespace(combinations=(1, 2))) == SINGLE_MEASURE_LEAD
    assert combination_lead(SimpleNamespace(combinations=(1, 2, 3))) == COMBINATION_LEAD


def _peak_stub(
    values: dict[str, float], *, demand_months: tuple[int, ...] = (7, 8, 9, 12, 1, 2)
) -> object:
    """월별 최대수요만 채운 진단 대역. **갈래 시험은 계산이 아니라 문장을 본다.**"""
    from types import SimpleNamespace

    import pandas as pd

    index = pd.PeriodIndex(list(values), freq="M")
    frame = pd.DataFrame(
        {"max_demand_kw": list(values.values()), "demand_basis_kw": list(values.values())},
        index=index,
    )
    return SimpleNamespace(
        peak=SimpleNamespace(monthly=frame, demand_months=demand_months, top_n=100)
    )


def _quality_stub(flagged: tuple[str, ...]) -> object:
    from types import SimpleNamespace

    import pandas as pd

    return SimpleNamespace(
        flagged_months=tuple(
            SimpleNamespace(month=pd.Period(month, freq="M")) for month in flagged
        )
    )


def test_5장이_그림과_같은_것을_말하고_세_갈래를_탄다(tariff: TariffTable) -> None:
    """**문장이 월별 최대수요를 말한다** (53세션 4-3).

    39세션은 여기에 정오 비중 판정을 적었는데 그림은 월별 최대수요였다.
    """
    from kwise.report.narrative import peak_month_lead

    # 기본 — 대상월이고 다음 달과 벌어져 있다.
    alone = peak_month_lead(
        _peak_stub({"2023-08": 5_300.0, "2023-07": 4_000.0}), None, tariff
    )
    assert alone == "8월에 최대수요가 가장 높고, 이 시기에 요금적용전력이 결정됩니다."

    # ③ 다음 달과 차이가 작으면 계절로 적는다.
    season = peak_month_lead(
        _peak_stub({"2023-08": 5_300.0, "2023-07": 5_290.0}), None, tariff
    )
    assert season.startswith("여름(6~8월)에 최대수요가 높고"), season
    # **한 해를 넘어가는 계절도 사람이 읽는 대로 적는다.**
    winter = peak_month_lead(
        _peak_stub({"2023-12": 6_100.0, "2024-01": 6_090.0}), None, tariff
    )
    assert winter.startswith("겨울(11~2월)"), winter

    # ① 대상월이 아니면 이월되지 않는다 — 실제로 정한 달을 함께 적는다.
    spring = peak_month_lead(
        _peak_stub({"2023-05": 5_300.0, "2023-08": 4_000.0}), None, tariff
    )
    assert "5월에 최대수요가 가장 높으나" in spring
    assert "요금적용전력은 8월 값으로 결정됩니다" in spring, spring

    # ② 신뢰 제한 달은 근거로 쓰지 않는다.
    limited = peak_month_lead(
        _peak_stub({"2023-08": 5_300.0, "2023-07": 4_000.0}),
        _quality_stub(("2023-08",)),
        tariff,
    )
    assert "결측이 많아 신뢰가 제한됩니다" in limited
    assert "결측이 적은 달 가운데는 7월이 가장 높습니다" in limited, limited

    # 값이 없으면 지어내지 않는다.
    assert "산출하지 못했습니다" in peak_month_lead(_peak_stub({}), None, tariff)


def test_DR_문장이_날_수로_세_갈래다(sample_diagnosis: Diagnosis) -> None:
    """**「245일 가운데 245일만」 은 성립하지 않는다** (53세션 4-10).

    C3 평탄형에서 거래 가능일이 전부 저부하로 잡힌다 — 「만」 은 적다는 뜻이라
    사실과 반대로 읽혔다.
    """
    from dataclasses import replace

    from kwise.report.narrative import dr_lead

    profile = sample_diagnosis.dr
    assert profile is not None
    assert dr_lead(None) == ""
    assert "내려오는 평일이 없어" in dr_lead(replace(profile, low_load_days=()))
    days = tuple(range(profile.eligible_days))
    assert "전부가 부하가 쉬는 날" in dr_lead(replace(profile, low_load_days=days))  # type: ignore[arg-type]
    assert "일만 부하가 쉬는 날" in dr_lead(replace(profile, low_load_days=days[:2]))  # type: ignore[arg-type]


# ===================================================================== 53세션 · 5절 차액 라벨


def test_차액_라벨이_0선_반대쪽에_선다() -> None:
    """**막대 안쪽에 있어 읽히지 않았다** (53세션 5절).

    「-0.54억」 이 파란 막대 위에 얹혀 있었다. 0 선 건너편은 어느 자료에서도
    비어 있으므로 그쪽에 두면 겹칠 일이 없다.
    """
    from kwise.report.figures import delta_label_place

    pad = 0.05
    down, down_align = delta_label_place(-0.54, pad)
    assert down > 0 and down_align == "bottom", (down, down_align)
    up, up_align = delta_label_place(0.31, pad)
    assert up < 0 and up_align == "top", (up, up_align)
    # 「현행」(0) 은 줄어드는 쪽과 같은 자리에 선다 — 겹칠 막대가 없다.
    zero, zero_align = delta_label_place(0.0, pad)
    assert zero > 0 and zero_align == "bottom"


# ===================================================================== 53세션 · 6절 계약전력


def test_계약전력_장이_근거를_먼저_세운다(
    sample_usage: UsageData, sample_bill: BillingResult, full_sections: DocumentSections
) -> None:
    """**근거가 결과보다 먼저다** (53세션 6-1·6-2·6-3).

    「6,000 kW 를 5,823 kW 로」 는 여유가 얼마나 있는지를 보고 나서야 판단할 수
    있는 값인데, 절감액·투자비·회수기간이 위에 서 있었다.
    """
    import dataclasses

    from pptx.util import Emu

    from kwise.measures import evaluate_contract_adjustment
    from kwise.report.document import measure_entries

    contract = evaluate_contract_adjustment(sample_usage, sample_bill, contract_kw=6_000.0)
    entry = next(
        item for item in measure_entries(contract=contract) if item.kind.key == "contract"
    )
    assert entry.facts_first
    labels = [name for name, _value in entry.facts]
    assert labels == ["현재 계약전력", "요금적용전력", "여유", "하향 여지"]
    # **여유는 %다** — 화면과 같은 산식이다.
    assert dict(entry.facts)["여유"].endswith("%"), entry.facts
    # **그림이 생겼다.**
    assert entry.figure is not None
    assert "계약전력과 요금적용전력의 틈" in entry.figure_caption

    sections = dataclasses.replace(full_sections, measures=(entry,))
    slide = _slide_by_key(build_slides(sections), sections, "measure_contract")
    tops = {
        shape.text_frame.text: Emu(shape.top).inches
        for shape in slide.shapes  # type: ignore[attr-defined]
        if shape.has_text_frame
    }
    assert tops["현재 계약전력"] < tops["절감액"], tops
    assert tops["하향 여지"] < tops["회수기간"], tops


def test_기본요금_변화없음도_사유를_각주로_내린다() -> None:
    """**큰 글씨 자리에 사유가 들어가 두 줄로 흘렀다** (53세션 6절).

    48세션이 「0원」 대신 이 말을 세우면서 사유가 통째로 지표 칸에 들어왔다.
    """
    from kwise.measures import BASE_FEE_UNCHANGED

    head, reason = split_reason(f"{BASE_FEE_UNCHANGED} — 요금적용전력 하한 30% 적용")
    assert head == BASE_FEE_UNCHANGED
    assert reason == "요금적용전력 하한 30% 적용"
    # 관계없는 값은 그대로 둔다.
    assert split_reason("53,575,000원") == ("53,575,000원", "")


# ===================================================================== 53세션 · 7절 3·4페이지


def test_3장_해석이_각주로_내려갔다(full_sections: DocumentSections) -> None:
    """**이 장의 본체는 표다** (53세션 7-1).

    자료가 얼마나 성한지는 표를 읽는 데 붙는 단서이지 제목 다음에 와야 할
    결론이 아니다.
    """
    from pptx.util import Emu

    from kwise.report import narrative
    from kwise.report.slides import NOTE_MARK

    guide = load_design_guide()
    slide = _slide_by_key(build_slides(full_sections), full_sections, "building")
    quality = full_sections.diagnosis.quality if full_sections.diagnosis else None
    lead = narrative.building_lead(quality)
    found = [
        (Emu(shape.top).inches, shape.text_frame.text)
        for shape in slide.shapes  # type: ignore[attr-defined]
        if shape.has_text_frame and lead in shape.text_frame.text
    ]
    assert len(found) == 1, found
    top, text = found[0]
    assert text.startswith(NOTE_MARK.strip()), text
    # 슬라이드 아래쪽이다 — 제목 바로 아래가 아니다.
    assert top > guide.slide.height_in * 0.8, top
    tables = [shape for shape in slide.shapes if shape.has_table]  # type: ignore[attr-defined]
    assert tables and Emu(tables[0].top).inches < top


def test_4장_캡션이_전제를_되풀이하지_않는다(full_sections: DocumentSections) -> None:
    """**캡션은 무엇을 그렸나를 적는다** (53세션 7-2)."""
    slide = _slide_by_key(build_slides(full_sections), full_sections, "usage_pattern")
    text = _slide_text(slide)
    assert "결측일은 그리지 않았습니다" not in text, text
    assert "일별 사용량" in text


def test_기온_그림은_범례가_아래다() -> None:
    """**오른쪽에는 기온 축의 이름이 이미 서 있다** (53세션 7-3).

    범례를 그 자리에 두면 축 이름 위에 겹쳐 둘 다 못 읽는다. 23세션의
    「바깥 오른쪽」 규약을 흔드는 것이 아니라 **그 자리가 비어 있지 않은
    유일한 그림**이다.
    """
    from pathlib import Path as _Path

    from kwise.report.slides import FULL_FIGURE, FULL_FIGURE_WITH_LEGEND

    # 범례가 높이를 먹으므로 **미리 낮게 굽는다.**
    assert FULL_FIGURE_WITH_LEGEND[1] < FULL_FIGURE[1]
    source = (SRC_ROOT / "report" / "figures.py").read_text(encoding="utf-8")
    body = source[
        source.index("def daily_temperature_png(") : source.index("def power_factor_day_png(")
    ]
    assert 'bbox_to_anchor=(0.5, -0.28)' in body, "범례가 아래로 내려가 있지 않습니다."
    assert "labelpad=8" in body, "기온 축 이름이 눈금에 붙어 있습니다."
    assert _Path(SRC_ROOT).is_dir()


# ===================================================================== 53세션 · 8절 각주와 부록


def test_역률_장에_용어_각주가_있다(
    sample_usage: UsageData, tariff: TariffTable, sample_bill: BillingResult,
    full_sections: DocumentSections,
) -> None:
    """**「지상역률」 을 세 번 말하면서 그것이 무엇인지는 없었다** (53세션 8-1)."""
    import dataclasses

    from kwise.measures import evaluate_power_factor
    from kwise.report.document import measure_entries
    from kwise.report.narrative import GLOSSARY_KEYS, terms
    from kwise.tariff import TariffSelection

    result = evaluate_power_factor(
        sample_usage,
        tariff,
        TariffSelection("general_b", "high_a", "I"),
        current_pct=92.0,
        target_pct=97.0,
        baseline=sample_bill,
    )
    entry = next(
        item for item in measure_entries(power_factor=result) if item.kind.key == "power_factor"
    )
    sections = dataclasses.replace(full_sections, measures=(entry,))
    text = _slide_text(
        _slide_by_key(build_slides(sections), sections, "measure_power_factor")
    )
    table = terms()
    for key in GLOSSARY_KEYS["measure_power_factor"]:
        assert f"{table[key].name} = " in text, key
    assert "지상역률" in text and "무효전력" in text and "유효전력" in text


def test_용어_각주가_없는_수단_장에는_까닭이_있다() -> None:
    """**전수 확인** (53세션 8-1). 빠진 셋은 그 자리를 다른 것이 쓰고 있다."""
    from kwise.measures import MEASURE_CATALOG
    from kwise.report.narrative import GLOSSARY_KEYS

    covered = {key for key in GLOSSARY_KEYS if key.startswith("measure_")}
    missing = {f"measure_{kind.key}" for kind in MEASURE_CATALOG} - covered
    assert missing == {"measure_tariff_switch", "measure_ess"}, missing
    source = (SRC_ROOT / "report" / "narrative.py").read_text(encoding="utf-8")
    assert "선택요금 전환   기본요금·전력량요금은 7장이 이미 깐다" in source
    assert "ESS            표식의 뜻이 표 아래 그 자리를 쓴다" in source


def test_부록이_현행과_최적을_장으로_가른다() -> None:
    """**한 장 맨 아래에 다음 요금제가 걸쳐 있었다** (53세션 8-2)."""
    from kwise.report.slides import APPENDIX_MIN_ROWS, appendix_chunks

    records = [(f"현행 {index}", "", "") for index in range(6)]
    records.append(("", "", ""))
    records += [(f"최적 {index}", "", "") for index in range(6)]
    chunks = appendix_chunks(records)
    assert len(chunks) == 2, chunks
    assert all(row[0].startswith("현행") for row in chunks[0])
    assert all(row[0].startswith("최적") for row in chunks[1])
    # **빈 줄은 싣지 않는다** — 묶음을 가르는 자리이지 내용이 아니다.
    assert not any(not any(v.strip() for v in row) for chunk in chunks for row in chunk)

    # **꼬리가 짧으면 앞장에 붙인다** — 세 줄짜리 장을 만들지 않는다.
    long = [(f"줄 {index}", "", "") for index in range(APPENDIX_ROW_LIMIT + 2)]
    merged = appendix_chunks(long)
    assert len(merged) == 1, [len(chunk) for chunk in merged]
    assert len(merged[0]) == APPENDIX_ROW_LIMIT + 2
    assert APPENDIX_MIN_ROWS > 2


def test_ESS_부록이_한_장이다(full_sections: DocumentSections) -> None:
    """**뒷장에 세 줄만 있었다** (53세션 8-4). 1-6 으로 자리가 생겼다."""
    import dataclasses

    from kwise.report.slides import APPENDIX_ROW_LIMIT, appendix_chunks
    from kwise.report.worksheet import Worksheet

    assert APPENDIX_ROW_LIMIT >= 13, "ESS 근거는 열셋이다."
    rows = [(f"줄 {index}", "", "") for index in range(13)]
    assert len(appendix_chunks(rows)) == 1

    # 제목이 **무엇의 근거인지** 밝힌다 (8-3).
    sheet = full_sections.worksheets[0]
    sections = dataclasses.replace(
        full_sections, worksheets=(Worksheet(sheet.key, sheet.title, sheet.rows),)
    )
    titles = [page.title for page in appendix_pages(sections)]
    assert all(" — " in title for title in titles), titles


# ===================================================================== 53세션 · 9절 렌더링


def _overlaps(first: tuple[float, float], second: tuple[float, float]) -> float:
    """두 세로 구간이 겹치는 길이 (in). 0 이면 안 겹친다."""
    return max(0.0, min(first[1], second[1]) - max(first[0], second[0]))


def _table_and_text_bands(
    slide: object,
) -> tuple[list[tuple[float, float]], list[tuple[float, float, str]]]:
    from pptx.util import Emu

    tables: list[tuple[float, float]] = []
    texts: list[tuple[float, float, str]] = []
    for shape in slide.shapes:  # type: ignore[attr-defined]
        if shape.top is None or not shape.height:
            continue
        band = (Emu(shape.top).inches, Emu(shape.top + shape.height).inches)
        if shape.has_table:
            tables.append(band)
        elif shape.has_text_frame and shape.text_frame.text.strip():
            texts.append((band[0], band[1], shape.text_frame.text))
    return tables, texts


def test_표가_아래_글씨를_덮지_않는다(full_sections: DocumentSections) -> None:
    """**실물을 렌더해서야 보였다** (53세션 9절).

    표 높이는 우리가 넘긴 값이 아니라 **줄 높이의 합**이라 계산과 어긋난다.
    개선안별 요약은 각주가 세 줄인데 예약이 두 줄어치였고, ESS 사양 표는
    캡션을 밑변에 딱 붙여 두어 마지막 줄에 얹혔다.
    """
    import dataclasses

    sections = dataclasses.replace(
        full_sections,
        measures=(*full_sections.measures, _ess_entry_with_spec()),
    )
    for index, slide in enumerate(build_slides(sections).slides, start=1):
        tables, texts = _table_and_text_bands(slide)
        for table in tables:
            for top, bottom, text in texts:
                # 표보다 위에서 시작하는 글(제목·해석·지표)은 대상이 아니다.
                if top < table[0]:
                    continue
                assert _overlaps(table, (top, bottom)) < 0.02, (
                    f"{index}장 — 표({table[0]:.2f}~{table[1]:.2f}in)가 "
                    f"「{text[:30]}」({top:.2f}in)를 덮습니다."
                )


#: 세 줄로 흐르는 각주. 실제 용어 각주보다 길게 지어 **줄 수가 늘 때**를 만든다.
_LONG_NOTE = " ".join(
    [
        "부하율 = 평균 수요 ÷ 최대 수요",
        "기저부하 비율 = 야간(22~8시) 평균 ÷ 주간 평균",
        "운영시간 외 부하 비중 = 운영시간(평일 9~18시) 밖 사용량 ÷ 전체",
        "요금적용전력 = 직전 12개월 최대수요로 매기는 기본요금의 기준 전력",
        "계시별 = 경부하·중간부하·최대부하로 단가가 갈리는 구간",
        "역률 = 유효전력 ÷ 피상전력, 무효전력 = 일을 하지 않고 계통을 오가는 전력",
    ]
)


def _note_bands(slide: object) -> tuple[float, list[tuple[float, float, str]]]:
    """각주 띠가 시작하는 y 와, 각주가 아닌 글의 세로 구간들."""
    from pptx.util import Emu

    from kwise.report.slides import NOTE_MARK

    note_top = float("inf")
    others: list[tuple[float, float, str]] = []
    for shape in slide.shapes:  # type: ignore[attr-defined]
        if shape.top is None or not shape.height:
            continue
        top = Emu(shape.top).inches
        bottom = top + Emu(shape.height).inches
        text = shape.text_frame.text.strip() if shape.has_text_frame else ""
        if text.startswith(NOTE_MARK.strip()):
            note_top = min(note_top, top)
        else:
            others.append((top, bottom, text or "그림·표"))
    return note_top, others


def test_각주가_세_줄로_흘러도_위를_덮지_않는다(
    full_sections: DocumentSections, monkeypatch: pytest.MonkeyPatch
) -> None:
    """**자리를 못박으면 줄이 늘 때 겹친다** (59세션 1절).

    53세션은 표 자리만 실측(:func:`~kwise.report.slides._note_top`)으로 옮기고
    **그림 덩어리 넷(4·5·6·7장)은 0.58in 고정** 그대로 두었다 — 그것이 두
    줄어치라, 각주가 세 줄로 흐르면 그림 캡션 위로 0.24in 이 얹혔다.

    실물에서는 지금 자료가 모두 한두 줄이라 보이지 않는다. **줄을 늘려 본다.**
    """
    from kwise.report import slides as slides_module

    monkeypatch.setattr(
        slides_module.narrative,
        "glossary_note",
        lambda *args, **kwargs: _LONG_NOTE,
    )

    guide = load_design_guide()
    slack = guide.slide.text_slack
    for index, slide in enumerate(build_slides(full_sections).slides, start=1):
        note_top, others = _note_bands(slide)
        if note_top == float("inf"):
            continue
        for top, bottom, text in others:
            # 상자 높이에는 렌더러 여유(slack)가 얹혀 있다. 글이 실제로 차지하는
            # 자리는 그만큼 위에서 끝난다.
            real_bottom = top + (bottom - top) / (1.0 + slack)
            assert real_bottom <= note_top + 0.01, (
                f"{index}장 — 「{text[:24]}」({real_bottom:.2f}in)가 "
                f"각주({note_top:.2f}in)를 덮습니다."
            )


def test_슬라이드_안에서_끝난다(full_sections: DocumentSections) -> None:
    """**각주는 아래에 붙어 있어 넘친 줄이 밖으로 나간다** (53세션 9절)."""
    import dataclasses

    guide = load_design_guide()
    sections = dataclasses.replace(
        full_sections,
        measures=(*full_sections.measures, _ess_entry_with_spec()),
    )
    for index, slide in enumerate(build_slides(sections).slides, start=1):
        tables, texts = _table_and_text_bands(slide)
        for _top, bottom, text in texts:
            assert bottom <= guide.slide.height_in + 0.01, f"{index}장 — 「{text[:20]}」"
        for table in tables:
            assert table[1] <= guide.slide.height_in + 0.01, f"{index}장 — 표"


# ============================================================ 60세션 — 각주 조립


def _priced_surplus_page() -> object:
    """**적용 단가가 붙은 잉여 장** — 각주가 두 줄이 되는 유일한 갈래다 (58세션).

    단가가 없으면 :meth:`SurplusResult.applied_price_note` 가 비어 각주가 한
    줄뿐이다. 「※」 가 두 번 서는 결함은 **둘째 줄이 있을 때만** 난다.
    """
    import dataclasses

    from kwise.measures.surplus import (
        EXTERNAL_SCENARIO,
        OFFSET_SCENARIO,
        SurplusScenario,
    )
    from kwise.report.document import surplus_page

    base = _surplus_result(weekday=105.0, weekend=20_000.0, holiday=3_311.0)
    priced = dataclasses.replace(
        base,  # type: ignore[type-var]
        external_price_won_per_kwh=140.0,
        scenarios=(
            SurplusScenario(OFFSET_SCENARIO, 2_545_000.0, "상계", "계약 변경"),
            SurplusScenario(EXTERNAL_SCENARIO, 3_267_824.0, "외부", "구매자 발굴"),
            base.scenarios[2],  # type: ignore[attr-defined]
        ),
    )
    page = surplus_page(priced, capacity_kwp=160.0, surplus_free_kwp=40.0)
    assert page is not None
    return page


def _unviable_ess_entry(
    sample_usage: UsageData, tariff: TariffTable, sample_bill: BillingResult
) -> MeasureEntry:
    """**ESS 성립 불가 장의 항목** — 그림도 표도 없고 지표만 두 줄인 장 (60세션).

    갑 종별은 기본요금이 계약전력에 붙어 성립할 수 없다 (56세션). 이 갈래는
    **대형·소형 어느 자료의 기본 계약에도 서지 않아서**, ``full_sections`` 만
    훑는 시험은 여기서 나는 결함을 못 잡는다 — 59세션 되짚기 ① 이 적은
    「갈래를 만든 세션과 그것이 실물에 서는지 본 세션이 달랐다」 가 이 자리다.
    """
    from kwise.measures.ess import ess_target_curve, refine_ess_target
    from kwise.report.document import measure_entries
    from kwise.tariff import TariffSelection

    selection = TariffSelection("general_a_1", "high_a", "I")
    curve = ess_target_curve(
        sample_usage.kw,
        15,
        baseline_demand_kw=5_293.44,
        base_fee_won_per_kw=float(tariff.rates(selection).base_won_per_kw),
    )
    optimum = refine_ess_target(
        sample_usage, tariff, selection, curve=curve, baseline=sample_bill
    )
    entry = next(
        item
        for item in measure_entries(ess_optimum=optimum, ess_curve=curve)
        if item.kind.key == "ess"
    )
    assert not entry.actionable and entry.facts and not entry.facts_first
    return entry


def test_각주_한_줄에_표식이_한_번만_선다(full_sections: DocumentSections) -> None:
    """**「※」 가 한 줄 가운데 또 서고 있었다** (60세션 1절).

    잉여 장이 자격요건 각주와 적용 단가 각주를 빈칸으로 이어 한 줄로 넘겼는데,
    **뒤 문장이 제 표식을 달고 온다.** :func:`mark_note` 는 맨 앞만 보므로
    가운데 것을 못 잡았다.

    조각을 잇는 것은 :func:`~kwise.report.narrative.note_line` 하나뿐이고,
    거기서 조각이 달고 온 표식을 뗀다.
    """
    from kwise.report.narrative import NOTE_MARK, note_line

    mark = NOTE_MARK.strip()
    # ① 잇는 자리가 표식을 뗀다.
    assert note_line(f"{mark} 앞", f"{mark} 뒤").count(mark) == 0
    assert note_line("앞입니다.", f"{mark} 뒤") == "앞입니다. 뒤"

    # ② **결함이 났던 그 장에서 확인한다.** 대형 자료에는 잉여 장이 없어
    #    ``full_sections`` 만 훑으면 이 시험은 뜰 수가 없다 (59세션 되짚기 ①).
    import dataclasses

    page = _priced_surplus_page()
    assert len(page.notes) >= 2, "잉여 장 각주가 둘이어야 이 시험이 뜻을 가집니다."
    assert any(mark in text for text in page.notes), "문장 하나가 표식을 달고 온다."

    guide = load_design_guide()
    sections = dataclasses.replace(
        full_sections, measures=_all_measures(full_sections), surplus=page
    )
    deck = build_slides(sections)
    surplus_lines = _note_lines(_slide_by_key(deck, sections, "surplus"), guide)
    assert len(surplus_lines) >= 2, f"잉여 장 각주가 두 줄이어야 합니다: {surplus_lines}"

    # ③ 덱 어느 줄에도 표식이 두 번 서지 않는다.
    for index, slide in enumerate(deck.slides, 1):
        for line in _note_lines(slide, guide):
            assert line.count(mark) <= 1, f"{index}장 각주에 {mark} 가 둘입니다: {line}"


def test_마침표로_끝난_조각_뒤에는_구분점이_없다(
    full_sections: DocumentSections,
    sample_usage: UsageData,
    tariff: TariffTable,
    sample_bill: BillingResult,
) -> None:
    """**마침표와 구분점이 겹쳤다** (60세션 2절).

    「… 사양을 정하지 않았습니다**. ·** 투자비 미산출」 — 사유가 마침표로 끝나는
    장이 ESS 성립 불가 하나뿐이라 다른 장에서는 드러나지 않았다.
    **문장이 이미 끝났으면 빈칸이면 족하다.**
    """
    import re

    from kwise.report.narrative import NOTE_JOIN, note_line

    # ① 마침표로 끝나면 빈칸, 아니면 구분점.
    assert note_line("않았습니다.", "투자비 미산출") == "않았습니다. 투자비 미산출"
    assert note_line("사유 (괄호로 끝남)", "다음") == f"사유 (괄호로 끝남){NOTE_JOIN}다음"

    # ② **마침표로 끝나는 사유가 있는 장에서 확인한다.** 그런 사유는 ESS 성립
    #    불가 장 하나뿐이라, 대형 자료만 훑으면 이 시험은 뜰 수가 없다.
    import dataclasses

    from kwise.report.slides import _measure_note

    entry = _unviable_ess_entry(sample_usage, tariff, sample_bill)
    reasons = [split_reason(entry.slide_saving)[1], split_reason(entry.investment)[1]]
    assert any(text and text.rstrip().endswith(".") for text in reasons), (
        f"마침표로 끝나는 사유가 있어야 이 시험이 뜻을 가집니다: {reasons}"
    )
    assert not re.search(r"[.。]\s*·", _measure_note(entry))

    guide = load_design_guide()
    measures = tuple(
        entry if item.kind.key == "ess" else item for item in _all_measures(full_sections)
    )
    sections = dataclasses.replace(full_sections, measures=measures)
    for index, slide in enumerate(build_slides(sections).slides, 1):
        for line in _note_lines(slide, guide):
            assert not re.search(r"[.。]\s*·", line), f"{index}장 각주가 겹칩니다: {line}"


def test_그림_없는_장은_지표_두_줄을_붙인다(
    full_sections: DocumentSections,
    sample_usage: UsageData,
    tariff: TariffTable,
    sample_bill: BillingResult,
) -> None:
    """**한 장이 둘로 읽혔다** (60세션 3절).

    36세션의 「남는 높이 **가운데** 앉힌다」 는 **그림 덩어리를 위해** 만든
    규약이다. 그림이 없는 장에는 채울 것이 없어 가운데가 곧 「멀리」가 된다 —
    ESS 성립 불가 장이 지표 셋과 지표 다섯 사이를 1.46in 벌리고 있었다.

    **간격은 :data:`~kwise.report.slides._STAT_ROW_GAP` 이 정한다.** 값을
    코드에 흩어 두면 다음에 같은 장이 하나 더 생길 때 또 벌어진다.
    """
    import dataclasses

    from pptx.util import Emu

    from kwise.report.slides import _STAT_ROW_GAP

    entry = _unviable_ess_entry(sample_usage, tariff, sample_bill)

    # 수단을 모두 켜 놓고 ESS 만 이 갈래로 바꿔 끼운다.
    measures = tuple(
        entry if item.kind.key == "ess" else item for item in _all_measures(full_sections)
    )
    sections = dataclasses.replace(full_sections, measures=measures)
    slide = _slide_by_key(build_slides(sections), sections, "measure_ess")
    # **라벨로 두 덩어리를 가른다.** 좌표만 훑으면 제목·결론 사이의 숨까지
    # 섞여 무엇을 재는지 흐려진다.
    upper = {"절감액", "투자비", "회수기간"}
    lower = {label for label, _value in entry.facts}
    assert lower, "아래 덩어리가 있어야 이 시험이 뜻을 가집니다."

    boxes = [
        (Emu(shape.top).inches, Emu(shape.top + shape.height).inches, shape.text_frame.text.strip())
        for shape in slide.shapes  # type: ignore[attr-defined]
        if shape.has_text_frame and shape.top is not None and shape.height
    ]
    assert any(text in upper for _t, _b, text in boxes), "위 덩어리를 못 찾았습니다."
    lower_top = min(top for top, _b, text in boxes if text in lower)
    # **값 줄까지 포함해 아래끝을 잰다** — 라벨만 보면 그 아래 값 칸이 틈으로
    # 잘못 셈된다.
    upper_bottom = max(bottom for _t, bottom, _x in boxes if bottom <= lower_top + 0.01)
    gap = lower_top - upper_bottom
    assert 0.0 <= gap <= _STAT_ROW_GAP + 0.02, (
        f"지표 두 줄이 {gap:.2f}in 갈렸습니다 — 규약은 {_STAT_ROW_GAP}in 입니다."
    )


# ==================================================== 60세션 10절 — 주의사항 표와 지표 자리


def _blind(sections: DocumentSections) -> DocumentSections:
    """**그림 굽기가 통째로 실패한 상태** — `_safe_figure` 가 전부 ``None`` 인 자리.

    13세션이 「차트 하나 때문에 문서 전체를 잃는 편보다 그림 없이 표만 내는 편이
    낫다」 로 만든 폴백이다. 자료로는 못 만든다 — 그림이 성공하는 한 수단 장은
    늘 그림 갈래로 돌아가기 때문이다.
    """
    import dataclasses

    return dataclasses.replace(
        sections,
        measures=tuple(
            dataclasses.replace(entry, figure=None, figures=(), spec_table=())
            for entry in _all_measures(sections)
        ),
    )


def _caution_table(slide: object) -> object | None:
    for shape in slide.shapes:  # type: ignore[attr-defined]
        if (
            getattr(shape, "has_table", False)
            and shape.has_table
            and shape.table.rows[0].cells[0].text.strip() == "주의사항"
        ):
            return shape.table
    return None


def test_그림이_안_구워지면_주의사항_표가_선다(full_sections: DocumentSections) -> None:
    """**뜨지 않는 갈래는 없는 갈래와 같다** (46세션 · 60세션 10절).

    벌 여섯 어디에도 이 표가 서지 않아 「죽은 갈래인가」 를 물었다. 죽지
    않았다 — **그림 굽기 실패 폴백**이라 정상 경로에서 안 설 뿐이다.
    조건이 실제로 있으므로 시험이 그 조건을 세운다.
    """
    sections = _blind(full_sections)
    entry = next(
        item for item in sections.measures if item.actionable and item.cautions
    )
    slide = _slide_by_key(build_slides(sections), sections, f"measure_{entry.kind.key}")
    table = _caution_table(slide)
    assert table is not None, f"{entry.kind.key} 장에 주의사항 표가 서야 합니다."
    body = [table.rows[index].cells[0].text.strip() for index in range(1, len(table.rows))]
    assert body and body != ["—"], f"표가 비어 있습니다: {body}"
    assert body[0] in entry.cautions, "표가 그 수단의 주의사항을 그대로 싣는다."


def test_실행할_것이_없으면_빈_주의사항_표를_그리지_않는다(
    full_sections: DocumentSections,
) -> None:
    """**「주의사항 / —」 만 남는 자리가 있었다** (60세션 10절).

    계약전력 조정이 하향 여지가 없으면 ``actionable=False`` 인데 **지표를 먼저
    세우는 장**(``facts_first``)이라, 그림이 실패하면 지표 갈래가 안 잡고
    주의사항 갈래로 흘렀다. `_cautions` 는 실행할 것이 없으면 빈 것을 돌려주므로
    머리글과 「—」 한 줄짜리 표가 섰다. **뜨지 않느니만 못한 표다.**
    """
    sections = _blind(full_sections)
    entry = next(
        item
        for item in sections.measures
        if not item.actionable and item.facts_first and item.facts
    )
    slide = _slide_by_key(build_slides(sections), sections, f"measure_{entry.kind.key}")
    assert _caution_table(slide) is None, "빈 주의사항 표가 서면 안 됩니다."
    # 지표는 위에서 이미 섰다 — 볼 것이 사라진 것이 아니다.
    text = _slide_text(slide)
    assert any(label in text for label, _value in entry.facts)


def test_지표_셋은_수단_장마다_같은_높이에_선다(full_sections: DocumentSections) -> None:
    """**여백보다 위계가 먼저다** (60세션 9-5 판정 · 10절에 규약으로).

    고객은 장을 넘기며 **같은 자리**에서 큰 숫자를 찾는다. 한 장만 내리면
    넘길 때마다 숫자가 위아래로 흔들린다 — ESS 성립 불가 장의 빈자리를 메우려고
    지표를 내리는 안을 버린 까닭이다.

    자리는 :func:`~kwise.report.slides._measure_stats_top` 이 쥔다. **결론이
    한 줄인 장과 두 줄인 장은 다르다** — 값이 아니라 **산식이 같아야 한다.**
    """
    import dataclasses

    from pptx.util import Emu

    from kwise.report.slides import _measure_stats_top

    sections = dataclasses.replace(full_sections, measures=_all_measures(full_sections))
    deck = build_slides(sections)
    keys = [spec.key for spec in slide_specs(sections)]
    seen: dict[float, list[str]] = {}
    for entry in sections.measures:
        key = f"measure_{entry.kind.key}"
        if key not in keys:
            continue
        slide = list(deck.slides)[keys.index(key)]
        tops = [
            Emu(shape.top).inches
            for shape in slide.shapes
            if shape.has_text_frame
            and shape.top is not None
            and shape.text_frame.text.strip() == "절감액"
        ]
        assert len(tops) == 1, f"{key} 에 「절감액」 라벨이 하나여야 합니다: {tops}"
        seen.setdefault(round(tops[0], 2), []).append(key)
    # 결론 줄 수가 갈래를 만든다 — 자리는 둘까지다 (한 줄 · 두 줄).
    assert len(seen) <= 2, f"지표 셋이 세 자리 넘게 흩어졌습니다: {seen}"
    assert _measure_stats_top(1.40, 0.72) == pytest.approx(2.12)
    assert _measure_stats_top(1.40, 0.44) == pytest.approx(1.84)


# ============================================= 60세션 11절 — 조용한 폴백에 기록을 남긴다


def test_그림이_안_구워지면_기록이_남는다() -> None:
    """**조용한 폴백은 21세션이 이미 걷어낸 모양이다** (60세션 11절 · T5).

    산출물 쪽에는 「그림 하나 때문에 덱 전체를 잃지 않는다」 는 폴백이 둘 있다.
    폴백 자체는 옳다 — 걷어내지 않는다. 다만 **조용하면 그림이 빠진 장을 받아도
    알 길이 없다.** 걷어내는 대신 남긴다.
    """
    from kwise.report.document import _safe_figure
    from kwise.report.figures import FigureFailureCollector

    def broken() -> bytes:
        raise RuntimeError("글꼴이 없다")

    with FigureFailureCollector() as collector:
        assert _safe_figure(broken, "태양광 · 일별 발전량") is None
        # **성공은 기록하지 않는다** — 실패만 센다.
        assert _safe_figure(lambda: b"png", "역률 개선 · 전력삼각형") == b"png"
    assert len(collector.messages) == 1, collector.messages
    line = collector.messages[0]
    assert "태양광 · 일별 발전량" in line, "어느 수단의 어느 그림인지 남아야 합니다."
    assert "RuntimeError" in line and "글꼴이 없다" in line, "무슨 예외였는지 남아야 합니다."


def test_손잡이는_나가면_떨어진다() -> None:
    """**벌을 잇달아 뽑을 때 앞 벌의 실패가 뒤 벌에 묻으면 안 된다** (60세션 11절)."""
    from kwise.report.document import _safe_figure
    from kwise.report.figures import FigureFailureCollector

    def broken() -> bytes:
        raise ValueError("깨졌다")

    with FigureFailureCollector() as first:
        _safe_figure(broken, "첫 벌")
    with FigureFailureCollector() as second:
        _safe_figure(broken, "둘째 벌")
    assert len(first.messages) == 1 and "첫 벌" in first.messages[0]
    assert len(second.messages) == 1 and "둘째 벌" in second.messages[0]


def test_덱_라벨이_그림_실패를_0_까지_적는다() -> None:
    """**0 이면 0 이라 적는다** (60세션 11절).

    줄이 아예 없으면 「기록을 안 남긴 것」 과 「실패가 없던 것」 이 갈리지 않는다.
    장 수가 같아도 속이 다를 수 있다는 것도 라벨이 말한다 — 소형 21장과 대형 을
    21장은 잉여 장이 붙고 부록 ESS 가 빠져 **상쇄된 것**이다.
    """
    import sys

    sys.path.insert(0, str(Path("tools").resolve()))
    import render_deck  # type: ignore[import-not-found]

    case = render_deck.BY_KEY["small-b"]
    quiet = render_deck.label_text(case, ("표지", "목차"), ())
    assert "그림 실패    0개" in quiet, quiet
    assert "장 수가 같아도 같은 덱이 아니다" in quiet
    assert "우연히 같다" in quiet

    noisy = render_deck.label_text(case, ("표지",), ("그림 실패 — 태양광 · 일별 발전량: ...",))
    assert "그림 실패    1개" in noisy
    assert "태양광 · 일별 발전량" in noisy


# ================================================= 60세션 12절 — T6 판정을 집행한다


def test_기온_그림이_실패하면_사용량만_그리고_기록을_남긴다(
    full_sections: DocumentSections,
) -> None:
    """**삼키던 자리 둘 가운데 하나** (60세션 11-3 의 2번 · 12절 집행).

    기온이 빠진 채 사용량만 그려도 장은 채워지므로 **아무도 모른 채 지나간다.**
    11절이 기록 통로를 붙였고, 이 시험이 그 통로를 실제로 타는지 본다.
    """
    import dataclasses

    import pandas as pd

    from kwise.report import figures
    from kwise.report.figures import FigureFailureCollector
    from kwise.report.slides import _usage_figure

    # **기온을 붙여야 갈래가 열린다.** 없으면 애초에 기온을 안 그리므로
    # 이 시험은 아무것도 못 본다 — 건너뛴 시험은 없는 시험이다.
    index = pd.date_range(full_sections.usage.meta.start, periods=48, freq="h")
    warm = dataclasses.replace(
        full_sections, temperature=pd.Series(range(48), index=index, dtype=float)
    )
    assert "기온" in _usage_figure(warm)[1], "성한 자료에서는 기온을 그린다."

    def boom(*_args: object, **_kwargs: object) -> bytes:
        raise RuntimeError("기온 축이 깨졌다")

    original = figures.daily_temperature_png
    try:
        figures.daily_temperature_png = boom  # type: ignore[assignment]
        with FigureFailureCollector() as collector:
            png, caption = _usage_figure(warm)
    finally:
        figures.daily_temperature_png = original  # type: ignore[assignment]

    assert png, "기온이 없어도 사용량 그림은 나온다 — 장을 비우지 않는다."
    assert "기온" not in caption, f"캡션이 사용량만 말해야 합니다: {caption}"
    assert len(collector.messages) == 1, collector.messages
    assert "일별 기온" in collector.messages[0]
    assert "RuntimeError" in collector.messages[0]


def test_진단이_없으면_제목만_선다(full_sections: DocumentSections) -> None:
    """**물음 넷을 한 자리로 모았다** (60세션 11-3 의 8~11 · 12절 집행).

    네 장이 저마다 「값이 없으면 조용히 돌아선다」 를 들고 있었고, 주석은 넷 다
    「진단 없이 부르지 않는다」 라고 적고 있었다. **틀린 말이었다** —
    :func:`slide_specs` 는 진단 유무를 보지 않아 넷을 늘 자리표에 넣는다.

    그래서 걷어내지 않고 **모으기만 했다.** 장은 그대로 서고, 비어 있다는
    사실도 그대로다 — 바뀐 것은 그 물음을 읽을 자리가 하나가 된 것뿐이다.
    """
    import dataclasses

    from kwise.report.slides import _NEEDS_DIAGNOSIS, slide_specs

    blind = dataclasses.replace(full_sections, diagnosis=None)
    keys = [spec.key for spec in slide_specs(blind)]
    # 자리표는 넷을 그대로 낸다 — 「부르지 않는다」 가 아니다.
    assert set(keys) >= _NEEDS_DIAGNOSIS, keys

    deck = build_slides(blind)
    for key in sorted(_NEEDS_DIAGNOSIS):
        slide = list(deck.slides)[keys.index(key)]
        texts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        assert texts == [SLIDE_TITLES[key]], f"{key} 는 제목만 서야 합니다: {texts}"
